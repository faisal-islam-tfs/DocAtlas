#!/usr/bin/env python3
"""
DocAtlas document processing pipeline:
- Extract text from PDF/DOCX/PPTX/XLS/XLSX
- Auto-unpack zip archives into a staging area and process supported contents
- Summarize, categorize, tag with Azure OpenAI
- Detect exact duplicates via hashes and near-duplicates via embeddings
- Output review Excel files plus a compressed full-text archive
- Move files into category folders (duplicates to <category>_Duplicate)
"""
from __future__ import annotations

import argparse
import gzip
import html
import io
import hashlib
import json
import logging
import math
import os
import random
import re
import shutil
import subprocess
import sys
import time
import tempfile
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Tuple
from concurrent.futures import ThreadPoolExecutor, as_completed
import threading
import queue

import numpy as np
import pandas as pd
import requests

try:
    import pdfplumber
except Exception:  # pragma: no cover
    pdfplumber = None

try:
    import docx
except Exception:  # pragma: no cover
    docx = None

try:
    import pptx
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except Exception:  # pragma: no cover
    pptx = None
    MSO_SHAPE_TYPE = None

try:
    import openpyxl
except Exception:  # pragma: no cover
    openpyxl = None

try:
    import pytesseract
except Exception:  # pragma: no cover
    pytesseract = None

try:
    from PIL import Image
except Exception:  # pragma: no cover
    Image = None

try:
    from pdf2image import convert_from_path
except Exception:  # pragma: no cover
    convert_from_path = None

try:
    import ocrmypdf
except Exception:  # pragma: no cover
    ocrmypdf = None

try:
    import tkinter as tk
    from tkinter import filedialog, messagebox
    from tkinter import ttk
except Exception:  # pragma: no cover
    tk = None
    messagebox = None
    ttk = None

try:
    from tqdm import tqdm
except Exception:  # pragma: no cover
    tqdm = None


SUPPORTED_EXTS = {".pdf", ".doc", ".docx", ".ppt", ".pptx", ".xls", ".xlsx"}
INVALID_WIN_CHARS = r'<>:"/\\|?*'
ILLEGAL_EXCEL_CHARS_RE = re.compile(r"[\x00-\x08\x0b\x0c\x0e-\x1f]")

DEFAULT_CHAT_BASE_URL = "https://api.geneai.thermofisher.com/dev/gpt5"
DEFAULT_EMBEDDINGS_BASE_URL = "https://api.geneai.thermofisher.com/dev/embeddingsv2"
DEFAULT_API_VERSION = "2025-03-01-preview"
DEFAULT_CHAT_DEPLOYMENT = "gpt-5.2"
DEFAULT_EMBEDDINGS_DEPLOYMENT = "text-embedding-3-small"
DEFAULT_API_KEY_HEADER = "api-key"
DEFAULT_CHAT_PATH = "/openai/deployments/{deployment}/chat/completions"
DEFAULT_EMBEDDINGS_PATH = "/openai/deployments/{deployment}/embeddings"
# Safe network defaults for long server runs with intermittent DNS/network instability.
DEFAULT_API_DELAY_SEC = 0.5
DEFAULT_API_MAX_RETRIES = 10
DEFAULT_API_RETRY_BASE_SEC = 2.0
DEFAULT_API_RETRY_MAX_SEC = 60.0
DEFAULT_API_TIMEOUT_SEC = 150

MAX_CHARS_PER_CHUNK = 12000
MAX_ARTICLE_CHARS = 20000
DUPLICATE_THRESHOLD = 0.97
NEAR_DUP_MUTUAL = 0.78
NEAR_DUP_STRONG = 0.84
MIN_EXTRACTED_CHARS = 200
MIN_ARTICLE_BODY_CHARS = 400
MIN_SPLIT_TOTAL_CHARS = 2500
MIN_SPLIT_SECTIONS = 2
MIN_SECTION_CHARS = 900
MAX_SINGLE_SECTION_SHARE = 0.82
MIN_BOUNDARY_GAP_LINES = 12
MIN_BOUNDARY_GAP_CHARS = 800
MIN_EMBEDDING_CHARS = 500
MIN_EMBEDDING_CHARS_SUMMARY = 200
MAX_TAGS = 10
FALLBACK_DOC_LONG_SENTENCES = 7
FALLBACK_DOC_SHORT_SENTENCES = 2
FALLBACK_ARTICLE_SENTENCES = 8
FALLBACK_MAX_SENTENCE_CHARS = 420
FALLBACK_MIN_SENTENCE_CHARS = 35
FALLBACK_MAX_TOTAL_CHARS_DOC_LONG = 2400
FALLBACK_MAX_TOTAL_CHARS_DOC_SHORT = 420
FALLBACK_MAX_TOTAL_CHARS_ARTICLE = 3200
MAX_DOC_SUMMARY_INPUT_CHARS = 180000
MAX_DOC_SUMMARY_CHUNKS = 12
MAX_DOC_SUMMARY_FALLBACK_SOURCE_CHARS = 140000
RESUME_FILENAME = "resume.json"
LAST_RUN_STATS_FILENAME = "last_run_stats.json"
DEFAULT_EMBEDDINGS_SOURCE = "full_text"
DEFAULT_CATEGORY_PATH_MAP_FILENAME = "category_path_map.json"
DEFAULT_INCLUDE_FULL_TEXT_OUTPUT = True
DEFAULT_ESTIMATE_SEC_PER_FILE = 50.0
DEFAULT_ESTIMATE_SEC_PER_MB = 1.5
EMBEDDINGS_SOURCE_NONE = "none"
UNREADABLE_CATEGORY = "Unreadable"
SUMMARY_STOPWORDS = {
    "a",
    "an",
    "and",
    "are",
    "as",
    "at",
    "be",
    "by",
    "for",
    "from",
    "has",
    "in",
    "into",
    "is",
    "it",
    "its",
    "of",
    "on",
    "or",
    "that",
    "the",
    "their",
    "this",
    "to",
    "was",
    "were",
    "will",
    "with",
    "using",
    "use",
    "used",
    "can",
    "may",
    "not",
    "which",
    "than",
    "we",
    "you",
    "your",
    "our",
    "these",
    "those",
    "such",
    "also",
    "per",
    "via",
    "within",
    "across",
    "about",
    "after",
    "before",
    "if",
    "when",
    "while",
    "do",
    "does",
    "did",
    "done",
    "no",
    "yes",
    "up",
    "out",
    "over",
    "under",
}

# Deterministic category hints tuned for qPCR-style corpora.
# Keys are normalized category names.
CATEGORY_HINT_PHRASES: Dict[str, List[Tuple[str, float]]] = {
    "agt media": [
        ("agt media", 4.5),
        ("advanced granulation technology", 2.8),
    ],
    "cell culture antibiotics": [
        ("antibiotic", 2.6),
        ("antibiotics", 2.6),
        ("antimycotic", 2.4),
        ("penicillin streptomycin", 3.4),
        ("gentamicin", 2.4),
    ],
    "cell isolation": [
        ("cell isolation", 4.5),
        ("primary cell isolation", 4.8),
        ("cell isolation kit", 4.2),
        ("isolation kit", 2.8),
        ("cardiomyocyte isolation", 4.0),
        ("neuron isolation", 4.0),
    ],
    "celleste image analysis software": [
        ("celleste", 5.0),
        ("celleste image analysis software", 4.6),
        ("evos imaging software", 3.2),
    ],
    "assay design": [
        ("assay design", 0.8),
        ("design tool", 2.5),
        ("custom assay", 1.8),
        ("mirbase", 2.5),
        ("target sequence", 1.8),
        ("context sequence", 1.6),
        ("amplicon", 1.4),
        ("cadt", 2.5),
    ],
    "copy number variation": [
        ("copy number", 4.0),
        ("cnv", 4.0),
        ("copycaller", 4.0),
        ("calibrator", 2.5),
        ("reference assay", 2.5),
        ("z-score", 2.0),
        ("delta ct", 1.8),
        ("ddct", 1.8),
    ],
    "dissociation reagents": [
        ("dissociation", 4.2),
        ("dissociation reagent", 4.2),
        ("detachment", 2.8),
        ("tryple", 4.0),
        ("trypsin", 3.2),
        ("accutase", 3.8),
        ("accumax", 3.2),
        ("cell stripper", 3.2),
    ],
    "dpm": [
        ("dpm media", 4.2),
        ("dpm", 2.6),
    ],
    "gene expression": [
        ("gene expression", 5.0),
        ("expression assay", 3.2),
        ("taqman gene expression", 4.0),
        ("mirna", 2.2),
        ("taqman mirna", 3.0),
        ("relative quantification", 3.0),
        ("ddct", 2.5),
        ("delta delta ct", 2.5),
        ("endogenous control", 3.0),
        ("housekeeping gene", 2.5),
        ("reference gene", 2.2),
        ("transcript", 2.0),
    ],
    "stepone": [
        ("stepone", 4.0),
        ("steponeplus", 4.0),
    ],
    "snp genotyping": [
        ("snp", 3.2),
        ("genotyping", 3.2),
        ("genotype", 2.8),
        ("allelic discrimination", 3.0),
        ("genotyper", 3.0),
        ("allele", 2.0),
        ("polymorphism", 2.0),
        ("cluster plot", 2.0),
    ],
    "seqstudio": [
        ("seqstudio", 4.0),
        ("sanger", 2.5),
        ("capillary electrophoresis", 3.0),
    ],
    "reagents": [
        ("reagent", 1.0),
        ("master mix", 1.2),
        ("enzyme", 0.8),
        ("buffer", 0.7),
        ("kit", 0.5),
    ],
    "instrumentation": [
        ("instrument", 2.8),
        ("real-time pcr system", 3.0),
        ("thermal cycler", 2.8),
        ("quantstudio", 3.0),
        ("platform", 1.6),
    ],
    "gibco manufacturing and packaging docs": [
        ("gibco manufacturing", 4.6),
        ("manufacturing and packaging", 4.2),
        ("packaging docs", 4.0),
        ("label claim", 2.6),
    ],
    "data analysis": [
        ("data analysis", 3.2),
        ("software", 2.0),
        ("threshold", 1.5),
        ("baseline", 1.4),
        ("normalization", 1.6),
    ],
    "hcs": [
        ("hcs", 4.6),
        ("high content screening", 5.0),
        ("high content analysis", 3.8),
    ],
    "liquid cell culture": [
        ("liquid media", 4.2),
        ("liquid cell culture", 4.6),
        ("basal medium", 2.8),
    ],
    "ecm and 3d culture": [
        ("extracellular matrix", 4.2),
        ("3d culture", 4.4),
        ("three dimensional culture", 4.0),
        ("matrix", 1.8),
        ("hydrogel", 2.4),
    ],
    "nutritional supplements and other reagents": [
        ("nutritional supplement", 4.2),
        ("nutritional supplements", 4.2),
        ("supplement", 1.8),
        ("feed supplement", 2.4),
    ],
    "qc": [
        ("quality control", 4.0),
        ("qc", 3.5),
        ("certificate of analysis", 4.0),
        ("coa", 3.0),
        ("specification", 1.8),
        ("lot release", 2.0),
    ],
    "troubleshooting": [
        ("troubleshooting", 4.0),
        ("troubleshoot", 3.6),
        ("low signal", 2.0),
        ("no amplification", 2.2),
    ],
    "arctus": [
        ("arcturus", 5.0),
        ("laser capture microdissection", 4.0),
        ("lcm", 2.2),
        ("pico pure", 3.0),
        ("histogene", 3.0),
    ],
    "custom dna oligos": [
        ("custom dna oligo", 4.2),
        ("custom oligo", 3.8),
        ("primer", 1.4),
        ("oligo", 1.8),
        ("sirna", 2.0),
        ("plate map", 2.4),
        ("well position", 2.0),
        ("extinction coefficient", 2.0),
        ("desalted", 1.8),
        ("tm", 1.2),
    ],
    "magmax and kingfisher": [
        ("magmax", 4.0),
        ("kingfisher", 4.0),
        ("bind wash elute", 1.8),
        ("magnetic particle processor", 2.0),
    ],
    "qubit and quant-it": [
        ("qubit", 4.5),
        ("quant-it", 3.8),
        ("quant it", 3.8),
        ("fluorometer", 2.0),
    ],
    "superscript reverse transcriptases": [
        ("superscript", 4.5),
        ("reverse transcriptase", 2.8),
    ],
    "microrna reverse transcription kits": [
        ("mirna reverse transcription", 4.2),
        ("microrna reverse transcription", 4.2),
        ("stem-loop", 1.8),
    ],
    "high-capacity reverse transcription kits": [
        ("high capacity reverse transcription", 5.0),
        ("high-capacity reverse transcription", 5.0),
        ("reverse transcription kit", 2.5),
    ],
    "thermal cycler plastics and reagents": [
        ("microamp", 4.0),
        ("optical adhesive film", 3.2),
        ("reaction plate", 2.2),
        ("strip tube", 2.0),
        ("optical plate", 2.0),
        ("pcr plate", 2.6),
        ("384 well", 1.8),
        ("96 well", 1.6),
        ("skirted plate", 2.0),
        ("optical adhesive", 3.0),
    ],
    "electrophoresis reagents and kits": [
        ("e gel", 4.0),
        ("e base", 4.0),
        ("safe imager", 4.0),
        ("ibase", 3.5),
        ("gel cassette", 2.4),
        ("dna ladder", 2.0),
        ("agarose gel", 1.8),
    ],
    "mol bio sample prep": [
        ("sample prep", 3.8),
        ("benchpro", 4.4),
        ("prepstation", 4.0),
        ("tempus", 3.4),
        ("vacuum manifold", 2.2),
        ("nucleic acid prepstation", 4.0),
        ("dna purification", 2.0),
        ("rna isolation", 2.0),
    ],
    "nanodrop": [
        ("nanodrop", 5.0),
        ("nanodrop one", 4.4),
        ("nanodrop eight", 4.4),
        ("microvolume", 2.4),
        ("spectrophotometer", 1.8),
        ("scivault", 1.6),
    ],
    "spectra": [
        ("spectra", 4.6),
        ("spectral", 2.2),
    ],
    "tali": [
        ("tali", 4.8),
        ("tali image based cytometer", 4.4),
    ],
    "protein expression": [
        ("protein expression", 4.6),
        ("bluegrass", 4.0),
        ("algae", 2.2),
        ("baculo", 2.8),
        ("baculodirect", 4.0),
        ("chlamydomonas", 2.4),
        ("synechococcus", 2.4),
    ],
    "protein affinity purification": [
        ("protein affinity purification", 4.8),
        ("affinity purification", 3.8),
        ("protein a", 2.4),
        ("protein g", 2.4),
        ("sepharose", 2.2),
        ("agarose", 1.8),
    ],
    "protein assays": [
        ("protein assay", 4.6),
        ("protein assays", 4.6),
        ("bca", 2.8),
        ("bradford", 2.4),
        ("660nm", 2.4),
        ("quant it protein", 2.8),
        ("quant-it protein", 2.8),
    ],
    "western blotting": [
        ("western blot", 4.8),
        ("western blotting", 4.8),
        ("immunoblot", 3.6),
        ("ecl", 2.0),
        ("blot", 1.4),
    ],
    "water": [
        ("water purification", 5.0),
        ("water purification system", 5.0),
        ("water analysis", 4.5),
        ("water quality", 3.5),
        ("ultrapure water", 4.5),
        ("ultra pure water", 4.5),
        ("nuclease free water", 4.0),
        ("rnase free water", 4.0),
        ("dnase free water", 4.0),
        ("deionized water", 3.5),
        ("distilled water", 3.5),
        ("molecular biology grade water", 4.5),
        ("pcr grade water", 4.0),
        ("water nuclease free", 4.0),
    ],
    "dynabeads": [
        ("beadretriever", 4.0),
        ("immunomagnetic separation", 2.8),
        ("chromagar", 2.2),
        ("o157", 1.8),
        ("giardia", 2.2),
        ("cryptosporidium", 2.2),
        ("legionella", 2.0),
    ],
    "transfection": [
        ("lipofectamine", 4.0),
        ("invivofectamine", 4.2),
        ("neon", 3.4),
        ("xenon", 3.4),
        ("electroporation", 2.4),
        ("microporation", 2.4),
        ("transfection reagent", 2.2),
    ],
}

CATEGORY_GENERIC_TOKENS = {
    "assay",
    "design",
    "data",
    "analysis",
    "quality",
    "control",
    "reagents",
    "instrumentation",
    "troubleshooting",
    "copy",
    "number",
    "variation",
    "gene",
    "expression",
    "snp",
    "genotyping",
}

CATEGORY_REQUIRED_PHRASES: Dict[str, Tuple[str, ...]] = {
    "water": (
        "water purification",
        "water purification system",
        "water analysis",
        "water quality",
        "ultrapure water",
        "ultra pure water",
        "nuclease free water",
        "rnase free water",
        "dnase free water",
        "deionized water",
        "distilled water",
        "molecular biology grade water",
        "pcr grade water",
        "water nuclease free",
    ),
}

CATEGORY_REQUIRED_TOKEN_HITS: Dict[str, int] = {
    "water": 2,
}

CATEGORY_PATH_COMPONENT_HINTS: Dict[str, List[Tuple[str, float]]] = {
    "agt media": [
        ("cell culture media", 2.2),
        ("agt media", 5.0),
    ],
    "arctus": [
        ("arcturus", 5.0),
        ("lcm", 2.5),
    ],
    "cell culture antibiotics": [
        ("cell culture reagents antibiotics and supplements", 2.8),
        ("antibiotics", 4.8),
    ],
    "cell isolation": [
        ("cell culture reagents antibiotics and supplements", 2.8),
        ("cell isolation and dissociation reagents", 4.8),
    ],
    "celleste image analysis software": [
        ("software celleste", 5.0),
        ("celleste", 4.6),
        ("ca evos imaging software celleste", 4.8),
    ],
    "ca instrument accessories and reagents": [
        ("cell analysis instrument accessories and reagents", 5.0),
    ],
    "ca instrument calibration tools and standards": [
        ("cell analysis instrument calibration tools and standards", 5.0),
    ],
    "ca microplate readers and washers": [
        ("instrument microplate readers and washers", 5.0),
        ("cell analysis microplate readers", 4.6),
        ("cell analysis microplate washers", 4.6),
    ],
    "countess 1": [
        ("instrument countess", 2.8),
        ("countess 1", 5.0),
    ],
    "countess 2 and 2fl": [
        ("instrument countess", 2.8),
        ("countess 2 and 2fl", 5.0),
    ],
    "countess 3 and 3fl": [
        ("instrument countess", 2.8),
        ("countess 3 and 3fl", 5.0),
    ],
    "custom dna oligos": [
        ("custom dna oligos", 4.5),
        ("oligo files", 4.0),
        ("randall primers", 3.2),
        ("primers", 2.2),
    ],
    "dissociation reagents": [
        ("cell culture reagents antibiotics and supplements", 2.8),
        ("cell isolation and dissociation reagents", 4.2),
    ],
    "dpm": [
        ("cell culture media", 2.2),
        ("dpm media", 5.0),
    ],
    "dynabeads": [
        ("dynabeads", 4.8),
        ("beadretriever", 4.0),
    ],
    "electrophoresis reagents and kits": [
        ("electrophoresis", 4.2),
        ("e gel", 4.0),
        ("e base", 4.0),
        ("safe imager", 3.8),
        ("ibases", 3.0),
        ("ibase", 3.8),
    ],
    "ecm and 3d culture": [
        ("extracellular matrices and 3d cultures", 5.0),
    ],
    "evos floid": [
        ("instrument evos", 2.8),
        ("cell analysis instrument evos floid", 5.0),
        ("evos floid", 4.6),
    ],
    "evos 5000": [
        ("instrument evos", 2.8),
        ("cell analysis instrument evos m5000", 5.0),
        ("evos m5000", 4.6),
    ],
    "evos1000": [
        ("instrument evos", 2.8),
        ("cell analysis instrument evos s1000", 5.0),
        ("evos s1000", 4.6),
    ],
    "gibco manufacturing and packaging docs": [
        ("gibco manufacturing and packaging docs", 5.0),
    ],
    "hcs": [
        ("instrument hcs", 5.0),
        ("hcs", 4.6),
    ],
    "liquid cell culture": [
        ("cell culture media", 2.2),
        ("liquid media", 5.0),
    ],
    "magmax and kingfisher": [
        ("magmax & kingfisher", 4.6),
        ("magmax and kingfisher", 4.6),
        ("magmax", 4.0),
        ("kingfisher", 4.0),
    ],
    "mol bio sample prep": [
        ("sample prep", 4.2),
        ("benchpro", 4.4),
        ("prepstation", 4.2),
        ("6100", 3.4),
        ("6700", 3.0),
        ("tempus", 3.2),
    ],
    "nanodrop": [
        ("nanodrop", 5.0),
        ("nanodrop eight", 4.0),
        ("nanodrop one onec", 4.0),
        ("nanodrop onec", 3.8),
    ],
    "nutritional supplements and other reagents": [
        ("cell culture reagents antibiotics and supplements", 2.8),
        ("nutritional supplements", 5.0),
    ],
    "protein expression": [
        ("protein expression", 4.8),
        ("algae", 2.6),
        ("bluegrass", 4.0),
    ],
    "protein affinity purification": [
        ("protein affinity purification", 5.0),
    ],
    "protein assays": [
        ("protein assays", 5.0),
        ("protein assay", 4.6),
    ],
    "western blotting": [
        ("western blotting", 5.0),
        ("western blot", 4.8),
    ],
    "antibodies - proteins": [
        ("antibodies and immunoassays antibodies", 5.0),
        ("antibodies and immunoassays antibodies primary", 5.0),
        ("primary antibodies education and training", 4.4),
        ("primary antibodies internal faqs", 4.2),
        ("secondary antibodies documentation confidential", 4.2),
        ("antibodies and recombinant proteins", 2.8),
    ],
    "pepro gmp recombinant proteins": [
        ("pepro gmp recombinant proteins", 5.0),
    ],
    "recombinant proteins": [
        ("recombinant proteins", 4.8),
        ("antibodies and recombinant proteins", 2.4),
    ],
    "epigenetics chip emsa and remsa": [
        ("epigenetics", 5.0),
    ],
    "ibright and my ecl": [
        ("ibright and myecl", 5.0),
        ("ibright", 3.8),
        ("myecl", 3.8),
    ],
    "bait ip": [
        ("education and training - bait", 5.0),
        ("bait", 2.6),
    ],
    "target ip": [
        ("application and product notes - target", 4.8),
        ("documentation - target", 4.8),
        ("education and training - target", 4.8),
        ("troubleshooting - target", 5.0),
        ("target", 2.6),
    ],
    "ms sample prep": [
        ("mass spectrometry reagents", 2.8),
        ("ms sample prep", 5.0),
    ],
    "ms sample quantitation": [
        ("mass spectrometry reagents", 2.8),
        ("ms sample quantitation", 5.0),
    ],
    "ms standards": [
        ("mass spectrometry reagents", 2.8),
        ("ms standards", 5.0),
    ],
    "beads - proteins": [
        ("magnetic beads", 5.0),
    ],
    "protein labeling xlinking and modification": [
        ("protein labeling xlinking and modification", 5.0),
        ("protein labeling xlinking modification", 5.0),
    ],
    "singleplex and multiplex protein immunoassays": [
        ("singleplex and multiplex immunoassays", 5.0),
        ("singleplex and multiplex protein immunoassays", 5.0),
    ],
    "qubit and quant-it": [
        ("qubit and quant it", 4.6),
        ("qubit and quant-it", 4.6),
        ("qubit", 4.0),
    ],
    "spectra": [
        ("ca reagents kits spectra", 5.0),
        ("ca spectra", 4.6),
        ("spectra", 4.0),
    ],
    "microspheres": [
        ("ca polystyrene microspheres", 5.0),
        ("microspheres", 4.4),
    ],
    "thermal cycler plastics and reagents": [
        ("pcr plastics", 4.8),
        ("microamp", 4.0),
        ("optical adhesive", 3.4),
    ],
    "tali": [
        ("instrument tali", 5.0),
        ("tali", 4.6),
    ],
    "transfection": [
        ("transfection", 4.4),
        ("neon", 4.0),
        ("xenon", 4.0),
        ("invivofectamine reagents", 3.8),
        ("lipid transfection reagents", 3.4),
    ],
}

CATEGORY_REQUIRED_PATH_COMPONENTS: Dict[str, Tuple[str, ...]] = {
    "water": ("water",),
}

ARTICLE_TYPE_PRIMARY_RULES: List[Tuple[str, Tuple[str, ...]]] = [
    ("FAQs", ("faq", "frequently asked questions")),
    (
        "Manuals and Guides",
        ("user guide", "manual", "protocol", "instructions for use", "instruction manual", "quick reference", "release notes"),
    ),
    (
        "Application and Product Notes",
        ("application note", "app note", "product note", "white paper", "brochure", "poster", "flyer"),
    ),
    (
        "Education & Training",
        ("training", "tutorial", "course", "workshop", "webinar", "slides", "slide deck", "deck"),
    ),
    (
        "References",
        ("certificate", "coa", "coc", "declaration", "specification", "datasheet", "compliance", "statement", "memo"),
    ),
    (
        "Troubleshooting",
        ("troubleshooting", "troubleshoot", "error", "failure", "issue", "problem", "fix", "resolve"),
    ),
]

ARTICLE_TYPE_SECONDARY_RULES: List[Tuple[str, Tuple[str, ...]]] = [
    ("FAQs", ("faq", "frequently asked questions")),
    ("Application and Product Notes", ("application note", "app note", "product note", "white paper")),
    ("Manuals and Guides", ("user guide", "manual", "protocol", "instructions for use", "quick reference")),
    ("Education & Training", ("training", "tutorial", "course", "workshop", "webinar", "slides", "deck")),
    ("References", ("certificate", "coa", "coc", "declaration", "specification", "datasheet", "compliance")),
]

ARTICLE_TYPE_TROUBLESHOOTING_SECONDARY_TERMS = (
    "troubleshooting",
    "troubleshoot",
    "error",
    "failure",
    "issue",
    "problem",
    "fix",
    "resolve",
)

NEAR_DUP_FILE_TOKEN_STOPWORDS = {
    "final",
    "draft",
    "copy",
    "version",
    "updated",
    "update",
    "notes",
    "document",
    "presentation",
    "sample",
    "guide",
    "manual",
}

USAGE_LOCK = threading.Lock()
USAGE: Dict[str, int] = {"chat_in": 0, "chat_out": 0, "embed_in": 0}

THEME = {
    # One Half Light-inspired
    "bg": "#FAFAFA",
    "panel": "#FFFFFF",
    "text_bg": "#FFFFFF",
    "fg": "#383A42",
    "muted": "#A0A1A7",
    "accent": "#61AFEF",
    "accent_dark": "#4A9FE0",
    "border": "#E5E5E5",
    "btn_bg": "#EFEFEF",
    "btn_fg": "#383A42",
}

FONT_FAMILY = "Cascadia Mono"
FONT_BASE_SIZE = 11
FONT_BASE = (FONT_FAMILY, FONT_BASE_SIZE)
FONT_SMALL = (FONT_FAMILY, FONT_BASE_SIZE - 1)
FONT_LABEL = (FONT_FAMILY, FONT_BASE_SIZE)
FONT_HEADER = (FONT_FAMILY, FONT_BASE_SIZE + 2, "bold")
FONT_BUTTON = (FONT_FAMILY, FONT_BASE_SIZE)


@dataclass
class AzureConfig:
    api_key: str
    chat_api_key: str
    embeddings_api_key: str
    api_version: str
    api_key_header: str
    chat_base_url: str
    chat_path: str
    chat_deployment: str
    embeddings_base_url: str
    embeddings_path: str
    embeddings_deployment: str
    include_model_in_body: bool


@dataclass
class DocRecord:
    doc_id: str
    file_key: str
    file_name: str
    file_path: str
    source_path: str
    file_ext: str
    category: str
    tags: List[str]
    short_summary: str
    long_summary: str
    word_count: int
    char_count: int
    extraction_status: str
    review_flags: str
    duplicate_of: str
    duplicate_score: Optional[float]
    duplicate_group_id: str
    near_duplicate_of: str
    near_duplicate_score: Optional[float]
    near_duplicate_group_id: str
    review_group_id: str
    duplicate_relation_type: str
    moved_to: str


@dataclass
class ArticleRecord:
    doc_id: str
    file_key: str
    file_name: str
    file_path: str
    article_index: int
    article_title: str
    article_summary: str
    duplicate_of: str
    duplicate_score: Optional[float]
    duplicate_group_id: str


@dataclass
class InputFile:
    source_path: Path
    display_path: str
    file_key: str


@dataclass
class UnsupportedFileRecord:
    file_name: str
    file_path: str
    file_type: str
    source_kind: str


def env(name: str, default: Optional[str] = None) -> Optional[str]:
    val = os.getenv(name)
    return val if val not in (None, "") else default


def reset_usage() -> None:
    with USAGE_LOCK:
        USAGE["chat_in"] = 0
        USAGE["chat_out"] = 0
        USAGE["embed_in"] = 0


def add_chat_usage(in_chars: int, out_chars: int) -> None:
    with USAGE_LOCK:
        USAGE["chat_in"] += in_chars
        USAGE["chat_out"] += out_chars


def add_embed_usage(in_chars: int) -> None:
    with USAGE_LOCK:
        USAGE["embed_in"] += in_chars


def get_usage() -> Dict[str, int]:
    with USAGE_LOCK:
        return dict(USAGE)


def resolve_embeddings_source(value: Optional[str]) -> str:
    if value in ("summary", "full_text", EMBEDDINGS_SOURCE_NONE):
        return value
    return DEFAULT_EMBEDDINGS_SOURCE


def apply_theme(root: tk.Tk) -> None:
    root.configure(bg=THEME["bg"])
    try:
        root.option_add("*Font", FONT_BASE)
    except Exception:
        pass
    if ttk is not None:
        try:
            style = ttk.Style(root)
            style.theme_use("clam")
            style.configure(
                "TProgressbar",
                troughcolor=THEME["panel"],
                background=THEME["accent"],
                bordercolor=THEME["border"],
                lightcolor=THEME["accent"],
                darkcolor=THEME["accent"],
            )
        except Exception:
            pass


def load_app_config(path: Path) -> Dict[str, List[str]]:
    if not path.exists():
        return {}
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return {}
    apps = data.get("applications", {})
    if isinstance(apps, dict):
        cleaned: Dict[str, List[str]] = {}
        for k, v in apps.items():
            if isinstance(v, list):
                cleaned[k] = [str(x).strip() for x in v if str(x).strip()]
        return cleaned
    return {}


def save_app_config(path: Path, apps: Dict[str, List[str]]) -> None:
    data = {"applications": apps}
    path.write_text(json.dumps(data, indent=2), encoding="utf-8")


def check_ocr_dependencies(ocrmypdf_enabled: bool) -> List[str]:
    missing: List[str] = []
    if ocrmypdf_enabled:
        if shutil.which("tesseract") is None:
            missing.append("tesseract")
        if shutil.which("qpdf") is None:
            missing.append("qpdf")
        # Ghostscript executable name varies on Windows
        if shutil.which("gswin64c") is None and shutil.which("gswin32c") is None and shutil.which("gs") is None:
            missing.append("ghostscript")
    # Poppler tools for pdf2image fallback
    if shutil.which("pdftoppm") is None:
        missing.append("poppler(pdftoppm)")
    return missing


def warn_missing_ocr_deps(ocrmypdf_enabled: bool) -> List[str]:
    missing = check_ocr_dependencies(ocrmypdf_enabled)
    if not missing:
        return []
    msg = "Missing OCR dependencies (OCR may be limited): " + ", ".join(missing)
    logging.warning(msg)
    if tk is not None and messagebox is not None:
        messagebox.showwarning("OCR Dependencies", msg)
    return missing


def sanitize_folder(name: str) -> str:
    name = (name or "").strip()
    if not name:
        return "uncategorized"
    name = "".join("_" if c in INVALID_WIN_CHARS else c for c in name)
    name = name.strip(" .")
    return name or "uncategorized"


def setup_logging(out_dir: Path) -> None:
    out_dir.mkdir(parents=True, exist_ok=True)
    log_path = out_dir / "docatlas.log"
    logging.basicConfig(
        level=logging.INFO,
        format="%(asctime)s [%(levelname)s] %(message)s",
        handlers=[logging.FileHandler(log_path, encoding="utf-8"), logging.StreamHandler(sys.stdout)],
        force=True,
    )


def build_url(base: str, path_template: str, deployment: str) -> str:
    if path_template.startswith("http://") or path_template.startswith("https://"):
        return path_template.format(deployment=deployment)
    base = base.rstrip("/")
    path = path_template.format(deployment=deployment)
    return f"{base}{path}"


def azure_config_from_env(require_key: bool = True) -> AzureConfig:
    api_key = env("AZURE_OPENAI_API_KEY", "")
    if require_key and not api_key and not env("AZURE_CHAT_API_KEY") and not env("AZURE_EMBEDDINGS_API_KEY"):
        raise ValueError("AZURE_OPENAI_API_KEY is not set")

    return AzureConfig(
        api_key=api_key,
        chat_api_key=env("AZURE_CHAT_API_KEY", api_key),
        embeddings_api_key=env("AZURE_EMBEDDINGS_API_KEY", api_key),
        api_version=env("AZURE_OPENAI_API_VERSION", DEFAULT_API_VERSION),
        api_key_header=env("AZURE_OPENAI_API_KEY_HEADER", DEFAULT_API_KEY_HEADER),
        chat_base_url=env("AZURE_CHAT_BASE_URL", DEFAULT_CHAT_BASE_URL),
        chat_path=env("AZURE_CHAT_PATH", DEFAULT_CHAT_PATH),
        chat_deployment=env("AZURE_CHAT_DEPLOYMENT", DEFAULT_CHAT_DEPLOYMENT),
        embeddings_base_url=env("AZURE_EMBEDDINGS_BASE_URL", DEFAULT_EMBEDDINGS_BASE_URL),
        embeddings_path=env("AZURE_EMBEDDINGS_PATH", DEFAULT_EMBEDDINGS_PATH),
        embeddings_deployment=env("AZURE_EMBEDDINGS_DEPLOYMENT", DEFAULT_EMBEDDINGS_DEPLOYMENT),
        include_model_in_body=(env("AZURE_INCLUDE_MODEL_IN_BODY", "1") == "1"),
    )


def api_delay_sec() -> float:
    val = env("DOCATLAS_API_DELAY", "")
    try:
        if val:
            return max(0.0, float(val))
    except Exception:
        pass
    return DEFAULT_API_DELAY_SEC


def api_max_retries() -> int:
    val = env("DOCATLAS_API_MAX_RETRIES", "")
    try:
        if val:
            return max(1, int(val))
    except Exception:
        pass
    return DEFAULT_API_MAX_RETRIES


def api_retry_base_sec() -> float:
    val = env("DOCATLAS_API_RETRY_BASE", "")
    try:
        if val:
            return max(0.1, float(val))
    except Exception:
        pass
    return DEFAULT_API_RETRY_BASE_SEC


def api_retry_max_sec() -> float:
    val = env("DOCATLAS_API_RETRY_MAX", "")
    try:
        if val:
            return max(1.0, float(val))
    except Exception:
        pass
    return DEFAULT_API_RETRY_MAX_SEC


def api_timeout_sec() -> float:
    val = env("DOCATLAS_API_TIMEOUT", "")
    try:
        if val:
            return max(5.0, float(val))
    except Exception:
        pass
    return float(DEFAULT_API_TIMEOUT_SEC)


def _retry_sleep_seconds(attempt: int) -> float:
    base = api_retry_base_sec()
    cap = api_retry_max_sec()
    backoff = min(cap, base * (2 ** max(0, attempt)))
    # Add jitter to avoid synchronized retries.
    jitter = random.uniform(0.0, min(2.0, 0.2 * backoff))
    return backoff + jitter


def _is_retryable_request_error(exc: Exception) -> bool:
    return isinstance(exc, (requests.exceptions.Timeout, requests.exceptions.ConnectionError))


def _post_with_retries(
    op_name: str,
    url: str,
    headers: Dict[str, str],
    params: Dict[str, Any],
    payload: Dict[str, Any],
) -> requests.Response:
    delay = api_delay_sec()
    retries = api_max_retries()
    timeout = api_timeout_sec()
    transient_statuses = {408, 425, 429, 500, 502, 503, 504}
    last_exc: Optional[Exception] = None

    for attempt in range(retries):
        if delay:
            time.sleep(delay)
        try:
            resp = requests.post(url, headers=headers, params=params, json=payload, timeout=timeout)
        except Exception as exc:
            if _is_retryable_request_error(exc) and attempt < retries - 1:
                wait_sec = _retry_sleep_seconds(attempt)
                logging.warning(
                    "%s request transient error on attempt %d/%d: %s; retrying in %.1fs",
                    op_name,
                    attempt + 1,
                    retries,
                    exc,
                    wait_sec,
                )
                time.sleep(wait_sec)
                last_exc = exc
                continue
            raise

        if resp.status_code in transient_statuses and attempt < retries - 1:
            wait_sec = _retry_sleep_seconds(attempt)
            logging.warning(
                "%s request transient status %s on attempt %d/%d; retrying in %.1fs",
                op_name,
                resp.status_code,
                attempt + 1,
                retries,
                wait_sec,
            )
            time.sleep(wait_sec)
            continue
        return resp

    if last_exc is not None:
        raise RuntimeError(f"{op_name} request failed after retries: {last_exc}")
    raise RuntimeError(f"{op_name} request failed after retries")


def file_key(path: Path, display_path: Optional[str] = None) -> str:
    st = path.stat()
    key_path = display_path or str(path)
    return f"{key_path}|{st.st_mtime_ns}|{st.st_size}"


def zip_member_key(zip_display_path: str, member_display_path: str, zip_path: Path, info: zipfile.ZipInfo) -> str:
    zip_stat = zip_path.stat()
    return (
        f"{member_display_path}|"
        f"zip={zip_display_path}|"
        f"zip_mtime={zip_stat.st_mtime_ns}|"
        f"zip_size={zip_stat.st_size}|"
        f"entry_crc={int(info.CRC)}|"
        f"entry_size={int(info.file_size)}"
    )


def sanitize_zip_member_segment(value: str) -> str:
    clean = "".join("_" if c in INVALID_WIN_CHARS else c for c in (value or ""))
    clean = clean.strip(" .")
    return clean or "_"


def build_safe_zip_target(base_dir: Path, member_name: str, used_targets: set[str]) -> Optional[Path]:
    raw_parts = [p for p in str(member_name).replace("\\", "/").split("/") if p not in ("", ".")]
    if not raw_parts:
        return None
    if any(part == ".." for part in raw_parts):
        raise ValueError(f"Zip member escapes extraction root: {member_name}")

    safe_parts = [sanitize_zip_member_segment(part) for part in raw_parts]
    target = base_dir.joinpath(*safe_parts)
    try:
        resolved_base = base_dir.resolve()
        resolved_target = target.resolve(strict=False)
    except Exception:
        resolved_base = base_dir.absolute()
        resolved_target = target.absolute()
    if not str(resolved_target).startswith(str(resolved_base)):
        raise ValueError(f"Zip member escapes extraction root: {member_name}")

    candidate = target
    if candidate.exists() or str(candidate).lower() in used_targets:
        stem = candidate.stem
        suffix = candidate.suffix
        idx = 1
        while True:
            alt = candidate.with_name(f"{stem}_{idx}{suffix}")
            if not alt.exists() and str(alt).lower() not in used_targets:
                candidate = alt
                break
            idx += 1

    used_targets.add(str(candidate).lower())
    return candidate


def _display_file_type(path_value: str) -> str:
    ext = Path(path_value or "").suffix.lower()
    return ext if ext else "[no extension]"


def extract_zip_inputs(
    zip_path: Path,
    zip_display_path: str,
    staging_root: Path,
) -> Tuple[List[InputFile], List[Tuple[Path, str]], List[UnsupportedFileRecord]]:
    extracted_files: List[InputFile] = []
    nested_zips: List[Tuple[Path, str]] = []
    unsupported_files: List[UnsupportedFileRecord] = []
    archive_slug = hashlib.sha1(zip_display_path.encode("utf-8", errors="ignore")).hexdigest()[:12]
    extract_dir = staging_root / archive_slug
    extract_dir.mkdir(parents=True, exist_ok=True)
    used_targets: set[str] = set()

    with zipfile.ZipFile(zip_path) as zf:
        for info in zf.infolist():
            if info.is_dir():
                continue
            target = build_safe_zip_target(extract_dir, info.filename, used_targets)
            if target is None:
                continue
            target.parent.mkdir(parents=True, exist_ok=True)
            with zf.open(info) as src, target.open("wb") as dst:
                shutil.copyfileobj(src, dst)

            try:
                member_dt = time.mktime(info.date_time + (0, 0, -1))
                os.utime(target, (member_dt, member_dt))
            except Exception:
                pass

            member_rel = str(info.filename).replace("\\", "/").lstrip("/")
            member_display_path = f"{zip_display_path}!/{member_rel}"
            ext = target.suffix.lower()
            if ext in SUPPORTED_EXTS:
                extracted_files.append(
                    InputFile(
                        source_path=target,
                        display_path=member_display_path,
                        file_key=zip_member_key(zip_display_path, member_display_path, zip_path, info),
                    )
                )
            elif ext == ".zip":
                nested_zips.append((target, member_display_path))
            else:
                unsupported_files.append(
                    UnsupportedFileRecord(
                        file_name=Path(member_rel).name,
                        file_path=member_display_path,
                        file_type=_display_file_type(member_rel),
                        source_kind="zip_member",
                    )
                )

    return extracted_files, nested_zips, unsupported_files


def relative_input_path(input_dir: Path, path: Path) -> str:
    try:
        rel = path.relative_to(input_dir)
    except Exception:
        rel = Path(path.name)
    return str(rel).replace("\\", "/")


def normalize_compare_path(value: str) -> str:
    return str(value or "").strip().replace("\\", "/").casefold()


def path_compare_aliases(value: str) -> set[str]:
    norm = normalize_compare_path(value)
    if not norm:
        return set()
    parts = [p for p in norm.split("/") if p]
    aliases = {norm}
    for idx in range(len(parts)):
        aliases.add("/".join(parts[idx:]))
    return aliases


def list_files(input_dir: Path) -> Tuple[List[InputFile], List[UnsupportedFileRecord], Optional[tempfile.TemporaryDirectory[str]]]:
    files: List[InputFile] = []
    zip_files: List[Tuple[Path, str]] = []
    unsupported_files: List[UnsupportedFileRecord] = []
    for p in input_dir.rglob("*"):
        if not p.is_file():
            continue
        ext = p.suffix.lower()
        display_path = relative_input_path(input_dir, p)
        if ext in SUPPORTED_EXTS:
            files.append(InputFile(source_path=p, display_path=display_path, file_key=file_key(p, display_path)))
        elif ext == ".zip":
            zip_files.append((p, display_path))
        else:
            unsupported_files.append(
                UnsupportedFileRecord(
                    file_name=p.name,
                    file_path=display_path,
                    file_type=_display_file_type(display_path),
                    source_kind="file",
                )
            )

    temp_dir: Optional[tempfile.TemporaryDirectory[str]] = None
    if zip_files:
        temp_dir = tempfile.TemporaryDirectory(prefix="docatlas_zip_stage_")
        staging_root = Path(temp_dir.name)
        pending = list(zip_files)
        while pending:
            zip_path, zip_display_path = pending.pop(0)
            try:
                staged_files, nested_zips, zip_unsupported = extract_zip_inputs(zip_path, zip_display_path, staging_root)
            except zipfile.BadZipFile:
                logging.warning("Skipping invalid zip archive: %s", zip_display_path)
                unsupported_files.append(
                    UnsupportedFileRecord(
                        file_name=Path(zip_display_path).name,
                        file_path=zip_display_path,
                        file_type=".zip",
                        source_kind="invalid_zip",
                    )
                )
                continue
            except Exception as exc:
                logging.warning("Failed to extract zip archive %s: %s", zip_display_path, exc)
                unsupported_files.append(
                    UnsupportedFileRecord(
                        file_name=Path(zip_display_path).name,
                        file_path=zip_display_path,
                        file_type=".zip",
                        source_kind="invalid_zip",
                    )
                )
                continue
            files.extend(staged_files)
            pending.extend(nested_zips)
            unsupported_files.extend(zip_unsupported)

    files.sort(key=lambda item: item.display_path.lower())
    unsupported_files.sort(key=lambda item: (item.file_type, item.file_path.lower()))
    return files, unsupported_files, temp_dir


def scan_input_stats(files: List[InputFile]) -> Dict[str, Any]:
    total_size = 0
    by_ext: Dict[str, int] = {}
    for item in files:
        p = item.source_path
        try:
            total_size += p.stat().st_size
        except Exception:
            pass
        ext = p.suffix.lower()
        by_ext[ext] = by_ext.get(ext, 0) + 1
    total_size_mb = total_size / (1024 * 1024)
    return {"count": len(files), "total_size_mb": total_size_mb, "by_ext": by_ext}


def unsupported_file_stats(items: List[UnsupportedFileRecord]) -> Dict[str, Any]:
    by_type: Dict[str, int] = {}
    by_source_kind: Dict[str, int] = {}
    by_source_folder: Dict[str, int] = {}
    for item in items:
        by_type[item.file_type] = by_type.get(item.file_type, 0) + 1
        by_source_kind[item.source_kind] = by_source_kind.get(item.source_kind, 0) + 1
        folder = unsupported_source_folder(str(item.file_path or ""))
        by_source_folder[folder] = by_source_folder.get(folder, 0) + 1
    return {
        "count": len(items),
        "by_type": by_type,
        "by_source_kind": by_source_kind,
        "by_source_folder": by_source_folder,
    }


def unsupported_source_folder(file_path: str) -> str:
    normalized = str(file_path or "").replace("\\", "/").rstrip("/")
    if "!/" in normalized:
        archive, member = normalized.split("!/", 1)
        member_parent = Path(member).parent.as_posix()
        folder = f"{archive}!/{member_parent}" if member_parent not in ("", ".") else f"{archive}!/[root]"
    else:
        folder = str(Path(normalized).parent).replace("\\", "/")
    return folder if folder not in ("", ".") else "[root]"


def write_unsupported_files_report(out_dir: Path, items: List[UnsupportedFileRecord]) -> Path:
    report_path = out_dir / "unsupported_files_report.txt"
    stats = unsupported_file_stats(items)
    lines = [
        "Unsupported Files Report",
        "========================",
        f"Total unsupported files: {stats['count']}",
        "",
        "Unsupported by Datatype:",
    ]
    for file_type in sorted(stats["by_type"].keys()):
        lines.append(f"- {file_type}: {stats['by_type'][file_type]}")
    lines.extend(["", "Unsupported by Source Kind:"])
    for source_kind in sorted(stats["by_source_kind"].keys()):
        lines.append(f"- {source_kind}: {stats['by_source_kind'][source_kind]}")
    if stats["by_source_folder"]:
        lines.extend(["", "Top Source Folders:"])
        for folder, count in sorted(stats["by_source_folder"].items(), key=lambda kv: (-kv[1], kv[0]))[:10]:
            lines.append(f"- {folder}: {count}")
    lines.extend(["", "Detailed List:"])
    for item in items:
        lines.append(f"- {item.file_type} | {item.file_name} | {item.file_path}")
    report_path.write_text("\n".join(lines), encoding="utf-8")
    return report_path


def write_unsupported_cleanup_workbook(out_dir: Path, items: List[UnsupportedFileRecord]) -> Path:
    workbook_path = out_dir / "unsupported_cleanup.xlsx"
    stats = unsupported_file_stats(items)
    by_source_folder = stats["by_source_folder"]

    cleanup_rows: List[Dict[str, Any]] = []
    for item in items:
        source_folder = unsupported_source_folder(item.file_path)
        cleanup_rows.append(
            {
                "FileType": item.file_type,
                "FileName": item.file_name,
                "FilePath": item.file_path,
                "SourceKind": item.source_kind,
                "SourceFolder": source_folder,
                "RecommendedAction": "Review",
                "DeleteCandidate": False,
                "Decision": "",
                "DecisionNotes": "",
                "ReviewedBy": "",
                "FinalDisposition": "",
            }
        )

    cleanup_rows.sort(
        key=lambda row: (
            -int(by_source_folder.get(str(row["SourceFolder"]), 0)),
            str(row["SourceFolder"]).lower(),
            str(row["FileType"]).lower(),
            str(row["FilePath"]).lower(),
        )
    )

    queue_columns = [
        "FileType",
        "FileName",
        "FilePath",
        "SourceKind",
        "SourceFolder",
        "RecommendedAction",
        "DeleteCandidate",
        "Decision",
        "DecisionNotes",
        "ReviewedBy",
        "FinalDisposition",
    ]
    queue_df = sanitize_excel_df(pd.DataFrame(cleanup_rows, columns=queue_columns))
    legend_rows = [
        {
            "Decision": "Keep",
            "Meaning": "Keep the unsupported file in place and do not delete it.",
            "Use": "Use when the file should remain at source or needs no action.",
            "Color": "Gray",
        },
        {
            "Decision": "Delete at Source",
            "Meaning": "Delete manually at the source after review approval.",
            "Use": "Use only after a human confirms the file is safe to remove.",
            "Color": "Red",
        },
        {
            "Decision": "Ignore",
            "Meaning": "Ignore the file for now and leave it out of the migration scope.",
            "Use": "Use for low-priority items that should not block migration.",
            "Color": "Yellow",
        },
        {
            "Decision": "Needs Follow-up",
            "Meaning": "Escalate for review because more context is needed.",
            "Use": "Use when ownership, scope, or disposition is unclear.",
            "Color": "Blue",
        },
    ]
    legend_df = pd.DataFrame(legend_rows, columns=["Decision", "Meaning", "Use", "Color"])

    with pd.ExcelWriter(workbook_path, engine="openpyxl") as writer:
        queue_df.to_excel(writer, index=False, sheet_name="cleanup_queue")
        legend_df.to_excel(writer, index=False, sheet_name="cleanup_legend")

    try:
        from openpyxl import load_workbook
        from openpyxl.styles import Alignment, Font, PatternFill
        from openpyxl.formatting.rule import FormulaRule
        from openpyxl.utils import get_column_letter
        from openpyxl.worksheet.datavalidation import DataValidation

        wb = load_workbook(workbook_path)
        queue_ws = wb["cleanup_queue"]
        legend_ws = wb["cleanup_legend"]

        queue_ws.freeze_panes = "A2"
        queue_ws.auto_filter.ref = f"A1:{get_column_letter(queue_ws.max_column)}{queue_ws.max_row}"
        queue_ws.row_dimensions[1].height = 28

        header_fill = PatternFill(fill_type="solid", fgColor="1F4E78")
        header_font = Font(color="FFFFFF", bold=True)
        header_alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        wrap_alignment = Alignment(wrap_text=True, vertical="top")
        top_alignment = Alignment(vertical="top")
        widths = {
            "A": 14,
            "B": 32,
            "C": 70,
            "D": 16,
            "E": 52,
            "F": 20,
            "G": 16,
            "H": 18,
            "I": 42,
            "J": 18,
            "K": 22,
        }
        wrap_cols = {"B", "C", "E", "F", "I", "K"}

        for cell in queue_ws[1]:
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = header_alignment
        for col_letter, width in widths.items():
            queue_ws.column_dimensions[col_letter].width = width
        for row in range(2, queue_ws.max_row + 1):
            for col in range(1, queue_ws.max_column + 1):
                col_letter = get_column_letter(col)
                queue_ws.cell(row=row, column=col).alignment = wrap_alignment if col_letter in wrap_cols else top_alignment

        decision_validation = DataValidation(
            type="list",
            formula1='"Keep,Delete at Source,Ignore,Needs Follow-up"',
            allow_blank=True,
        )
        decision_validation.prompt = "Choose the reviewed disposition for this unsupported file."
        decision_validation.error = "Select one of the allowed decision values."
        queue_ws.add_data_validation(decision_validation)
        decision_validation.add(f"H2:H{queue_ws.max_row}")

        queue_range = f"A2:{get_column_letter(queue_ws.max_column)}{queue_ws.max_row}"
        for formula, color in [
            ('=$H2="Keep"', "E7E6E6"),
            ('=$H2="Delete at Source"', "FDE9D9"),
            ('=$H2="Ignore"', "FFF2CC"),
            ('=$H2="Needs Follow-up"', "D9EAF7"),
        ]:
            queue_ws.conditional_formatting.add(
                queue_range,
                FormulaRule(formula=[formula], fill=PatternFill(fill_type="solid", fgColor=color)),
            )

        legend_ws.freeze_panes = "A2"
        legend_ws.auto_filter.ref = f"A1:{get_column_letter(legend_ws.max_column)}{legend_ws.max_row}"
        legend_widths = {"A": 20, "B": 44, "C": 44, "D": 12}
        for cell in legend_ws[1]:
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = header_alignment
        for col_letter, width in legend_widths.items():
            legend_ws.column_dimensions[col_letter].width = width
        for row in range(2, legend_ws.max_row + 1):
            for col in range(1, legend_ws.max_column + 1):
                legend_ws.cell(row=row, column=col).alignment = wrap_alignment
        legend_color_map = {
            "Keep": "E7E6E6",
            "Delete at Source": "FDE9D9",
            "Ignore": "FFF2CC",
            "Needs Follow-up": "D9EAF7",
        }
        for row in range(2, legend_ws.max_row + 1):
            decision = str(legend_ws.cell(row=row, column=1).value or "")
            color = legend_color_map.get(decision)
            if color:
                legend_ws.cell(row=row, column=1).fill = PatternFill(fill_type="solid", fgColor=color)

        wb.save(workbook_path)
    except Exception:
        pass

    return workbook_path


def load_last_run_stats(out_dir: Path) -> Optional[Dict[str, Any]]:
    stats_path = out_dir / LAST_RUN_STATS_FILENAME
    if not stats_path.exists():
        return None
    try:
        with stats_path.open("r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return None


def save_last_run_stats(out_dir: Path, stats: Dict[str, Any]) -> None:
    stats_path = out_dir / LAST_RUN_STATS_FILENAME
    try:
        with stats_path.open("w", encoding="utf-8") as f:
            json.dump(stats, f, indent=2)
    except Exception:
        pass


def format_duration(seconds: float) -> str:
    seconds = max(0, int(seconds))
    minutes = seconds / 60
    hours = minutes / 60
    days = hours / 24
    if days >= 1:
        return f"{days:.1f} days"
    if hours >= 1:
        return f"{hours:.1f} hours"
    if minutes >= 1:
        return f"{minutes:.1f} minutes"
    return f"{seconds} seconds"


def quick_estimate_runtime(
    input_stats: Dict[str, Any],
    output_dir: Path,
    ocrmypdf_enabled: bool,
    embeddings_source: str,
    chat_deployment: str,
) -> Tuple[Optional[float], str, bool]:
    baseline = load_last_run_stats(output_dir)
    count = input_stats.get("count", 0)
    total_mb = input_stats.get("total_size_mb", 0.0)
    if count <= 0:
        return None, "none", False

    if baseline and baseline.get("processed_files", 0) > 0:
        sec_per_file = baseline["elapsed_sec"] / baseline["processed_files"]
        sec_per_mb = 0.0
        if baseline.get("total_size_mb", 0) > 0:
            sec_per_mb = baseline["elapsed_sec"] / baseline["total_size_mb"]
        est_by_file = sec_per_file * count
        est_by_mb = sec_per_mb * total_mb if sec_per_mb > 0 else 0.0
        est_sec = max(est_by_file, est_by_mb)
        settings_match = (
            baseline.get("ocr_enabled") == ocrmypdf_enabled
            and baseline.get("embeddings_source") == embeddings_source
            and baseline.get("chat_deployment") == chat_deployment
        )
        logging.info(
            "Quick estimate (baseline): files=%d, size=%.1f MB -> ~%ds",
            count,
            total_mb,
            int(est_sec),
        )
        if not settings_match:
            logging.info("Quick estimate note: baseline settings differ from current run.")
        return est_sec, "baseline", settings_match

    est_by_file = DEFAULT_ESTIMATE_SEC_PER_FILE * count
    est_by_mb = DEFAULT_ESTIMATE_SEC_PER_MB * total_mb
    est_sec = max(est_by_file, est_by_mb)
    logging.info(
        "Quick estimate (heuristic): files=%d, size=%.1f MB -> ~%ds",
        count,
        total_mb,
        int(est_sec),
    )
    return est_sec, "heuristic", True


def normalize_text(text: str) -> str:
    text = re.sub(r"\s+", " ", text or "").strip().lower()
    return text


def hash_text(text: str) -> str:
    return hashlib.sha256(text.encode("utf-8", errors="ignore")).hexdigest()


def hash_file_bytes(path: Path, chunk_size: int = 1024 * 1024) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as fh:
        while True:
            chunk = fh.read(chunk_size)
            if not chunk:
                break
            digest.update(chunk)
    return digest.hexdigest()


def exact_hash_for_file(path: Path, fallback_key: str) -> str:
    try:
        return hash_file_bytes(path)
    except Exception as exc:
        logging.warning("Falling back to file-key exact hash for %s: %s", path, exc)
        return f"filekey:{hash_text(fallback_key)}"


def exact_hash_for_text(text: str, fallback_key: str) -> str:
    normalized = normalize_text(text)
    if normalized:
        return hash_text(normalized)
    return f"empty:{hash_text(fallback_key)}"


def embedding_text_for_doc(embeddings_source: str, normalized_text: str, long_summary: str, short_summary: str) -> str:
    if embeddings_source == "full_text":
        return normalized_text
    summary = long_summary or short_summary or ""
    return normalize_text(summary)


def embedding_text_for_article(embeddings_source: str, body: str, article_summary: str) -> str:
    if embeddings_source == "full_text":
        return normalize_text(body)
    summary = article_summary or ""
    return normalize_text(summary)


def min_embedding_chars_for_source(embeddings_source: str) -> int:
    return MIN_EMBEDDING_CHARS if embeddings_source == "full_text" else MIN_EMBEDDING_CHARS_SUMMARY


def split_text(text: str, max_chars: int) -> List[str]:
    text = text or ""
    if len(text) <= max_chars:
        return [text]
    chunks: List[str] = []
    start = 0
    while start < len(text):
        end = min(len(text), start + max_chars)
        chunk = text[start:end]
        chunks.append(chunk)
        start = end
    return chunks


def is_content_filter_error(exc: Exception) -> bool:
    msg = str(exc or "").lower()
    return (
        "content_filter" in msg
        or "responsibleaipolicyviolation" in msg
        or ("chat api error 400" in msg and "filtered" in msg)
    )


def _summary_tokens(text: str) -> List[str]:
    return [
        tok
        for tok in re.findall(r"[a-z0-9][a-z0-9\-]{2,}", normalize_text(text))
        if tok not in SUMMARY_STOPWORDS and not tok.isdigit()
    ]


def _split_into_sentences(text: str) -> List[str]:
    if not text:
        return []
    compact = re.sub(r"[ \t]+", " ", text)
    compact = re.sub(r"\n{2,}", "\n", compact)
    parts = re.split(r"(?<=[.!?])\s+|\n", compact)
    out: List[str] = []
    for part in parts:
        s = re.sub(r"\s+", " ", (part or "").strip())
        if not s:
            continue
        if len(s) > FALLBACK_MAX_SENTENCE_CHARS:
            s = s[:FALLBACK_MAX_SENTENCE_CHARS].rstrip()
        if len(s) < FALLBACK_MIN_SENTENCE_CHARS:
            continue
        out.append(s)
    return out


def _sentence_similarity(tokens_a: set[str], tokens_b: set[str]) -> float:
    if not tokens_a or not tokens_b:
        return 0.0
    inter = len(tokens_a.intersection(tokens_b))
    union = len(tokens_a.union(tokens_b))
    if union == 0:
        return 0.0
    return inter / union


def _extractive_summary_sentences(text: str, max_sentences: int, max_total_chars: int) -> List[str]:
    sentences = _split_into_sentences(text)
    if not sentences:
        fallback = re.sub(r"\s+", " ", (text or "").strip())
        if not fallback:
            return []
        return [fallback[:max_total_chars].rstrip()]

    freq: Dict[str, int] = {}
    sentence_tokens: List[set[str]] = []
    for s in sentences:
        toks = set(_summary_tokens(s))
        sentence_tokens.append(toks)
        for t in toks:
            freq[t] = freq.get(t, 0) + 1

    scored: List[Tuple[float, int, str]] = []
    total = len(sentences)
    for i, s in enumerate(sentences):
        toks = sentence_tokens[i]
        if toks:
            density = sum(freq.get(t, 0) for t in toks) / max(1.0, len(toks) ** 0.5)
        else:
            density = 0.0
        # Keep early context, but let high-information later lines win.
        position_bonus = max(0.0, 1.0 - (i / max(1, total)))
        score = density + (0.5 * position_bonus)
        scored.append((score, i, s))

    scored.sort(key=lambda x: (-x[0], x[1]))
    selected: List[Tuple[int, str]] = []
    selected_tokens: List[set[str]] = []
    char_count = 0
    for _score, i, s in scored:
        toks = sentence_tokens[i]
        if any(_sentence_similarity(toks, prev) > 0.72 for prev in selected_tokens):
            continue
        projected = char_count + len(s) + (1 if selected else 0)
        if projected > max_total_chars:
            continue
        selected.append((i, s))
        selected_tokens.append(toks)
        char_count = projected
        if len(selected) >= max_sentences:
            break

    if not selected:
        best = sorted(scored, key=lambda x: (-x[0], x[1]))[0][2]
        return [best[:max_total_chars].rstrip()]

    selected.sort(key=lambda x: x[0])
    return [s for _, s in selected]


def _extractive_summary_text(text: str, max_sentences: int, max_total_chars: int) -> str:
    return " ".join(_extractive_summary_sentences(text, max_sentences, max_total_chars)).strip()


def _extract_top_tags(text: str, max_tags: int = MAX_TAGS) -> List[str]:
    freq: Dict[str, int] = {}
    for tok in _summary_tokens(text):
        if len(tok) < 4:
            continue
        freq[tok] = freq.get(tok, 0) + 1
    if not freq:
        return []
    ranked = sorted(freq.items(), key=lambda kv: (-kv[1], kv[0]))
    return [t for t, _n in ranked[:max_tags]]


def _normalized_category_path_components(file_name: str = "", file_path: str = "") -> List[str]:
    raw_parts: List[str] = []
    if file_path:
        raw_parts.extend(re.split(r"[\\/!]+", str(file_path)))
    if file_name:
        raw_parts.append(Path(str(file_name)).stem)
    normalized: List[str] = []
    for part in raw_parts:
        part_norm = normalize_text(part)
        if part_norm:
            normalized.append(part_norm)
    return normalized


def _infer_category_from_text(
    text: str,
    categories: List[str],
    file_name: str = "",
    file_path: str = "",
) -> str:
    """
    Deterministic category resolver:
    - Scores each candidate category from text only (no model output dependency)
    - Uses phrase hints + category-token evidence
    - Uses stable tie-break by original category order
    """
    candidates = [c for c in categories if c and c not in ("Other", UNREADABLE_CATEGORY)]
    if not candidates:
        return "Other"

    text_norm = normalize_text(text)[:200000]
    if not text_norm:
        return "Other"

    text_tokens = _summary_tokens(text_norm)
    token_freq: Dict[str, int] = {}
    for tok in text_tokens:
        token_freq[tok] = token_freq.get(tok, 0) + 1

    path_components = _normalized_category_path_components(file_name=file_name, file_path=file_path)
    path_norm = " ".join(path_components)
    path_tokens = _summary_tokens(path_norm)
    path_token_freq: Dict[str, int] = {}
    for tok in path_tokens:
        path_token_freq[tok] = path_token_freq.get(tok, 0) + 1
    path_component_set = set(path_components)

    scored: List[Tuple[float, int, str]] = []
    for i, cat in enumerate(candidates):
        cat_norm = normalize_text(cat)
        if not cat_norm:
            continue

        score = 0.0
        if cat_norm in text_norm:
            score += 2.0

        # Category-name token signal (e.g., "gene", "expression", "genotyping").
        cat_tokens = [
            t
            for t in re.findall(r"[a-z0-9]+", cat_norm)
            if t not in SUMMARY_STOPWORDS and t not in CATEGORY_GENERIC_TOKENS
        ]
        for tok in cat_tokens:
            score += min(1.0, float(token_freq.get(tok, 0)) * 0.2)
            score += min(1.4, float(path_token_freq.get(tok, 0)) * 0.35)

        # Tuned phrase hints for known qPCR categories.
        for phrase, weight in CATEGORY_HINT_PHRASES.get(cat_norm, []):
            hit_count = text_norm.count(phrase)
            if hit_count > 0:
                score += min(3.0, float(hit_count)) * weight
            path_hit_count = path_norm.count(phrase)
            if path_hit_count > 0:
                score += min(2.0, float(path_hit_count)) * min(weight, 4.0) * 0.6

        for component, weight in CATEGORY_PATH_COMPONENT_HINTS.get(cat_norm, []):
            if component in path_component_set:
                score += weight

        required_phrases = CATEGORY_REQUIRED_PHRASES.get(cat_norm, ())
        if required_phrases:
            has_required_phrase = any(phrase in text_norm or phrase in path_norm for phrase in required_phrases)
            required_components = CATEGORY_REQUIRED_PATH_COMPONENTS.get(cat_norm, ())
            has_required_component = any(component in path_component_set for component in required_components)
            if cat_norm == "water":
                if not has_required_phrase and not has_required_component:
                    score = 0.0
            else:
                required_hits = max(token_freq.get(tok, 0) for tok in cat_tokens) if cat_tokens else 0
                min_token_hits = CATEGORY_REQUIRED_TOKEN_HITS.get(cat_norm, 1)
                if not has_required_phrase and required_hits < min_token_hits:
                    score = min(score, 0.25)

        scored.append((score, i, cat))

    if not scored:
        return "Other"

    scored.sort(key=lambda x: (-x[0], x[1]))
    best_score, _best_idx, best_cat = scored[0]
    # If no meaningful signal, keep deterministic safe fallback.
    if best_score < 1.0:
        return "Other"
    return best_cat


def build_fallback_document_summary(
    text: str,
    categories: List[str],
    file_name: str = "",
    file_path: str = "",
) -> Dict[str, Any]:
    long_summary = _extractive_summary_text(
        text,
        max_sentences=FALLBACK_DOC_LONG_SENTENCES,
        max_total_chars=FALLBACK_MAX_TOTAL_CHARS_DOC_LONG,
    )
    short_summary = _extractive_summary_text(
        text,
        max_sentences=FALLBACK_DOC_SHORT_SENTENCES,
        max_total_chars=FALLBACK_MAX_TOTAL_CHARS_DOC_SHORT,
    )
    if not short_summary:
        short_summary = long_summary[:FALLBACK_MAX_TOTAL_CHARS_DOC_SHORT].strip()
    return {
        "long_summary": long_summary,
        "short_summary": short_summary,
        "category": _infer_category_from_text(text, categories, file_name=file_name, file_path=file_path),
        "tags": _extract_top_tags(text, MAX_TAGS),
    }


def build_fallback_article_summary(text: str) -> str:
    return _extractive_summary_text(
        text,
        max_sentences=FALLBACK_ARTICLE_SENTENCES,
        max_total_chars=FALLBACK_MAX_TOTAL_CHARS_ARTICLE,
    )


def _summary_fallback_source_text(text: str) -> str:
    if len(text) <= MAX_DOC_SUMMARY_FALLBACK_SOURCE_CHARS:
        return text
    head = int(MAX_DOC_SUMMARY_FALLBACK_SOURCE_CHARS * 0.75)
    tail = MAX_DOC_SUMMARY_FALLBACK_SOURCE_CHARS - head
    return f"{text[:head]}\n\n[... truncated for local summary fallback ...]\n\n{text[-tail:]}"


def summary_guard_reason(text: str) -> str:
    if not text:
        return ""
    if len(text) > MAX_DOC_SUMMARY_INPUT_CHARS:
        return "summary_truncated_large_doc"
    chunk_count = max(1, math.ceil(len(text) / MAX_CHARS_PER_CHUNK))
    if chunk_count > MAX_DOC_SUMMARY_CHUNKS:
        return "summary_truncated_large_doc"
    return ""


def summarize_document_safe(
    cfg: AzureConfig,
    text: str,
    categories: List[str],
    file_label: str,
    file_path: str = "",
) -> Tuple[Dict[str, Any], str]:
    guard_reason = summary_guard_reason(text)
    if guard_reason:
        logging.warning(
            "Large-document summary guard triggered for %s (chars=%d); using local extractive fallback",
            file_label,
            len(text),
        )
        return (
            build_fallback_document_summary(
                _summary_fallback_source_text(text),
                categories,
                file_name=file_label,
                file_path=file_path,
            ),
            guard_reason,
        )
    try:
        return summarize_document(cfg, text, categories, file_name=file_label, file_path=file_path), ""
    except Exception as exc:
        if not is_content_filter_error(exc):
            raise
        logging.warning(
            "Content filter during document summary for %s; using local extractive fallback",
            file_label,
        )
        return (
            build_fallback_document_summary(
                _summary_fallback_source_text(text),
                categories,
                file_name=file_label,
                file_path=file_path,
            ),
            "summary_fallback_content_filter",
        )


def summarize_article_safe(cfg: AzureConfig, text: str, file_label: str) -> Tuple[str, bool]:
    try:
        return summarize_article(cfg, text), False
    except Exception as exc:
        if not is_content_filter_error(exc):
            raise
        logging.warning(
            "Content filter during article summary for %s; using local extractive fallback",
            file_label,
        )
        return build_fallback_article_summary(text), True


def split_for_excel(text: str, max_chars: int = 32767) -> List[str]:
    if not text:
        return [""]
    return split_text(text, max_chars)


def sanitize_excel_value(val: Any) -> Any:
    if isinstance(val, str):
        return ILLEGAL_EXCEL_CHARS_RE.sub("", val)
    return val


def sanitize_excel_df(df: pd.DataFrame) -> pd.DataFrame:
    try:
        # Include pandas "string" dtype to stay compatible across pandas 2.x/3.x.
        obj_cols = df.select_dtypes(include=["object", "string"]).columns
        for col in obj_cols:
            df[col] = df[col].apply(sanitize_excel_value)
    except Exception:
        pass
    return df


def drop_duplicate_header_rows(df: pd.DataFrame) -> pd.DataFrame:
    if df.empty:
        return df
    columns = list(df.columns)
    if not columns:
        return df

    keep_mask: List[bool] = []
    for row in df.itertuples(index=False, name=None):
        is_header_row = True
        for col_name, value in zip(columns, row):
            if str(value) != str(col_name):
                is_header_row = False
                break
        keep_mask.append(not is_header_row)

    if all(keep_mask):
        return df
    return df.loc[keep_mask].reset_index(drop=True)


def ensure_duplicate_review_columns(df: pd.DataFrame) -> pd.DataFrame:
    review_cols = ["GroupReviewed", "Decision", "DecisionNotes", "ReviewedBy"]
    for col in review_cols:
        if col not in df.columns:
            df[col] = ""
    return df


def load_category_path_map(path: Path) -> Dict[str, Any]:
    if not path.exists():
        return {}
    try:
        with path.open("r", encoding="utf-8") as f:
            data = json.load(f)
        return data if isinstance(data, dict) else {}
    except Exception:
        return {}


def validate_app_and_category_map(app_config: Dict[str, List[str]], category_path_map: Dict[str, Any]) -> None:
    errors: List[str] = []

    app_names = [str(a).strip() for a in app_config.keys() if str(a).strip()]
    app_names_cf = [a.casefold() for a in app_names]
    if len(app_names_cf) != len(set(app_names_cf)):
        errors.append("applications.json has duplicate application names (case-insensitive).")

    for app, cats in app_config.items():
        if not isinstance(cats, list):
            errors.append(f"applications.json app '{app}' must map to a list of categories.")
            continue
        seen_cf: set[str] = set()
        for c in cats:
            cs = str(c).strip()
            if not cs:
                errors.append(f"applications.json app '{app}' contains an empty category.")
                continue
            ccf = cs.casefold()
            if ccf in seen_cf:
                errors.append(f"applications.json app '{app}' has duplicate category '{cs}' (case-insensitive).")
            seen_cf.add(ccf)

    if not category_path_map:
        errors.append("category_path_map.json is empty or invalid JSON object.")
    else:
        map_apps = [str(a).strip() for a in category_path_map.keys() if str(a).strip()]
        map_apps_cf = {a.casefold(): a for a in map_apps}
        app_cfg_cf = {a.casefold(): a for a in app_names}

        for app_cf, app in app_cfg_cf.items():
            if app_cf not in map_apps_cf:
                errors.append(f"category_path_map.json missing application key '{app}'.")
                continue
            app_map = category_path_map.get(map_apps_cf[app_cf])
            if not isinstance(app_map, dict):
                errors.append(f"category_path_map.json application '{app}' must map to an object.")
                continue
            cfg_cats = app_config[app]
            cfg_cats_cf = {str(c).strip().casefold(): str(c).strip() for c in cfg_cats if str(c).strip()}
            map_cats_cf = {str(k).strip().casefold(): str(k).strip() for k in app_map.keys() if str(k).strip()}

            for cat_cf, cat in cfg_cats_cf.items():
                if cat_cf not in map_cats_cf:
                    errors.append(f"category_path_map.json missing category '{cat}' under application '{app}'.")
                    continue
                val = app_map.get(map_cats_cf[cat_cf])
                if not isinstance(val, str) or not val.strip():
                    errors.append(f"category_path_map.json has empty path for '{app}' -> '{cat}'.")

            extra = sorted(set(map_cats_cf.keys()) - set(cfg_cats_cf.keys()))
            for e in extra:
                errors.append(
                    f"category_path_map.json has extra category '{map_cats_cf[e]}' under application '{app}' not present in applications.json."
                )

        extra_apps = sorted(set(map_apps_cf.keys()) - set(app_cfg_cf.keys()))
        for e in extra_apps:
            errors.append(
                f"category_path_map.json has extra application '{map_apps_cf[e]}' not present in applications.json."
            )

    if errors:
        msg = "Configuration validation failed:\n- " + "\n- ".join(errors)
        raise ValueError(msg)


def resolve_category_path(category_path_map: Dict[str, Any], app_name: Optional[str], category: str) -> str:
    app_key = (app_name or "uncategorized").strip()
    cat = (category or "").strip()
    app_seg = sanitize_path_segment(app_key)
    cat_seg = sanitize_path_segment(cat)
    fallback = f"/Life_Sciences/Life_Science_Applications/{app_seg}"
    if cat_seg:
        fallback = f"{fallback}/{cat_seg}"

    def normalize_with_app_prefix(raw_path: str) -> str:
        noisy = {"? to", "top level (1)", "mid level (2)", "mid level (3)", "bottom level (4)"}
        parts = [
            str(p).strip()
            for p in str(raw_path).replace("\\", "/").split("/")
            if str(p).strip() and str(p).strip().casefold() not in noisy
        ]
        if not parts:
            return fallback
        # Absolute tool-template path: /Life_Sciences/Life_Science_Applications/...
        if parts[0].casefold() == "life_sciences":
            return "/" + "/".join(parts)
        if parts[0].casefold() == app_key.casefold():
            parts[0] = app_key
            return "/".join(parts)
        return f"{app_key}/{'/'.join(parts)}"

    if not category_path_map:
        return fallback

    app_map = category_path_map.get(app_key)
    if isinstance(app_map, dict):
        val = app_map.get(cat)
        if isinstance(val, str) and val.strip():
            return normalize_with_app_prefix(val.strip().strip("/\\"))

    val = category_path_map.get(cat)
    if isinstance(val, str) and val.strip():
        return normalize_with_app_prefix(val.strip().strip("/\\"))

    return fallback


def stable_import_id(file_path: str, title: str) -> str:
    base = f"{normalize_text(file_path)}|{normalize_text(title)}"
    return f"imp_{hashlib.sha1(base.encode('utf-8', errors='ignore')).hexdigest()[:16]}"


def sanitize_path_segment(value: str) -> str:
    s = (value or "").strip()
    s = re.sub(r"\s*\([^)]*\)", "", s)
    s = s.replace("&", " and ")
    s = s.replace("/", " ").replace("\\", " ").replace("'", "")
    s = re.sub(r"\s+", "_", s)
    s = re.sub(r"[^A-Za-z0-9_\-]", "", s)
    s = re.sub(r"_+", "_", s).strip("_")
    return s or "Uncategorized"


ARTICLE_TYPE_PATH_SEGMENTS = {
    "Application and Product Notes": "Application_and_Product_Notes",
    "Documentation": "Documentation",
    "Education & Training": "Education_and_Training",
    "Manuals and Guides": "Manuals_and_Guides",
    "FAQs": "FAQs",
    "Troubleshooting": "Troubleshooting",
    "References": "References",
}


def article_type_path_segment(article_type: str) -> str:
    return ARTICLE_TYPE_PATH_SEGMENTS.get((article_type or "").strip(), sanitize_path_segment(article_type))


def build_import_path(base_path: str, category: str, article_type: str) -> str:
    base = (base_path or "").strip().replace("\\", "/")
    if not base:
        base = "/Life_Sciences/Life_Science_Applications"
    if not base.startswith("/"):
        base = "/" + base
    base = re.sub(r"/+", "/", base).rstrip("/")
    guide_seg = f"{sanitize_path_segment(category)}_{article_type_path_segment(article_type)}"
    return f"{base}/{guide_seg}"


def display_doc_ref(doc_ref: str) -> str:
    ref = (doc_ref or "").strip()
    if not ref:
        return ""
    m = re.search(r"DOC-(\d+)$", ref)
    if not m:
        return ref
    try:
        return f"DOC-{int(m.group(1)):04d}"
    except ValueError:
        return f"DOC-{m.group(1)}"


def text_to_html(text: str) -> str:
    text = (text or "").replace("\r\n", "\n").replace("\r", "\n").strip()
    if not text:
        return "<p></p>"
    paragraphs = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]
    html_parts: List[str] = []
    for p in paragraphs:
        escaped = html.escape(p, quote=False).replace("\n", "<br>")
        html_parts.append(f"<p>{escaped}</p>")
    return "\n".join(html_parts) if html_parts else "<p></p>"


def classify_article_type_by_content(
    file_name: str,
    title: str,
    short_summary: str,
    long_summary: str,
    content_text: str,
) -> str:
    file_stem = Path(file_name or "").stem
    primary_text = normalize_text(" ".join([file_stem, title or "", short_summary or ""]))
    secondary_text = normalize_text(" ".join([long_summary or "", (content_text or "")[:12000]]))

    for label, needles in ARTICLE_TYPE_PRIMARY_RULES:
        if any(needle in primary_text for needle in needles):
            return label

    for label, needles in ARTICLE_TYPE_SECONDARY_RULES:
        if any(needle in secondary_text for needle in needles):
            return label

    troubleshooting_hits = sum(1 for needle in ARTICLE_TYPE_TROUBLESHOOTING_SECONDARY_TERMS if needle in secondary_text)
    if troubleshooting_hits >= 2:
        return "Troubleshooting"

    return "Documentation"


def attachment_path_for_doc(file_name: str, file_path: str) -> str:
    ext = Path(file_name or file_path or "").suffix.lower().lstrip(".")
    return f"attachments/{ext or 'file'}"


def write_import_excel(
    out_dir: Path,
    app_name: Optional[str],
    import_rows: List[Dict[str, Any]],
    append_excel: bool,
) -> Path:
    app_slug = sanitize_folder(app_name or "uncategorized")
    import_path = out_dir / f"{app_slug}__docatlas_import.xlsx"
    columns = ["Id", "Path", "Title", "Content", "Summary", "Tags", "Attachments", "AutoPublish", "ArticleType"]
    new_df = sanitize_excel_df(pd.DataFrame(import_rows, columns=columns))
    new_df = drop_duplicate_header_rows(new_df)

    if append_excel and import_path.exists():
        try:
            existing_df = pd.read_excel(import_path, sheet_name="import")
        except Exception:
            existing_df = pd.DataFrame(columns=columns)
        existing_df = sanitize_excel_df(existing_df)
        existing_df = drop_duplicate_header_rows(existing_df)
        if "Id" in existing_df.columns and "Id" in new_df.columns:
            existing_ids = set(existing_df["Id"].astype(str))
            new_df = new_df[~new_df["Id"].astype(str).isin(existing_ids)]
        out_df = pd.concat([existing_df, new_df], ignore_index=True)
    else:
        out_df = new_df
    out_df = drop_duplicate_header_rows(out_df)

    with pd.ExcelWriter(import_path, engine="openpyxl") as writer:
        out_df.to_excel(writer, index=False, sheet_name="import")
    return import_path


def write_full_text_legacy_structure_note(out_dir: Path) -> Path:
    note_path = out_dir / "full_text_legacy_structure.txt"
    text = (
        "Legacy Full Text Workbook Structure\n"
        "==================================\n\n"
        "Current status:\n"
        "- DocAtlas now writes full text as <app_slug>__docatlas_full_text.jsonl.gz by default.\n"
        "- The Excel full-text workbook below is retained here for historical reference only.\n\n"
        "Workbook name pattern:\n"
        "- <app_slug>__docatlas_full_text.xlsx\n\n"
        "Sheet:\n"
        "- FullText\n\n"
        "Typical columns:\n"
        "- doc_id, file_key, category, short_summary, long_summary, tags,\n"
        "  word_count, char_count, extraction_status, review_flags,\n"
        "  full_text, full_text_parts_count, full_text_part_1..N,\n"
        "  moved_to, file_name, file_path\n\n"
        "Purpose:\n"
        "- Historical workbook layout used before the default JSONL.GZ archive.\n"
    )
    note_path.write_text(text, encoding="utf-8")
    return note_path


def iter_full_text_archive_records(path: Path) -> Iterable[Dict[str, Any]]:
    if not path.exists():
        return
    with gzip.open(path, "rt", encoding="utf-8") as fh:
        for line in fh:
            line = line.strip()
            if not line:
                continue
            try:
                record = json.loads(line)
            except json.JSONDecodeError:
                continue
            if isinstance(record, dict):
                yield record


def write_full_text_archive(
    out_dir: Path,
    app_name: Optional[str],
    full_text_rows: List[Dict[str, Any]],
    append_excel: bool,
) -> Path:
    app_slug = sanitize_folder(app_name or "uncategorized")
    full_text_path = out_dir / f"{app_slug}__docatlas_full_text.jsonl.gz"
    rows_to_write = full_text_rows

    if append_excel and full_text_path.exists():
        existing_keys: set[str] = set()
        try:
            for record in iter_full_text_archive_records(full_text_path):
                file_key = str(record.get("file_key", "")).strip()
                if file_key:
                    existing_keys.add(file_key)
        except OSError:
            existing_keys = set()
            append_excel = False
        if existing_keys:
            rows_to_write = [row for row in full_text_rows if str(row.get("file_key", "")).strip() not in existing_keys]

    mode = "at" if append_excel and full_text_path.exists() else "wt"
    with gzip.open(full_text_path, mode, encoding="utf-8") as fh:
        for row in rows_to_write:
            fh.write(json.dumps(row, ensure_ascii=False, separators=(",", ":")))
            fh.write("\n")
    return full_text_path


def ocr_image_bytes(image_bytes: bytes) -> str:
    if pytesseract is None or Image is None:
        return ""
    try:
        with Image.open(io.BytesIO(image_bytes)) as img:
            return pytesseract.image_to_string(img) or ""
    except Exception:
        return ""


def extract_text_docx(path: Path, ocr_images: bool = False) -> str:
    if docx is None:
        raise RuntimeError("python-docx is not installed")
    doc = docx.Document(str(path))
    parts: List[str] = []
    for p in doc.paragraphs:
        if p.text:
            parts.append(p.text)
    for table in doc.tables:
        for row in table.rows:
            row_text = "\t".join(cell.text or "" for cell in row.cells)
            if row_text.strip():
                parts.append(row_text)
    if ocr_images:
        try:
            ocr_texts = []
            for shape in doc.inline_shapes:
                try:
                    r_id = shape._inline.graphic.graphicData.pic.blipFill.blip.embed
                    image_part = doc.part.related_parts.get(r_id)
                    if image_part is None:
                        continue
                    txt = ocr_image_bytes(image_part.blob)
                    if txt.strip():
                        ocr_texts.append(txt.strip())
                except Exception:
                    continue
            if ocr_texts:
                logging.info("OCR extracted text from %d images in %s", len(ocr_texts), path)
                parts.append("\n".join(ocr_texts))
        except Exception:
            pass
    return "\n".join(parts)


def convert_office_with_soffice(path: Path, target_ext: str) -> Optional[Path]:
    soffice = shutil.which("soffice") or shutil.which("soffice.exe")
    if not soffice:
        logging.warning("LibreOffice (soffice) not found; cannot convert %s", path)
        return None
    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            out_dir = Path(tmpdir)
            profile_dir = out_dir / "lo_profile"
            profile_dir.mkdir(parents=True, exist_ok=True)
            cmd = [
                soffice,
                "--headless",
                "--nologo",
                "--nodefault",
                "--nolockcheck",
                "--norestore",
                f"-env:UserInstallation={profile_dir.resolve().as_uri()}",
                "--convert-to",
                target_ext,
                "--outdir",
                str(out_dir),
                str(path),
            ]
            subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
            candidates = list(out_dir.glob(f"*.{target_ext}"))
            if not candidates:
                return None
            fd, temp_path = tempfile.mkstemp(suffix=f".{target_ext}")
            os.close(fd)
            target = Path(temp_path)
            shutil.copy2(candidates[0], target)
            return target
    except Exception as exc:
        logging.exception("Failed to convert %s: %s", path, exc)
        return None


def convert_doc_to_docx(path: Path) -> Optional[Path]:
    """Convert legacy .doc to .docx using LibreOffice (soffice)."""
    return convert_office_with_soffice(path, "docx")


def convert_ppt_to_pptx(path: Path) -> Optional[Path]:
    """Convert legacy .ppt to .pptx using LibreOffice (soffice)."""
    return convert_office_with_soffice(path, "pptx")


def convert_xls_to_xlsx(path: Path) -> Optional[Path]:
    """Convert legacy .xls to .xlsx using LibreOffice (soffice)."""
    return convert_office_with_soffice(path, "xlsx")


def extract_text_pptx(path: Path, ocr_images: bool = False) -> str:
    if pptx is None:
        raise RuntimeError("python-pptx is not installed")
    pres = pptx.Presentation(str(path))
    parts: List[str] = []
    for slide in pres.slides:
        for shape in slide.shapes:
            text = ""
            if hasattr(shape, "text"):
                text = shape.text or ""
            elif hasattr(shape, "text_frame") and shape.text_frame is not None:
                text = shape.text_frame.text or ""
            if text.strip():
                parts.append(text.strip())
            if ocr_images and MSO_SHAPE_TYPE is not None:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        img_blob = shape.image.blob
                        txt = ocr_image_bytes(img_blob)
                        if txt.strip():
                            parts.append(txt.strip())
                    except Exception:
                        pass
    return "\n".join(parts)


def extract_text_xlsx(path: Path) -> str:
    if openpyxl is None:
        raise RuntimeError("openpyxl is not installed")
    wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
    try:
        parts: List[str] = []
        for ws in wb.worksheets:
            parts.append(f"[Sheet: {ws.title}]")
            for row in ws.iter_rows(values_only=True):
                row_vals = [str(v) for v in row if v not in (None, "")]
                if row_vals:
                    parts.append("\t".join(row_vals))
        return "\n".join(parts)
    finally:
        try:
            wb.close()
        except Exception:
            pass


def extract_text_pdf(path: Path, ocrmypdf_enabled: bool) -> Tuple[str, List[str], str]:
    if pdfplumber is None:
        raise RuntimeError("pdfplumber is not installed")
    text, page_texts = extract_pdf_text_with_pdfplumber(path)
    logging.info("PDF text extracted (chars=%d) for %s", len(text.strip()), path)
    if len(text.strip()) >= MIN_EXTRACTED_CHARS:
        return text, page_texts, "ok"

    # OCRmyPDF default (if available):
    # 1) try without force first (better for tagged/text PDFs),
    # 2) force OCR only if still low-text.
    if ocrmypdf_enabled and ocrmypdf is not None:
        logging.info("OCR triggered (OCRmyPDF, non-forced pass) for %s", path)
        ocr_text, ocr_pages, ocr_status = ocrmypdf_ocr_pdf(path, force_ocr=False)
        if len(ocr_text.strip()) >= MIN_EXTRACTED_CHARS:
            return ocr_text, ocr_pages, "ocrmypdf_used"

        logging.info("OCR triggered (OCRmyPDF, forced pass) for %s", path)
        forced_text, forced_pages, forced_status = ocrmypdf_ocr_pdf(path, force_ocr=True)
        if len(forced_text.strip()) >= MIN_EXTRACTED_CHARS:
            return forced_text, forced_pages, "ocrmypdf_used_forced"

        # If OCRmyPDF failed or produced no text, try Tesseract fallback
        logging.info("OCR triggered (Tesseract fallback) for %s", path)
        ocr_texts, status = ocr_pdf(path)
        if ocr_texts:
            return "\n".join(ocr_texts), ocr_texts, "ocrmypdf_failed_then_ocr_used"
        # Preserve the most specific OCRmyPDF status from the forced pass.
        return text, page_texts, forced_status if forced_status else ocr_status

    # Fallback OCR (Tesseract)
    if ocrmypdf_enabled:
        logging.info("OCR triggered (Tesseract fallback; OCRmyPDF unavailable) for %s", path)
    else:
        logging.info("OCR triggered (Tesseract fallback; OCRmyPDF disabled) for %s", path)
    ocr_texts, status = ocr_pdf(path)
    if ocr_texts:
        return "\n".join(ocr_texts), ocr_texts, status
    return text, page_texts, status


def extract_pdf_text_with_pdfplumber(path: Path) -> Tuple[str, List[str]]:
    parts: List[str] = []
    page_texts: List[str] = []
    with pdfplumber.open(str(path)) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ""
            page_texts.append(text)
            if text.strip():
                parts.append(text)
    return "\n".join(parts), page_texts


def ocrmypdf_ocr_pdf(path: Path, force_ocr: bool = False) -> Tuple[str, List[str], str]:
    if ocrmypdf is None:
        return "", [], "no_text_ocrmypdf_unavailable"
    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            out_pdf = Path(tmpdir) / "ocr.pdf"
            kwargs = {
                "skip_text": not force_ocr,
                "output_type": "pdf",
                "progress_bar": False,
            }
            if force_ocr:
                kwargs["force_ocr"] = True
            try:
                ocrmypdf.ocr(str(path), str(out_pdf), **kwargs)
            except TypeError:
                # Older ocrmypdf versions may not support force_ocr
                kwargs.pop("force_ocr", None)
                ocrmypdf.ocr(str(path), str(out_pdf), **kwargs)
            text, page_texts = extract_pdf_text_with_pdfplumber(out_pdf)
            if len(text.strip()) >= MIN_EXTRACTED_CHARS:
                return text, page_texts, "ocrmypdf_used"
            return text, page_texts, "no_text_ocrmypdf"
    except Exception:
        return "", [], "no_text_ocrmypdf_failed"


def ocr_pdf(path: Path) -> Tuple[List[str], str]:
    if pytesseract is None or convert_from_path is None:
        return [], "no_text_ocr_unavailable"
    try:
        images = convert_from_path(str(path))
    except Exception:
        return [], "no_text_ocr_failed"
    texts: List[str] = []
    for img in images:
        try:
            txt = pytesseract.image_to_string(img)
        except Exception:
            txt = ""
        texts.append(txt or "")
    combined = "\n".join(texts).strip()
    if len(combined) >= MIN_EXTRACTED_CHARS:
        return texts, "ocr_used"
    return texts, "no_text"


def is_strong_heading(line: str) -> bool:
    line = (line or "").strip()
    if not line:
        return False
    if len(line) < 8 or len(line) > 110:
        return False
    if not re.search(r"[A-Za-z]", line):
        return False
    if line.endswith((".", ",", ";")) and not re.match(r"^[A-Z0-9][A-Z0-9 \-_/()]+$", line):
        return False

    if re.match(r"^(article|section|chapter)\s+([0-9]+|[ivxlcdm]+)\b", line, re.IGNORECASE):
        return True
    if re.match(r"^\d+(\.\d+){0,2}[)\.\-:]?\s+[A-Z][^\n]{2,}$", line):
        return True
    if re.match(r"^\([0-9]+\)\s+[A-Z][^\n]{2,}$", line):
        return True

    # All-caps headings are strong signals if length/word-count are sane.
    if line.isupper():
        words = line.split()
        if 2 <= len(words) <= 14:
            return True
    return False


def is_heading(line: str) -> bool:
    # Backward-compatible alias used by older call sites.
    return is_strong_heading(line)


def _single_article(lines: List[str]) -> List[Tuple[str, str]]:
    if not lines:
        return []
    return [("Article 1", "\n".join(lines))]


def _log_split_decision(
    source_label: str,
    decision: str,
    candidate_count: int,
    final_sections: int,
    total_chars: int,
    reason: str,
) -> None:
    logging.info(
        "Article split: decision=%s source=%s candidates=%d final_sections=%d total_chars=%d reason=%s",
        decision,
        source_label,
        candidate_count,
        final_sections,
        total_chars,
        reason,
    )


def split_pdf_into_articles(page_texts: List[str], source_label: str = "<pdf>") -> List[Tuple[str, str]]:
    lines: List[str] = []
    for page in page_texts:
        for line in (page or "").splitlines():
            cleaned = line.strip()
            if cleaned:
                lines.append(cleaned)
    if not lines:
        _log_split_decision(source_label, "single", 0, 0, 0, "no_lines")
        return []

    total_chars = sum(len(line) for line in lines)
    if total_chars < MIN_SPLIT_TOTAL_CHARS:
        _log_split_decision(source_label, "single", 0, 1, total_chars, "below_min_total_chars")
        return _single_article(lines)

    line_char_offsets: List[int] = []
    running = 0
    for line in lines:
        line_char_offsets.append(running)
        running += len(line) + 1

    candidates: List[Tuple[int, str]] = []
    for i, line in enumerate(lines):
        if is_strong_heading(line):
            if candidates:
                prev_i = candidates[-1][0]
                if i - prev_i < MIN_BOUNDARY_GAP_LINES:
                    continue
                if line_char_offsets[i] - line_char_offsets[prev_i] < MIN_BOUNDARY_GAP_CHARS:
                    continue
            candidates.append((i, line))

    if len(candidates) < MIN_SPLIT_SECTIONS:
        _log_split_decision(
            source_label,
            "single",
            len(candidates),
            1,
            total_chars,
            "insufficient_heading_candidates",
        )
        return _single_article(lines)

    sections: List[Tuple[str, str]] = []
    for idx, (start, title) in enumerate(candidates):
        end = candidates[idx + 1][0] if idx + 1 < len(candidates) else len(lines)
        body = "\n".join(lines[start + 1 : end]).strip()
        if not body:
            body = "\n".join(lines[start:end]).strip()
        sections.append((title or f"Article {idx + 1}", body))

    # Merge very small sections into neighbors to avoid fragmented article tabs.
    merged: List[Tuple[str, str]] = []
    for title, body in sections:
        clean_body = (body or "").strip()
        if not clean_body:
            continue
        if not merged:
            merged.append((title, clean_body))
            continue
        if len(clean_body) < MIN_SECTION_CHARS:
            prev_title, prev_body = merged[-1]
            merged[-1] = (prev_title, (prev_body + "\n" + clean_body).strip())
            continue
        merged.append((title, clean_body))

    final: List[Tuple[str, str]] = []
    for title, body in merged:
        if not final:
            final.append((title, body))
            continue
        if len(body) < MIN_SECTION_CHARS:
            prev_title, prev_body = final[-1]
            final[-1] = (prev_title, (prev_body + "\n" + body).strip())
        else:
            final.append((title, body))

    if len(final) < MIN_SPLIT_SECTIONS:
        _log_split_decision(
            source_label,
            "single",
            len(candidates),
            len(final),
            total_chars,
            "insufficient_final_sections",
        )
        return _single_article(lines)

    section_sizes = [len((body or "").strip()) for _, body in final if (body or "").strip()]
    if not section_sizes:
        _log_split_decision(source_label, "single", len(candidates), len(final), total_chars, "empty_sections")
        return _single_article(lines)

    strong_sections = sum(1 for n in section_sizes if n >= MIN_SECTION_CHARS)
    if strong_sections < MIN_SPLIT_SECTIONS:
        _log_split_decision(
            source_label,
            "single",
            len(candidates),
            len(final),
            total_chars,
            "insufficient_strong_sections",
        )
        return _single_article(lines)

    dominant_share = max(section_sizes) / max(1, sum(section_sizes))
    if dominant_share > MAX_SINGLE_SECTION_SHARE:
        _log_split_decision(
            source_label,
            "single",
            len(candidates),
            len(final),
            total_chars,
            "dominant_single_section",
        )
        return _single_article(lines)

    _log_split_decision(source_label, "multi", len(candidates), len(final), total_chars, "passed_conservative_checks")
    return final


def call_azure_chat(cfg: AzureConfig, messages: List[Dict[str, str]], temperature: float = 0.2) -> str:
    url = build_url(cfg.chat_base_url, cfg.chat_path, cfg.chat_deployment)
    headers = {cfg.api_key_header: cfg.chat_api_key, "Content-Type": "application/json"}
    payload: Dict[str, Any] = {
        "messages": messages,
        "temperature": float(temperature),
    }
    if cfg.include_model_in_body:
        payload["model"] = cfg.chat_deployment

    params = {"api-version": cfg.api_version}

    resp = _post_with_retries("Chat API", url, headers, params, payload)
    if resp.status_code >= 400:
        raise RuntimeError(f"Chat API error {resp.status_code}: {resp.text}")
    data = resp.json()
    choices = data.get("choices") or []
    if not choices:
        raise RuntimeError("Chat API returned no choices")
    message = choices[0].get("message", {})
    content = message.get("content", "") or ""
    in_chars = sum(len(m.get("content", "") or "") for m in messages)
    add_chat_usage(in_chars, len(content))
    return content


def call_azure_embeddings(cfg: AzureConfig, text: str) -> List[float]:
    url = build_url(cfg.embeddings_base_url, cfg.embeddings_path, cfg.embeddings_deployment)
    headers = {cfg.api_key_header: cfg.embeddings_api_key, "Content-Type": "application/json"}
    payload: Dict[str, Any] = {"input": text}
    if cfg.include_model_in_body:
        payload["model"] = cfg.embeddings_deployment
    params = {"api-version": cfg.api_version}
    resp = _post_with_retries("Embeddings API", url, headers, params, payload)
    if resp.status_code >= 400:
        raise RuntimeError(f"Embeddings API error {resp.status_code}: {resp.text}")
    data = resp.json()
    data_list = data.get("data") or []
    if not data_list:
        raise RuntimeError("Embeddings API returned no data")
    add_embed_usage(len(text))
    return data_list[0].get("embedding")


def extract_json(text: str) -> Dict[str, Any]:
    text = text.strip()
    if text.startswith("{"):
        try:
            return json.loads(text)
        except json.JSONDecodeError:
            pass
    # Try to find JSON block
    match = re.search(r"\{.*\}", text, re.DOTALL)
    if match:
        return json.loads(match.group(0))
    raise ValueError("Failed to parse JSON")


def summarize_document(
    cfg: AzureConfig,
    text: str,
    categories: List[str],
    file_name: str = "",
    file_path: str = "",
) -> Dict[str, Any]:
    categories_list = categories + (["Other"] if "Other" not in categories else [])
    if UNREADABLE_CATEGORY not in categories_list:
        categories_list.append(UNREADABLE_CATEGORY)
    if len(text) <= MAX_CHARS_PER_CHUNK:
        summary = summarize_with_model(cfg, text, categories_list)
        summary["category"] = _infer_category_from_text(
            text,
            categories_list,
            file_name=file_name,
            file_path=file_path,
        )
        return summary

    chunk_summaries: List[str] = []
    for chunk in split_text(text, MAX_CHARS_PER_CHUNK):
        chunk_summary = summarize_chunk(cfg, chunk)
        chunk_summaries.append(chunk_summary)
    combined = "\n".join(chunk_summaries)
    summary = summarize_with_model(cfg, combined, categories_list)
    summary["category"] = _infer_category_from_text(
        text,
        categories_list,
        file_name=file_name,
        file_path=file_path,
    )
    return summary


def summarize_chunk(cfg: AzureConfig, text: str) -> str:
    messages = [
        {"role": "system", "content": "You are a precise summarizer."},
        {
            "role": "user",
            "content": (
                "Summarize this chunk in 5-8 bullet points. Only output bullet points.\n\n"
                f"Chunk:\n{text}"
            ),
        },
    ]
    return call_azure_chat(cfg, messages)


def summarize_with_model(cfg: AzureConfig, text: str, categories: List[str]) -> Dict[str, Any]:
    categories_str = ", ".join(categories)
    messages = [
        {"role": "system", "content": "You are an expert analyst. Output JSON only."},
        {
            "role": "user",
            "content": (
                "Given the document text, produce JSON with keys: "
                "long_summary (5-7 sentences), short_summary (1-2 sentences), "
                "category (one of the provided categories), tags (array of strings). "
                "Tags can be as many as needed but should be specific and not redundant. "
                "If multiple categories could apply, prefer the most specific product/application "
                "category over broad process or issue buckets (e.g., prefer 'SeqStudio' over "
                "'Troubleshooting' when both fit). "
                f"Categories: {categories_str}.\n\n"
                f"Document:\n{text}"
            ),
        },
    ]
    # Keep category + summary selection deterministic across runs.
    content = call_azure_chat(cfg, messages, temperature=0.0)
    return extract_json(content)


def summarize_article(cfg: AzureConfig, text: str) -> str:
    text = text[:MAX_ARTICLE_CHARS]
    messages = [
        {"role": "system", "content": "You are a precise technical writer."},
        {
            "role": "user",
            "content": (
                "Write a condensed article that preserves all key facts, findings, and conclusions. "
                "Length should be as long as needed to capture every important point (do not force "
                "a fixed sentence count). No bullets. Do not start with phrases like "
                "'The article is about' or 'This article'. Write directly.\n\n"
                f"Article:\n{text}"
            ),
        },
    ]
    return call_azure_chat(cfg, messages).strip()


def cosine_similarity(a: np.ndarray, b: np.ndarray) -> float:
    if a.size == 0 or b.size == 0:
        return 0.0
    denom = np.linalg.norm(a) * np.linalg.norm(b)
    if denom == 0:
        return 0.0
    return float(np.dot(a, b) / denom)


def detect_duplicates(
    items: List[Tuple[str, str, Optional[np.ndarray]]], threshold: float
) -> Tuple[Dict[str, str], Dict[str, float], Dict[str, str]]:
    """
    items: list of (item_id, exact_hash, embedding)
    returns: duplicate_of, duplicate_score, duplicate_group_id
    """
    duplicate_of: Dict[str, str] = {}
    duplicate_score: Dict[str, float] = {}
    duplicate_group_id: Dict[str, str] = {}

    hash_to_primary: Dict[str, str] = {}
    group_counter = 1

    for item_id, hsh, _emb in items:
        if hsh in hash_to_primary:
            primary = hash_to_primary[hsh]
            duplicate_of[item_id] = primary
            duplicate_score[item_id] = 1.0
            group = duplicate_group_id.get(primary)
            if not group:
                group = f"DUP-{group_counter:04d}"
                group_counter += 1
                duplicate_group_id[primary] = group
            duplicate_group_id[item_id] = group
            continue

        hash_to_primary[hsh] = item_id

    return duplicate_of, duplicate_score, duplicate_group_id


def _sorted_pair(a: str, b: str) -> Tuple[str, str]:
    return (a, b) if a <= b else (b, a)


def _near_dup_name_tokens(file_name: str) -> set[str]:
    stem = Path(file_name or "").stem
    return {
        tok
        for tok in re.findall(r"[a-z0-9][a-z0-9\-]{2,}", normalize_text(stem))
        if tok not in SUMMARY_STOPWORDS and tok not in NEAR_DUP_FILE_TOKEN_STOPWORDS and (len(tok) >= 4 or any(ch.isdigit() for ch in tok))
    }


def _near_dup_parent_tokens(file_path: str) -> set[str]:
    normalized = str(file_path or "").replace("\\", "/")
    if "!/" in normalized:
        normalized = normalized.split("!/", 1)[1]
    parts = [p for p in normalized.split("/") if p]
    if len(parts) <= 1:
        return set()
    parent_text = " ".join(parts[:-1])
    return {
        tok
        for tok in re.findall(r"[a-z0-9][a-z0-9\-]{2,}", normalize_text(parent_text))
        if tok not in SUMMARY_STOPWORDS and len(tok) >= 4
    }


def _near_dup_path_family(file_path: str) -> str:
    normalized = str(file_path or "").replace("\\", "/")
    if "!/" in normalized:
        normalized = normalized.split("!/", 1)[1]
    parts = [p.casefold() for p in normalized.split("/") if p]
    if len(parts) <= 1:
        return ""
    return "/".join(parts[: min(2, len(parts) - 1)])


def has_weak_near_duplicate_signal(doc_a: DocRecord, doc_b: DocRecord) -> bool:
    if _near_dup_path_family(doc_a.file_path) and _near_dup_path_family(doc_a.file_path) == _near_dup_path_family(doc_b.file_path):
        return True
    if _near_dup_name_tokens(doc_a.file_name).intersection(_near_dup_name_tokens(doc_b.file_name)):
        return True
    if _near_dup_parent_tokens(doc_a.file_path).intersection(_near_dup_parent_tokens(doc_b.file_path)):
        return True
    return False


def detect_near_duplicates_docs(
    docs: List[DocRecord],
    doc_embeddings: Dict[str, Optional[np.ndarray]],
    min_mutual: float = NEAR_DUP_MUTUAL,
    strong_score: float = NEAR_DUP_STRONG,
) -> Tuple[Dict[str, str], Dict[str, float], Dict[str, str], Dict[str, set[str]], set[Tuple[str, str]]]:
    """
    Detect near-duplicate candidate edges within the same category.
    Returns:
    - near_duplicate_of
    - near_duplicate_score
    - near_duplicate_group_id
    - near_adjacency (doc_id -> neighbor ids)
    - near_edges (undirected doc_id pairs)
    """
    docs_by_id = {d.doc_id: d for d in docs}
    by_category: Dict[str, List[str]] = {}
    for d in docs:
        if d.category == UNREADABLE_CATEGORY:
            continue
        emb = doc_embeddings.get(d.doc_id)
        if emb is None:
            continue
        by_category.setdefault(d.category, []).append(d.doc_id)

    near_edges: set[Tuple[str, str]] = set()
    near_edge_scores: Dict[Tuple[str, str], float] = {}
    near_adjacency: Dict[str, set[str]] = {}
    near_duplicate_of: Dict[str, str] = {}
    near_duplicate_score: Dict[str, float] = {}
    near_duplicate_group_id: Dict[str, str] = {}

    for cat, doc_ids in by_category.items():
        if len(doc_ids) < 2:
            continue
        ids_sorted = sorted(doc_ids, key=lambda x: docs_by_id[x].file_path.lower())
        vectors: Dict[str, np.ndarray] = {}
        for doc_id in ids_sorted:
            vec = doc_embeddings.get(doc_id)
            if vec is None:
                continue
            norm = np.linalg.norm(vec)
            if norm == 0:
                continue
            vectors[doc_id] = vec / norm
        ids = [i for i in ids_sorted if i in vectors]
        if len(ids) < 2:
            continue

        top1: Dict[str, Tuple[str, float]] = {}
        for i in range(len(ids)):
            id_i = ids[i]
            vec_i = vectors[id_i]
            for j in range(i + 1, len(ids)):
                id_j = ids[j]
                vec_j = vectors[id_j]
                score = float(np.dot(vec_i, vec_j))
                pair = _sorted_pair(id_i, id_j)
                if score >= strong_score:
                    near_edges.add(pair)
                    near_edge_scores[pair] = max(score, near_edge_scores.get(pair, 0.0))
                prev_i = top1.get(id_i)
                if prev_i is None or score > prev_i[1]:
                    top1[id_i] = (id_j, score)
                prev_j = top1.get(id_j)
                if prev_j is None or score > prev_j[1]:
                    top1[id_j] = (id_i, score)

        for id_i, (id_j, score_i) in top1.items():
            if score_i < min_mutual:
                continue
            top_j = top1.get(id_j)
            if top_j is None:
                continue
            if top_j[0] != id_i or top_j[1] < min_mutual:
                continue
            if not has_weak_near_duplicate_signal(docs_by_id[id_i], docs_by_id[id_j]):
                continue
            pair = _sorted_pair(id_i, id_j)
            near_edges.add(pair)
            near_edge_scores[pair] = max(score_i, near_edge_scores.get(pair, 0.0))

    for a, b in near_edges:
        near_adjacency.setdefault(a, set()).add(b)
        near_adjacency.setdefault(b, set()).add(a)

    # Near group IDs (deterministic)
    seen: set[str] = set()
    components: List[List[str]] = []
    for d in sorted(docs, key=lambda x: x.file_path.lower()):
        if d.doc_id in seen:
            continue
        if d.doc_id not in near_adjacency:
            continue
        stack = [d.doc_id]
        comp: List[str] = []
        seen.add(d.doc_id)
        while stack:
            cur = stack.pop()
            comp.append(cur)
            for nei in sorted(near_adjacency.get(cur, set())):
                if nei in seen:
                    continue
                seen.add(nei)
                stack.append(nei)
        if len(comp) >= 2:
            components.append(sorted(comp, key=lambda x: docs_by_id[x].file_path.lower()))

    components.sort(key=lambda comp: docs_by_id[comp[0]].file_path.lower())
    for idx, comp in enumerate(components, start=1):
        gid = f"NDUP-{idx:04d}"
        for doc_id in comp:
            near_duplicate_group_id[doc_id] = gid

    # Per-doc best near link
    for d in docs:
        neighbors = near_adjacency.get(d.doc_id, set())
        if not neighbors:
            continue
        best_id = ""
        best_score = -1.0
        for nei in neighbors:
            pair = _sorted_pair(d.doc_id, nei)
            score = float(near_edge_scores.get(pair, 0.0))
            if score > best_score:
                best_score = score
                best_id = nei
            elif abs(score - best_score) < 1e-12:
                if best_id and docs_by_id[nei].file_path.lower() < docs_by_id[best_id].file_path.lower():
                    best_id = nei
        if best_id:
            near_duplicate_of[d.doc_id] = best_id
            near_duplicate_score[d.doc_id] = best_score

    return near_duplicate_of, near_duplicate_score, near_duplicate_group_id, near_adjacency, near_edges


def build_unified_review_groups(
    docs: List[DocRecord],
    exact_pairs: set[Tuple[str, str]],
    near_pairs: set[Tuple[str, str]],
) -> Tuple[Dict[str, str], Dict[str, str]]:
    """
    Build connected review groups from exact + near relations (within same category).
    Returns:
    - review_group_id by doc_id
    - duplicate_relation_type by doc_id
    """
    docs_by_id = {d.doc_id: d for d in docs}
    review_adj: Dict[str, set[str]] = {}
    exact_adj: Dict[str, set[str]] = {}
    near_adj: Dict[str, set[str]] = {}

    def _add_pair(pair: Tuple[str, str], target_adj: Dict[str, set[str]]) -> None:
        a, b = pair
        da = docs_by_id.get(a)
        db = docs_by_id.get(b)
        if da is None or db is None:
            return
        if da.category != db.category:
            return
        review_adj.setdefault(a, set()).add(b)
        review_adj.setdefault(b, set()).add(a)
        target_adj.setdefault(a, set()).add(b)
        target_adj.setdefault(b, set()).add(a)

    for pair in sorted(exact_pairs):
        _add_pair(pair, exact_adj)
    for pair in sorted(near_pairs):
        _add_pair(pair, near_adj)

    review_group_id: Dict[str, str] = {}
    relation_type: Dict[str, str] = {}

    seen: set[str] = set()
    components: List[List[str]] = []
    for d in sorted(docs, key=lambda x: x.file_path.lower()):
        doc_id = d.doc_id
        if doc_id in seen or doc_id not in review_adj:
            continue
        stack = [doc_id]
        comp: List[str] = []
        seen.add(doc_id)
        while stack:
            cur = stack.pop()
            comp.append(cur)
            for nei in sorted(review_adj.get(cur, set())):
                if nei in seen:
                    continue
                seen.add(nei)
                stack.append(nei)
        if len(comp) >= 2:
            components.append(sorted(comp, key=lambda x: docs_by_id[x].file_path.lower()))

    components.sort(key=lambda comp: (docs_by_id[comp[0]].category.lower(), docs_by_id[comp[0]].file_path.lower()))
    for idx, comp in enumerate(components, start=1):
        gid = f"RGRP-{idx:04d}"
        for doc_id in comp:
            review_group_id[doc_id] = gid
            has_exact = bool(exact_adj.get(doc_id))
            has_near = bool(near_adj.get(doc_id))
            if has_exact and has_near:
                relation_type[doc_id] = "exact+near"
            elif has_exact:
                relation_type[doc_id] = "exact"
            elif has_near:
                relation_type[doc_id] = "near"
            else:
                relation_type[doc_id] = ""

    return review_group_id, relation_type


def get_categories_gui(
    app_config: Dict[str, List[str]],
    config_path: Path,
) -> Tuple[List[str], Optional[str]]:
    if tk is None:
        raise RuntimeError("tkinter is not available")

    root = tk.Tk()
    root.title("DocAtlas")
    root.geometry("720x560")
    apply_theme(root)

    selected_app = tk.StringVar(value="")

    container = tk.Frame(root, bg=THEME["bg"])
    container.pack(fill=tk.BOTH, expand=True, padx=16, pady=16)

    header = tk.Label(container, text="DocAtlas", bg=THEME["bg"], fg=THEME["fg"], font=FONT_HEADER)
    header.pack(anchor="w", pady=(0, 6))
    logo = tk.Label(container, text="— DocAtlas —", bg=THEME["bg"], fg=THEME["muted"], font=FONT_SMALL)
    logo.pack(anchor="w", pady=(0, 6))
    sub = tk.Label(
        container,
        text="Select Application & Categories",
        bg=THEME["bg"],
        fg=THEME["muted"],
        font=FONT_SMALL,
    )
    sub.pack(anchor="w", pady=(0, 12))

    btn_frame = tk.Frame(container, bg=THEME["bg"])
    btn_frame.pack(side=tk.BOTTOM, pady=12, fill=tk.X)

    panel = tk.Frame(container, bg=THEME["panel"], highlightbackground=THEME["border"], highlightthickness=1)
    panel.pack(fill=tk.BOTH, expand=True)

    if app_config:
        label_app = tk.Label(panel, text="Application", bg=THEME["panel"], fg=THEME["fg"], font=FONT_LABEL)
        label_app.pack(anchor="w", padx=12, pady=(12, 4))
        options = ["(Custom)"] + sorted(app_config.keys())
        dropdown = tk.OptionMenu(panel, selected_app, *options)
        dropdown.config(bg=THEME["text_bg"], fg=THEME["fg"], highlightthickness=0, font=FONT_BASE)
        try:
            dropdown["menu"].config(bg=THEME["text_bg"], fg=THEME["fg"], font=FONT_BASE)
        except Exception:
            pass
        dropdown.pack(anchor="w", padx=12, pady=(0, 8))
        selected_app.set(options[0])

    label = tk.Label(panel, text="Categories (one per line)", bg=THEME["panel"], fg=THEME["fg"], font=FONT_LABEL)
    label.pack(anchor="w", padx=12, pady=(8, 4))

    text_widget = tk.Text(
        panel,
        height=16,
        width=78,
        bg=THEME["text_bg"],
        fg=THEME["fg"],
        insertbackground=THEME["fg"],
        highlightbackground=THEME["border"],
        font=FONT_BASE,
    )
    text_widget.pack(padx=12, pady=(0, 12), fill=tk.BOTH, expand=True)

    result: List[str] = []
    result_app: Optional[str] = None

    def apply_app_categories(*_args: Any) -> None:
        app = selected_app.get()
        if app and app in app_config:
            text_widget.delete("1.0", tk.END)
            text_widget.insert(tk.END, "\n".join(app_config[app]))

    if app_config:
        selected_app.trace_add("write", apply_app_categories)

    def on_ok() -> None:
        content = text_widget.get("1.0", tk.END)
        lines = [line.strip() for line in content.splitlines() if line.strip()]
        if not lines:
            messagebox.showerror("Error", "Please enter at least one category.")
            return
        nonlocal result
        result = lines
        nonlocal result_app
        app = selected_app.get().strip()
        if app and app != "(Custom)" and app in app_config:
            result_app = app
        root.destroy()

    def on_cancel() -> None:
        root.destroy()

    if app_config:
        tk.Button(
            btn_frame,
            text="Edit Apps",
            command=lambda: edit_applications_gui(config_path, app_config, root),
            width=12,
            bg=THEME["btn_bg"],
            fg=THEME["btn_fg"],
            relief=tk.FLAT,
            borderwidth=1,
            font=FONT_BUTTON,
        ).pack(side=tk.RIGHT, padx=8, ipady=6)
    tk.Button(btn_frame, text="Cancel", command=on_cancel, width=12, bg=THEME["btn_bg"], fg=THEME["btn_fg"], relief=tk.RAISED, borderwidth=1, font=FONT_BUTTON).pack(side=tk.RIGHT, padx=8, ipady=6)
    tk.Button(btn_frame, text="OK", command=on_ok, width=12, bg=THEME["accent"], fg="#ffffff", relief=tk.RAISED, borderwidth=1, font=FONT_BUTTON).pack(side=tk.RIGHT, padx=8, ipady=6)

    root.mainloop()
    return result, result_app


def pick_directories_gui() -> Tuple[Path, Path]:
    if tk is None:
        raise RuntimeError("tkinter is not available")

    root = tk.Tk()
    root.withdraw()
    input_dir = filedialog.askdirectory(title="Select Input Folder")
    if not input_dir:
        raise RuntimeError("No input folder selected")

    output_dir = filedialog.askdirectory(
        title="Select Output Folder (Cancel = use input folder)",
        initialdir=input_dir,
    )
    if not output_dir:
        output_dir = input_dir

    root.destroy()
    return Path(input_dir), Path(output_dir)


def get_ocrmypdf_gui() -> bool:
    if tk is None:
        return True
    root = tk.Tk()
    root.title("DocAtlas - OCR Options")
    root.geometry("520x260")
    apply_theme(root)

    var = tk.BooleanVar(value=True)
    container = tk.Frame(root, bg=THEME["bg"])
    container.pack(fill=tk.BOTH, expand=True, padx=16, pady=16)

    label = tk.Label(container, text="OCR Settings", bg=THEME["bg"], fg=THEME["fg"], font=FONT_HEADER)
    label.pack(anchor="w", pady=(0, 6))
    sub = tk.Label(container, text="OCR runs only when PDFs have little or no text.", bg=THEME["bg"], fg=THEME["muted"], font=FONT_SMALL)
    sub.pack(anchor="w", pady=(0, 12))

    chk = tk.Checkbutton(container, text="Use OCRmyPDF (recommended)", variable=var, bg=THEME["bg"], fg=THEME["fg"], font=FONT_LABEL)
    chk.pack(anchor="w", pady=6)

    result: List[bool] = []

    def on_ok() -> None:
        result.append(bool(var.get()))
        root.destroy()

    btn_frame = tk.Frame(container, bg=THEME["bg"])
    btn_frame.pack(pady=12, anchor="e", fill=tk.X)

    def on_test() -> None:
        missing = warn_missing_ocr_deps(bool(var.get()))
        if not missing and messagebox is not None:
            messagebox.showinfo("OCR Dependencies", "All OCR dependencies found.")

    tk.Button(btn_frame, text="Test OCR", command=on_test, width=12, bg=THEME["btn_bg"], fg=THEME["btn_fg"], relief=tk.FLAT, height=1, font=FONT_BUTTON).pack(side=tk.RIGHT, padx=8, ipady=4)
    tk.Button(btn_frame, text="OK", command=on_ok, width=12, bg=THEME["accent"], fg="#ffffff", relief=tk.FLAT, height=1, font=FONT_BUTTON).pack(side=tk.RIGHT, padx=8, ipady=4)
    root.mainloop()
    return result[0] if result else True


def get_embeddings_source_gui() -> Tuple[str, bool]:
    if tk is None:
        return DEFAULT_EMBEDDINGS_SOURCE, True
    root = tk.Tk()
    root.title("DocAtlas - Embeddings Source")
    root.geometry("520x320")
    apply_theme(root)

    var = tk.StringVar(value=DEFAULT_EMBEDDINGS_SOURCE)
    append_var = tk.BooleanVar(value=True)
    container = tk.Frame(root, bg=THEME["bg"])
    container.pack(fill=tk.BOTH, expand=True, padx=16, pady=16)

    label = tk.Label(container, text="Embeddings Source", bg=THEME["bg"], fg=THEME["fg"], font=FONT_HEADER)
    label.pack(anchor="w", pady=(0, 6))
    sub = tk.Label(
        container,
        text="Choose what text is embedded for duplicate detection.",
        bg=THEME["bg"],
        fg=THEME["muted"],
        font=FONT_SMALL,
    )
    sub.pack(anchor="w", pady=(0, 12))

    tk.Radiobutton(
        container,
        text="Long summary (lower cost)",
        variable=var,
        value="summary",
        bg=THEME["bg"],
        fg=THEME["fg"],
        selectcolor=THEME["panel"],
        font=FONT_LABEL,
    ).pack(anchor="w", pady=4)
    tk.Radiobutton(
        container,
        text="Full text (recommended, stricter)",
        variable=var,
        value="full_text",
        bg=THEME["bg"],
        fg=THEME["fg"],
        selectcolor=THEME["panel"],
        font=FONT_LABEL,
    ).pack(anchor="w", pady=4)
    tk.Radiobutton(
        container,
        text="Disable embeddings (hash-only duplicates)",
        variable=var,
        value=EMBEDDINGS_SOURCE_NONE,
        bg=THEME["bg"],
        fg=THEME["fg"],
        selectcolor=THEME["panel"],
        font=FONT_LABEL,
    ).pack(anchor="w", pady=4)

    chk = tk.Checkbutton(
        container,
        text="Append to existing Excel (recommended)",
        variable=append_var,
        bg=THEME["bg"],
        fg=THEME["fg"],
        font=FONT_LABEL,
    )
    chk.pack(anchor="w", pady=(8, 4))

    result: List[Tuple[str, bool]] = []

    def on_ok() -> None:
        result.append((var.get(), bool(append_var.get())))
        root.destroy()

    btn_frame = tk.Frame(container, bg=THEME["bg"])
    btn_frame.pack(pady=12, anchor="e", fill=tk.X)
    tk.Button(btn_frame, text="OK", command=on_ok, width=12, bg=THEME["accent"], fg="#ffffff", relief=tk.FLAT, height=1, font=FONT_BUTTON).pack(
        side=tk.RIGHT, padx=8, ipady=4
    )
    root.mainloop()
    return result[0] if result else (DEFAULT_EMBEDDINGS_SOURCE, True)


def get_run_mode_gui() -> Optional[bool]:
    if tk is None:
        return True
    root = tk.Tk()
    root.title("DocAtlas - Run Mode")
    root.geometry("760x320")
    apply_theme(root)

    var = tk.StringVar(value="charter")
    container = tk.Frame(root, bg=THEME["bg"])
    container.pack(fill=tk.BOTH, expand=True, padx=16, pady=16)

    label = tk.Label(container, text="Run Mode", bg=THEME["bg"], fg=THEME["fg"], font=FONT_HEADER)
    label.pack(anchor="w", pady=(0, 6))
    sub = tk.Label(
        container,
        text="Choose how to run this job.",
        bg=THEME["bg"],
        fg=THEME["muted"],
        font=FONT_SMALL,
    )
    sub.pack(anchor="w", pady=(0, 12))

    tk.Radiobutton(
        container,
        text="Charter Mode (summaries, tags, duplicates; no file moves)",
        variable=var,
        value="charter",
        bg=THEME["bg"],
        fg=THEME["fg"],
        selectcolor=THEME["panel"],
        font=FONT_LABEL,
    ).pack(anchor="w", pady=4)
    tk.Radiobutton(
        container,
        text="Atlas Mode (summaries, tags, duplicates; move files)",
        variable=var,
        value="atlas",
        bg=THEME["bg"],
        fg=THEME["fg"],
        selectcolor=THEME["panel"],
        font=FONT_LABEL,
    ).pack(anchor="w", pady=4)

    result: List[Optional[bool]] = []

    def on_ok() -> None:
        result.append(var.get() == "charter")
        root.destroy()

    def on_close() -> None:
        result.append(None)
        root.destroy()

    root.protocol("WM_DELETE_WINDOW", on_close)

    btn_frame = tk.Frame(container, bg=THEME["bg"])
    btn_frame.pack(pady=12, anchor="e", fill=tk.X)
    tk.Button(btn_frame, text="OK", command=on_ok, width=12, bg=THEME["accent"], fg="#ffffff", relief=tk.FLAT, height=1, font=FONT_BUTTON).pack(
        side=tk.RIGHT, padx=8, ipady=4
    )
    root.mainloop()
    return result[0] if result else True


def get_articles_mode_gui() -> bool:
    if tk is None:
        return False
    root = tk.Tk()
    root.title("DocAtlas - Article Generation")
    root.geometry("560x250")
    apply_theme(root)

    var = tk.BooleanVar(value=False)
    container = tk.Frame(root, bg=THEME["bg"])
    container.pack(fill=tk.BOTH, expand=True, padx=16, pady=16)

    label = tk.Label(container, text="Article Generation", bg=THEME["bg"], fg=THEME["fg"], font=FONT_HEADER)
    label.pack(anchor="w", pady=(0, 6))
    sub = tk.Label(
        container,
        text="Generate per-article rows for PDFs with clear multi-article structure.",
        bg=THEME["bg"],
        fg=THEME["muted"],
        font=FONT_SMALL,
    )
    sub.pack(anchor="w", pady=(0, 12))

    tk.Checkbutton(
        container,
        text="Enable article generation (PDF-only)",
        variable=var,
        bg=THEME["bg"],
        fg=THEME["fg"],
        font=FONT_LABEL,
    ).pack(anchor="w", pady=6)

    result: List[bool] = []

    def on_ok() -> None:
        result.append(bool(var.get()))
        root.destroy()

    def on_close() -> None:
        result.append(False)
        root.destroy()

    root.protocol("WM_DELETE_WINDOW", on_close)

    btn_frame = tk.Frame(container, bg=THEME["bg"])
    btn_frame.pack(pady=12, anchor="e", fill=tk.X)
    tk.Button(
        btn_frame,
        text="OK",
        command=on_ok,
        width=12,
        bg=THEME["accent"],
        fg="#ffffff",
        relief=tk.FLAT,
        height=1,
        font=FONT_BUTTON,
    ).pack(side=tk.RIGHT, padx=8, ipady=4)
    root.mainloop()
    return result[0] if result else False


def argv_has_flag(argv_tokens: List[str], flag: str) -> bool:
    prefix = flag + "="
    return any(tok == flag or tok.startswith(prefix) for tok in argv_tokens)


def prompt_text_cli(prompt: str, default: Optional[str] = None, allow_empty: bool = False) -> str:
    suffix = f" [{default}]" if default not in (None, "") else ""
    while True:
        value = input(f"{prompt}{suffix}: ").strip()
        if not value and default is not None:
            value = default
        if value or allow_empty:
            return value
        print("Value required.")


def prompt_path_cli(prompt: str, default: Optional[str] = None, must_exist: bool = False) -> Path:
    while True:
        path = Path(prompt_text_cli(prompt, default=default)).expanduser()
        if must_exist and not path.exists():
            print(f"Path not found: {path}")
            continue
        return path


def prompt_choice_cli(prompt: str, options: List[Tuple[str, Any]], default_index: int = 1) -> Any:
    print(prompt)
    for idx, (label, _value) in enumerate(options, start=1):
        print(f"  {idx}) {label}")
    while True:
        raw = input(f"Select option [{default_index}]: ").strip()
        if not raw:
            return options[default_index - 1][1]
        if raw.isdigit():
            idx = int(raw)
            if 1 <= idx <= len(options):
                return options[idx - 1][1]
        print("Enter a valid option number.")


def prompt_yes_no_cli(prompt: str, default: bool) -> bool:
    suffix = "[Y/n]" if default else "[y/N]"
    while True:
        raw = input(f"{prompt} {suffix}: ").strip().lower()
        if not raw:
            return default
        if raw in ("y", "yes"):
            return True
        if raw in ("n", "no"):
            return False
        print("Enter y or n.")


def prompt_int_cli(prompt: str, default: int, min_value: int = 1) -> int:
    while True:
        raw = prompt_text_cli(prompt, default=str(default))
        try:
            value = int(raw)
        except ValueError:
            print("Enter a whole number.")
            continue
        if value < min_value:
            print(f"Value must be >= {min_value}.")
            continue
        return value


def get_categories_cli(app_config: Dict[str, List[str]]) -> Tuple[List[str], Optional[str]]:
    if app_config:
        options: List[Tuple[str, str]] = [(app, app) for app in sorted(app_config.keys())]
        options.append(("Custom categories (semicolon-separated)", "__custom__"))
        selected = prompt_choice_cli("Application", options, default_index=1)
        if selected != "__custom__":
            return list(app_config[selected]), selected
    raw = prompt_text_cli("Categories separated by semicolons")
    categories = [c.strip() for c in raw.split(";") if c.strip()]
    if not categories:
        raise ValueError("At least one category is required")
    return categories, None


def resolve_cli_interactive_inputs(
    args: argparse.Namespace,
    argv_tokens: List[str],
    app_config: Dict[str, List[str]],
) -> Tuple[Path, Path, List[str], Optional[str], bool, str, bool, bool, bool, int]:
    if not sys.stdin.isatty():
        raise RuntimeError("--interactive requires a terminal")

    input_dir = Path(args.input).expanduser() if args.input else prompt_path_cli("Input folder", must_exist=True)
    output_dir = Path(args.output).expanduser() if args.output else prompt_path_cli("Output folder", default=str(input_dir))

    if args.categories:
        categories = [c.strip() for c in args.categories.split(";") if c.strip()]
        if not categories:
            raise ValueError("Provide at least one category in --categories")
        app_name = None
    elif args.app and args.app in app_config:
        categories = list(app_config[args.app])
        app_name = args.app
    else:
        if args.app:
            print(f"Application not found in config: {args.app}")
        categories, app_name = get_categories_cli(app_config)

    if args.charter_mode or args.signal_scan or args.no_move:
        no_move = True
    else:
        no_move = prompt_choice_cli(
            "Run mode",
            [
                ("Charter mode (recommended: no file moves)", True),
                ("Atlas mode (move files into output folders)", False),
            ],
            default_index=1,
        )

    ocrmypdf_enabled = (not args.no_ocrmypdf) if argv_has_flag(argv_tokens, "--no-ocrmypdf") else prompt_yes_no_cli(
        "Enable OCR when PDFs have little or no text?",
        True,
    )

    embeddings_source = (
        resolve_embeddings_source(args.embeddings_source)
        if argv_has_flag(argv_tokens, "--embeddings-source")
        else prompt_choice_cli(
            "Embeddings source",
            [
                ("Full text (recommended, stricter duplicate detection)", "full_text"),
                ("Long summary (lower cost)", "summary"),
                ("Disable embeddings (hash-only duplicates)", EMBEDDINGS_SOURCE_NONE),
            ],
            default_index=1,
        )
    )

    append_excel = (not args.overwrite_excel) if argv_has_flag(argv_tokens, "--overwrite-excel") else prompt_choice_cli(
        "Excel write mode",
        [
            ("Append to existing outputs (recommended)", True),
            ("Overwrite existing outputs", False),
        ],
        default_index=1,
    )

    articles_enabled = bool(args.articles) if argv_has_flag(argv_tokens, "--articles") else prompt_yes_no_cli(
        "Enable article generation for PDFs?",
        False,
    )

    workers = args.workers if argv_has_flag(argv_tokens, "--workers") else prompt_int_cli("Workers", default=1, min_value=1)

    return input_dir, output_dir, categories, app_name, ocrmypdf_enabled, embeddings_source, append_excel, no_move, articles_enabled, workers


def edit_applications_gui(config_path: Path, app_config: Dict[str, List[str]], parent: tk.Tk) -> None:
    if tk is None:
        return
    win = tk.Toplevel(parent)
    win.title("DocAtlas - Edit Applications")
    win.geometry("820x560")
    apply_theme(win)

    apps = dict(app_config)

    left_frame = tk.Frame(win, bg=THEME["panel"], highlightbackground=THEME["border"], highlightthickness=1)
    left_frame.pack(side=tk.LEFT, fill=tk.Y, padx=10, pady=10)

    right_frame = tk.Frame(win, bg=THEME["bg"])
    right_frame.pack(side=tk.RIGHT, fill=tk.BOTH, expand=True, padx=10, pady=10)

    listbox = tk.Listbox(
        left_frame,
        height=20,
        width=25,
        bg=THEME["text_bg"],
        fg=THEME["fg"],
        selectbackground=THEME["accent"],
        selectforeground="#ffffff",
        font=FONT_BASE,
    )
    listbox.pack(pady=5)

    def refresh_list() -> None:
        listbox.delete(0, tk.END)
        for name in sorted(apps.keys()):
            listbox.insert(tk.END, name)

    def load_selected(event: Any = None) -> None:
        selection = listbox.curselection()
        if not selection:
            return
        name = listbox.get(selection[0])
        name_entry.delete(0, tk.END)
        name_entry.insert(0, name)
        cat_text.delete("1.0", tk.END)
        cat_text.insert(tk.END, "\n".join(apps.get(name, [])))

    def add_new() -> None:
        name_entry.delete(0, tk.END)
        cat_text.delete("1.0", tk.END)

    def save_current() -> None:
        name = name_entry.get().strip()
        if not name:
            if messagebox:
                messagebox.showerror("Error", "Application name is required.")
            return
        cats = [c.strip() for c in cat_text.get("1.0", tk.END).splitlines() if c.strip()]
        apps[name] = cats
        refresh_list()

    def delete_current() -> None:
        name = name_entry.get().strip()
        if not name:
            return
        if name in apps:
            del apps[name]
        refresh_list()
        name_entry.delete(0, tk.END)
        cat_text.delete("1.0", tk.END)

    def save_all() -> None:
        save_current()
        save_app_config(config_path, apps)
        app_config.clear()
        app_config.update(apps)
        if messagebox:
            messagebox.showinfo("Saved", f"Saved to {config_path}")

    def save_and_close() -> None:
        save_all()
        win.destroy()

    listbox.bind("<<ListboxSelect>>", load_selected)

    tk.Label(right_frame, text="Application Name", bg=THEME["bg"], fg=THEME["fg"], font=FONT_LABEL).pack(anchor="w")
    name_entry = tk.Entry(right_frame, width=40, bg=THEME["text_bg"], fg=THEME["fg"], font=FONT_BASE)
    name_entry.pack(fill=tk.X, pady=4)

    tk.Label(right_frame, text="Categories (one per line)", bg=THEME["bg"], fg=THEME["fg"], font=FONT_LABEL).pack(anchor="w")
    cat_text = tk.Text(
        right_frame,
        height=15,
        bg=THEME["text_bg"],
        fg=THEME["fg"],
        insertbackground=THEME["fg"],
        highlightbackground=THEME["border"],
        font=FONT_BASE,
    )
    cat_text.pack(fill=tk.BOTH, expand=True)

    btn_row = tk.Frame(right_frame, bg=THEME["bg"])
    btn_row.pack(pady=8, fill=tk.X)
    tk.Button(btn_row, text="Save & Close", command=save_and_close, width=14, bg=THEME["accent"], fg="#ffffff", relief=tk.FLAT, height=1, padx=10, pady=6, font=FONT_BUTTON).pack(side=tk.RIGHT, padx=8)
    tk.Button(btn_row, text="Save All", command=save_all, width=12, bg=THEME["btn_bg"], fg=THEME["btn_fg"], relief=tk.FLAT, height=1, padx=10, pady=6, font=FONT_BUTTON).pack(side=tk.RIGHT, padx=8)
    tk.Button(btn_row, text="Delete", command=delete_current, width=10, bg=THEME["btn_bg"], fg=THEME["btn_fg"], relief=tk.FLAT, height=1, padx=10, pady=6, font=FONT_BUTTON).pack(side=tk.RIGHT, padx=8)
    tk.Button(btn_row, text="Save App", command=save_current, width=10, bg=THEME["btn_bg"], fg=THEME["btn_fg"], relief=tk.FLAT, height=1, padx=10, pady=6, font=FONT_BUTTON).pack(side=tk.RIGHT, padx=8)
    tk.Button(btn_row, text="New", command=add_new, width=10, bg=THEME["btn_bg"], fg=THEME["btn_fg"], relief=tk.FLAT, height=1, padx=10, pady=6, font=FONT_BUTTON).pack(side=tk.RIGHT, padx=8)

    refresh_list()


def process_file(
    path: Path,
    ocrmypdf_enabled: bool,
    articles_enabled: bool = True,
    source_label: Optional[str] = None,
) -> Tuple[str, List[Tuple[str, str]], str]:
    ext = path.suffix.lower()
    try:
        if path.stat().st_size == 0:
            return "", [], "no_text_empty"
    except Exception:
        pass
    if ext == ".docx":
        text = extract_text_docx(path, ocrmypdf_enabled)
        status = "ok" if len(text.strip()) >= MIN_EXTRACTED_CHARS else "no_text"
        return text, [], status
    if ext == ".doc":
        docx_path = convert_doc_to_docx(path)
        if docx_path is None:
            return "", [], "no_text_doc_convert_failed"
        try:
            text = extract_text_docx(docx_path, ocrmypdf_enabled)
        finally:
            try:
                docx_path.unlink(missing_ok=True)
            except Exception:
                pass
        status = "ok" if len(text.strip()) >= MIN_EXTRACTED_CHARS else "no_text"
        return text, [], status
    if ext == ".pptx":
        text = extract_text_pptx(path, ocrmypdf_enabled)
        status = "ok" if len(text.strip()) >= MIN_EXTRACTED_CHARS else "no_text"
        return text, [], status
    if ext == ".ppt":
        pptx_path = convert_ppt_to_pptx(path)
        if pptx_path is None:
            return "", [], "no_text_ppt_convert_failed"
        try:
            text = extract_text_pptx(pptx_path, ocrmypdf_enabled)
        finally:
            try:
                pptx_path.unlink(missing_ok=True)
            except Exception:
                pass
        status = "ok" if len(text.strip()) >= MIN_EXTRACTED_CHARS else "no_text"
        return text, [], status
    if ext == ".xlsx":
        text = extract_text_xlsx(path)
        status = "ok" if len(text.strip()) >= MIN_EXTRACTED_CHARS else "no_text"
        return text, [], status
    if ext == ".xls":
        xlsx_path = convert_xls_to_xlsx(path)
        if xlsx_path is None:
            return "", [], "no_text_xls_convert_failed"
        try:
            text = extract_text_xlsx(xlsx_path)
        finally:
            try:
                xlsx_path.unlink(missing_ok=True)
            except Exception:
                pass
        status = "ok" if len(text.strip()) >= MIN_EXTRACTED_CHARS else "no_text"
        return text, [], status
    if ext == ".pdf":
        text, pages, status = extract_text_pdf(path, ocrmypdf_enabled)
        articles = split_pdf_into_articles(pages, source_label=source_label or str(path)) if articles_enabled else []
        return text, articles, status
    raise ValueError(f"Unsupported file type: {ext}")


def write_excels(
    out_dir: Path,
    docs: List[DocRecord],
    articles: List[ArticleRecord],
    full_text_rows: List[Dict[str, Any]],
    app_name: Optional[str],
    append_excel: bool,
    category_path_map: Dict[str, Any],
    include_full_text_output: bool = DEFAULT_INCLUDE_FULL_TEXT_OUTPUT,
    articles_enabled: bool = False,
) -> Tuple[Path, Optional[Path], Path]:
    app_slug = sanitize_folder(app_name or "uncategorized")
    peers_path = out_dir / f"{app_slug}__docatlas_summaries.xlsx"
    full_text_path = out_dir / f"{app_slug}__docatlas_full_text.jsonl.gz"

    existing_docs_df = None
    existing_dups_df = None
    existing_articles_df = None
    existing_doc_keys: set[str] = set()
    existing_doc_paths: set[str] = set()

    if append_excel and peers_path.exists():
        try:
            existing_docs_df = pd.read_excel(peers_path, sheet_name="Documents")
            existing_docs_df = drop_duplicate_header_rows(sanitize_excel_df(existing_docs_df))
            if "file_key" in existing_docs_df.columns:
                existing_doc_keys = set(existing_docs_df["file_key"].astype(str))
            if "file_path" in existing_docs_df.columns:
                for val in existing_docs_df["file_path"].astype(str):
                    existing_doc_paths.update(path_compare_aliases(val))
            if "FilePath" in existing_docs_df.columns:
                for val in existing_docs_df["FilePath"].astype(str):
                    existing_doc_paths.update(path_compare_aliases(val))
        except Exception:
            existing_docs_df = None
            existing_doc_keys = set()
            existing_doc_paths = set()
        try:
            existing_dups_df = pd.read_excel(peers_path, sheet_name="Duplicates")
            existing_dups_df = drop_duplicate_header_rows(sanitize_excel_df(existing_dups_df))
        except Exception:
            existing_dups_df = None
        try:
            existing_articles_df = pd.read_excel(peers_path, sheet_name="Articles")
            existing_articles_df = drop_duplicate_header_rows(sanitize_excel_df(existing_articles_df))
        except Exception:
            existing_articles_df = None

    docs_rows: List[Dict[str, Any]] = []
    new_docs: List[DocRecord] = []
    for d in docs:
        if append_excel and (d.file_key in existing_doc_keys or normalize_compare_path(d.file_path) in existing_doc_paths):
            continue
        new_docs.append(d)
        docs_rows.append(
            {
                "Category": d.category,
                "FilePath": d.file_path,
                "FileName": d.file_name,
                "DuplicateOf": display_doc_ref(d.duplicate_of),
                "DupScore": float(d.duplicate_score) if d.duplicate_score is not None else 0.0,
                "NearDuplicateOf": display_doc_ref(d.near_duplicate_of),
                "NearDupScore": float(d.near_duplicate_score) if d.near_duplicate_score is not None else 0.0,
                "LongSummary": d.long_summary,
                "ShortSummary": d.short_summary,
                "ReviewFlag": d.review_flags,
                "ExtractionStatus": d.extraction_status,
                "DuplicateClusterID": d.duplicate_group_id,
                "NearDuplicateClusterID": d.near_duplicate_group_id,
                "ReviewGroupID": d.review_group_id,
                "DuplicateRelationType": d.duplicate_relation_type,
            }
        )

    dup_rows: List[Dict[str, Any]] = []
    for d in new_docs:
        if not d.review_group_id:
            continue
        dup_rows.append(
            {
                "ReviewGroupID": d.review_group_id,
                "Category": d.category,
                "FilePath": d.file_path,
                "FileName": d.file_name,
                "DuplicateRelationType": d.duplicate_relation_type,
                "DupScore": float(d.duplicate_score) if d.duplicate_score is not None else 0.0,
                "NearDupScore": float(d.near_duplicate_score) if d.near_duplicate_score is not None else 0.0,
                "DuplicateOf": display_doc_ref(d.duplicate_of),
                "NearDuplicateOf": display_doc_ref(d.near_duplicate_of),
                "ReviewFlag": d.review_flags,
            }
        )
    relation_rank = {"exact": 0, "exact+near": 1, "near": 2}
    dup_rows.sort(
        key=lambda r: (
            str(r.get("Category", "")).lower(),
            str(r.get("ReviewGroupID", "")).lower(),
            relation_rank.get(str(r.get("DuplicateRelationType", "")).lower(), 3),
            -(
                float(r.get("DupScore", 0.0))
                if float(r.get("DupScore", 0.0)) > 0
                else float(r.get("NearDupScore", 0.0))
            ),
            str(r.get("FileName", "")).lower(),
        )
    )

    new_doc_ids = {d.doc_id for d in new_docs}
    doc_by_id = {d.doc_id: d for d in docs}
    article_rows: List[Dict[str, Any]] = []
    if articles_enabled:
        for a in articles:
            if append_excel and a.doc_id not in new_doc_ids:
                continue
            parent_doc = doc_by_id.get(a.doc_id)
            article_rows.append(
                {
                    "Category": parent_doc.category if parent_doc else "",
                    "FilePath": a.file_path,
                    "FileName": a.file_name,
                    "ParentDocID": a.doc_id,
                    "ArticleIndex": a.article_index,
                    "ArticleTitle": a.article_title,
                    "ArticleSummary": a.article_summary,
                    "DuplicateClusterID": a.duplicate_group_id,
                    "DupScore": float(a.duplicate_score) if a.duplicate_score is not None else 0.0,
                    "DuplicateOf": a.duplicate_of,
                }
            )
        article_rows.sort(
            key=lambda r: (
                str(r.get("Category", "")).lower(),
                str(r.get("FileName", "")).lower(),
                int(r.get("ArticleIndex", 0)),
            )
        )

    docs_columns = [
        "Category",
        "FilePath",
        "FileName",
        "DuplicateOf",
        "DupScore",
        "NearDuplicateOf",
        "NearDupScore",
        "LongSummary",
        "ShortSummary",
        "ReviewFlag",
        "ExtractionStatus",
        "DuplicateClusterID",
        "NearDuplicateClusterID",
        "ReviewGroupID",
        "DuplicateRelationType",
    ]
    dups_columns = [
        "ReviewGroupID",
        "Category",
        "FilePath",
        "FileName",
        "DuplicateRelationType",
        "DupScore",
        "NearDupScore",
        "DuplicateOf",
        "NearDuplicateOf",
        "ReviewFlag",
    ]
    articles_columns = [
        "Category",
        "FilePath",
        "FileName",
        "ParentDocID",
        "ArticleIndex",
        "ArticleTitle",
        "ArticleSummary",
        "DuplicateClusterID",
        "DupScore",
        "DuplicateOf",
    ]

    docs_df = drop_duplicate_header_rows(sanitize_excel_df(pd.DataFrame(docs_rows, columns=docs_columns)))
    dups_df = drop_duplicate_header_rows(sanitize_excel_df(pd.DataFrame(dup_rows, columns=dups_columns)))
    articles_df: Optional[pd.DataFrame] = None
    if articles_enabled:
        articles_df = drop_duplicate_header_rows(sanitize_excel_df(pd.DataFrame(article_rows, columns=articles_columns)))
    if append_excel and existing_docs_df is not None:
        docs_df = pd.concat([existing_docs_df, docs_df], ignore_index=True)
    if append_excel and existing_dups_df is not None:
        dups_df = pd.concat([existing_dups_df, dups_df], ignore_index=True)
    if articles_enabled and append_excel and existing_articles_df is not None and articles_df is not None:
        articles_df = pd.concat([existing_articles_df, articles_df], ignore_index=True)
    docs_df = drop_duplicate_header_rows(docs_df)
    dups_df = drop_duplicate_header_rows(dups_df)
    if articles_df is not None:
        articles_df = drop_duplicate_header_rows(articles_df)
    dups_df = ensure_duplicate_review_columns(dups_df)
    if not dups_df.empty:
        if "DupScore" in dups_df.columns:
            dups_df["DupScore"] = pd.to_numeric(dups_df["DupScore"], errors="coerce").fillna(0.0)
        if "NearDupScore" in dups_df.columns:
            dups_df["NearDupScore"] = pd.to_numeric(dups_df["NearDupScore"], errors="coerce").fillna(0.0)
        if "DuplicateRelationType" not in dups_df.columns:
            dups_df["DuplicateRelationType"] = ""
        dups_df["_relation_rank"] = dups_df["DuplicateRelationType"].astype(str).str.lower().map(relation_rank).fillna(3)
        dups_df["_best_score"] = np.where(
            dups_df["DupScore"] > 0,
            dups_df["DupScore"],
            dups_df.get("NearDupScore", 0.0),
        )
        sort_cols = [c for c in ["Category", "ReviewGroupID", "_relation_rank", "_best_score", "FileName"] if c in dups_df.columns]
        if sort_cols:
            asc_map = {"Category": True, "ReviewGroupID": True, "_relation_rank": True, "_best_score": False, "FileName": True}
            ascending = [asc_map.get(c, True) for c in sort_cols]
            dups_df = dups_df.sort_values(sort_cols, ascending=ascending, kind="mergesort", na_position="last").reset_index(drop=True)
        dups_df = dups_df.drop(columns=[c for c in ["_relation_rank", "_best_score"] if c in dups_df.columns], errors="ignore")

    articles_sheet_status = "skipped"
    with pd.ExcelWriter(peers_path, engine="openpyxl") as writer:
        docs_df.to_excel(writer, index=False, sheet_name="Documents")
        dups_df.to_excel(writer, index=False, sheet_name="Duplicates")
        pd.DataFrame(
            [
                {
                    "Field": "GroupReviewed",
                    "Allowed Values": "Reviewed",
                    "Meaning": "Mark one row in a review group as Reviewed; all rows in the same ReviewGroupID will be highlighted.",
                },
                {
                    "Field": "GroupReviewed",
                    "Allowed Values": "Unfinished",
                    "Meaning": "Use when a duplicate group has been opened but is not finalized yet.",
                },
                {
                    "Field": "Decision",
                    "Allowed Values": "Primary",
                    "Meaning": "Chosen file to keep as the main migration candidate in a duplicate group.",
                },
                {
                    "Field": "Decision",
                    "Allowed Values": "Keep",
                    "Meaning": "Keep this file as well.",
                },
                {
                    "Field": "Decision",
                    "Allowed Values": "Drop",
                    "Meaning": "Do not migrate this file.",
                },
                {
                    "Field": "Decision",
                    "Allowed Values": "Needs Review",
                    "Meaning": "Still unresolved.",
                },
            ]
        ).to_excel(writer, index=False, sheet_name="duplicate_review_legend")
        if articles_enabled and articles_df is not None:
            articles_df.to_excel(writer, index=False, sheet_name="Articles")
            articles_sheet_status = "written"
        elif append_excel and existing_articles_df is not None:
            existing_articles_df.to_excel(writer, index=False, sheet_name="Articles")
            articles_sheet_status = "preserved"
    logging.info("Articles sheet in summaries workbook: %s", articles_sheet_status)

    text_by_doc_id = {str(r.get("doc_id", "")): str(r.get("full_text", "") or "") for r in full_text_rows}
    import_rows: List[Dict[str, Any]] = []
    for d in docs:
        if d.category == UNREADABLE_CATEGORY:
            continue
        content_text = text_by_doc_id.get(d.doc_id, "")
        title = d.short_summary or d.file_name
        article_type = classify_article_type_by_content(d.file_name, title, d.short_summary, d.long_summary, content_text)
        base_path = resolve_category_path(category_path_map, app_name, d.category)
        import_rows.append(
            {
                "Id": stable_import_id(d.file_path, title),
                "Path": build_import_path(base_path, d.category, article_type),
                "Title": title,
                "Content": text_to_html(content_text),
                "Summary": d.long_summary,
                "Tags": ", ".join(d.tags),
                "Attachments": attachment_path_for_doc(d.file_name, d.file_path),
                "AutoPublish": True,
                "ArticleType": article_type,
            }
        )
    import_path = write_import_excel(out_dir, app_name, import_rows, append_excel)

    full_text_path = write_full_text_archive(out_dir, app_name, full_text_rows, append_excel)
    write_full_text_legacy_structure_note(out_dir)

    # Apply formatting: wrap summaries, add reviewer UX, and widen columns
    try:
        from openpyxl import load_workbook
        from openpyxl.styles import Alignment, Font, PatternFill
        from openpyxl.formatting.rule import FormulaRule
        from openpyxl.utils import get_column_letter
        from openpyxl.worksheet.datavalidation import DataValidation

        def format_sheet(path: Path, sheet_name: str, wrap_cols: List[str]) -> None:
            wb = load_workbook(path)
            ws = wb[sheet_name]
            header = {cell.value: cell.column for cell in ws[1]}
            wrap = Alignment(wrap_text=True, vertical="top")
            center = Alignment(horizontal="center", vertical="top")
            dark_fill = PatternFill(fill_type="solid", fgColor="1F2937")
            review_fill = PatternFill(fill_type="solid", fgColor="0F766E")
            white_bold = Font(color="FFFFFF", bold=True)

            # Style and preserve the visible header/filter/freeze UX without Excel table objects.
            for cell in ws[1]:
                cell.fill = dark_fill
                cell.font = white_bold
                cell.alignment = wrap
            ws.row_dimensions[1].height = 28
            for col_name in wrap_cols:
                col_idx = header.get(col_name)
                if not col_idx:
                    continue
                for row in ws.iter_rows(min_row=2, min_col=col_idx, max_col=col_idx):
                    for cell in row:
                        cell.alignment = wrap
                # widen column
                col_letter = ws.cell(row=1, column=col_idx).column_letter
                ws.column_dimensions[col_letter].width = 60

            # Apply filter/freeze only; avoid Table objects because they trigger repair warnings after rewrites.
            max_row = ws.max_row
            max_col = ws.max_column
            if max_row >= 2 and max_col >= 1:
                ref = f"A1:{ws.cell(row=max_row, column=max_col).coordinate}"
                ws.auto_filter.ref = ref
                ws.freeze_panes = "A2"

            if sheet_name == "Duplicates":
                review_headers = {
                    "GroupReviewed": 16,
                    "Decision": 18,
                    "DecisionNotes": 36,
                    "ReviewedBy": 18,
                }
                for col_name, width in review_headers.items():
                    col_idx = header.get(col_name)
                    if not col_idx:
                        continue
                    col_letter = get_column_letter(col_idx)
                    ws.column_dimensions[col_letter].width = width
                    ws.cell(row=1, column=col_idx).fill = review_fill
                    for row in ws.iter_rows(min_row=2, min_col=col_idx, max_col=col_idx):
                        for cell in row:
                            cell.alignment = center if col_name in ("GroupReviewed", "Decision") else wrap

                if max_row >= 2:
                    group_col = get_column_letter(header["GroupReviewed"])
                    decision_col = get_column_letter(header["Decision"])
                    reviewed_dv = DataValidation(type="list", formula1='"Reviewed,Unfinished"', allow_blank=True)
                    decision_dv = DataValidation(
                        type="list",
                        formula1='"Primary,Keep,Drop,Needs Review"',
                        allow_blank=True,
                    )
                    ws.add_data_validation(reviewed_dv)
                    ws.add_data_validation(decision_dv)
                    reviewed_dv.add(f"{group_col}2:{group_col}{max_row}")
                    decision_dv.add(f"{decision_col}2:{decision_col}{max_row}")

                    row_fill = PatternFill(fill_type="solid", fgColor="E2F0D9")
                    primary_fill = PatternFill(fill_type="solid", fgColor="89B36D")
                    keep_fill = PatternFill(fill_type="solid", fgColor="F3E7D2")
                    drop_fill = PatternFill(fill_type="solid", fgColor="E7D7F4")
                    needs_fill = PatternFill(fill_type="solid", fgColor="DCEAF7")
                    primary_font = Font(color="1F3B08", bold=True)
                    keep_font = Font(color="5B4B2A", bold=True)
                    drop_font = Font(color="5A2740", bold=True)
                    needs_font = Font(color="234A68", bold=True)
                    row_ref = f"A2:{ws.cell(row=max_row, column=max_col).coordinate}"
                    decision_ref = f"{decision_col}2:{decision_col}{max_row}"
                    reviewed_formula = (
                        f'COUNTIFS($A$2:$A${max_row},$A2,${group_col}$2:${group_col}${max_row},"Reviewed")>0'
                    )
                    ws.conditional_formatting.add(
                        row_ref,
                        FormulaRule(formula=[reviewed_formula], fill=row_fill),
                    )
                    ws.conditional_formatting.add(
                        decision_ref,
                        FormulaRule(formula=[f'${decision_col}2="Primary"'], fill=primary_fill, font=primary_font),
                    )
                    ws.conditional_formatting.add(
                        decision_ref,
                        FormulaRule(formula=[f'${decision_col}2="Keep"'], fill=keep_fill, font=keep_font),
                    )
                    ws.conditional_formatting.add(
                        decision_ref,
                        FormulaRule(formula=[f'${decision_col}2="Drop"'], fill=drop_fill, font=drop_font),
                    )
                    ws.conditional_formatting.add(
                        decision_ref,
                        FormulaRule(formula=[f'${decision_col}2="Needs Review"'], fill=needs_fill, font=needs_font),
                    )

            if sheet_name == "duplicate_review_legend":
                for cell in ws[1]:
                    cell.fill = dark_fill
                    cell.font = white_bold
                    cell.alignment = wrap
                ws.row_dimensions[1].height = 28
                widths = {"A": 20, "B": 24, "C": 96}
                for col_letter, width in widths.items():
                    ws.column_dimensions[col_letter].width = width
                for row in ws.iter_rows(min_row=2):
                    for cell in row:
                        cell.alignment = wrap
                if max_row >= 2:
                    ws.auto_filter.ref = f"A1:C{max_row}"
                    ws.freeze_panes = "A2"

            wb.save(path)

        format_sheet(peers_path, "Documents", ["LongSummary", "ShortSummary", "FilePath"])
        format_sheet(peers_path, "Duplicates", ["FilePath"])
        format_sheet(peers_path, "duplicate_review_legend", ["Meaning"])
        if articles_sheet_status in ("written", "preserved"):
            format_sheet(peers_path, "Articles", ["FilePath", "ArticleTitle", "ArticleSummary"])
        format_sheet(import_path, "import", ["Path", "Title", "Content", "Summary", "Tags", "Attachments", "ArticleType"])
    except Exception:
        pass

    return peers_path, full_text_path, import_path


def write_summary_report(
    out_dir: Path,
    docs: List[DocRecord],
    articles: List[ArticleRecord],
    unsupported_files: Optional[List[UnsupportedFileRecord]] = None,
    errors: Optional[List[Dict[str, str]]] = None,
    usage: Optional[Dict[str, int]] = None,
    total_files: Optional[int] = None,
    processed_files: Optional[int] = None,
    limit: Optional[int] = None,
) -> Path:
    report_path = out_dir / "summary_report.txt"
    total_docs = len(docs)
    total_articles = len(articles)
    dup_docs = sum(1 for d in docs if d.duplicate_of)
    dup_articles = sum(1 for a in articles if a.duplicate_of)

    categories: Dict[str, int] = {}
    for d in docs:
        categories[d.category] = categories.get(d.category, 0) + 1

    ext_counts: Dict[str, int] = {}
    for d in docs:
        ext_counts[d.file_ext] = ext_counts.get(d.file_ext, 0) + 1

    extraction_status: Dict[str, int] = {}
    for d in docs:
        extraction_status[d.extraction_status] = extraction_status.get(d.extraction_status, 0) + 1
    no_text_docs = [d for d in docs if d.extraction_status != "ok"]
    ocr_docs = [d for d in docs if d.extraction_status.startswith("ocr") or "ocr" in d.extraction_status]

    word_counts = [d.word_count for d in docs if d.word_count is not None]
    char_counts = [d.char_count for d in docs if d.char_count is not None]
    avg_words = int(sum(word_counts) / len(word_counts)) if word_counts else 0
    avg_chars = int(sum(char_counts) / len(char_counts)) if char_counts else 0
    longest = max(docs, key=lambda d: d.char_count or 0, default=None)
    shortest = min(docs, key=lambda d: d.char_count or 0, default=None)

    dup_group_sizes: Dict[str, int] = {}
    for d in docs:
        if d.duplicate_group_id:
            dup_group_sizes[d.duplicate_group_id] = dup_group_sizes.get(d.duplicate_group_id, 0) + 1
    dup_group_count = len(dup_group_sizes)
    avg_dup_group_size = int(sum(dup_group_sizes.values()) / dup_group_count) if dup_group_count else 0
    near_dup_docs = sum(1 for d in docs if d.near_duplicate_of or d.near_duplicate_group_id)
    near_dup_group_count = len({d.near_duplicate_group_id for d in docs if d.near_duplicate_group_id})
    review_group_count = len({d.review_group_id for d in docs if d.review_group_id})
    unsupported_stats = unsupported_file_stats(unsupported_files or [])

    lines = []
    lines.append("Summary Report")
    lines.append("================")
    lines.append(f"Total documents: {total_docs}")
    lines.append(f"Total articles: {total_articles}")
    lines.append(f"Duplicate documents: {dup_docs}")
    lines.append(f"Duplicate articles: {dup_articles}")
    lines.append("")
    lines.append("Documents by Category:")
    for k in sorted(categories.keys()):
        pct = (categories[k] / total_docs * 100) if total_docs else 0
        lines.append(f"- {k}: {categories[k]} ({pct:.1f}%)")
    lines.append("")
    lines.append("Documents by File Type:")
    for k in sorted(ext_counts.keys()):
        pct = (ext_counts[k] / total_docs * 100) if total_docs else 0
        lines.append(f"- {k}: {ext_counts[k]} ({pct:.1f}%)")
    lines.append("")
    lines.append("Unsupported Files:")
    lines.append(f"- total_unsupported_files: {unsupported_stats['count']}")
    if unsupported_stats["by_source_kind"]:
        lines.append(f"- normal_folder_count: {unsupported_stats['by_source_kind'].get('file', 0)}")
        lines.append(f"- zip_member_count: {unsupported_stats['by_source_kind'].get('zip_member', 0)}")
        lines.append(f"- invalid_zip_count: {unsupported_stats['by_source_kind'].get('invalid_zip', 0)}")
    for file_type in sorted(unsupported_stats["by_type"].keys()):
        lines.append(f"- {file_type}: {unsupported_stats['by_type'][file_type]}")
    if unsupported_stats["by_source_folder"]:
        lines.append("- top_source_folders:")
        for folder, count in sorted(unsupported_stats["by_source_folder"].items(), key=lambda kv: (-kv[1], kv[0]))[:10]:
            lines.append(f"  - {folder}: {count}")
    lines.append("")
    lines.append("Document Length (approx):")
    lines.append(f"- avg_words: {avg_words}")
    lines.append(f"- avg_chars: {avg_chars}")
    if longest:
        lines.append(f"- longest: {longest.file_name} ({longest.char_count} chars)")
    if shortest:
        lines.append(f"- shortest: {shortest.file_name} ({shortest.char_count} chars)")
    lines.append("")
    lines.append("OCR Usage:")
    lines.append(f"- ocr_used: {len(ocr_docs)}")
    lines.append("")
    lines.append("Duplicate Groups:")
    lines.append(f"- duplicate_group_count: {dup_group_count}")
    lines.append(f"- avg_duplicate_group_size: {avg_dup_group_size}")
    lines.append(f"- near_duplicate_document_count: {near_dup_docs}")
    lines.append(f"- near_duplicate_group_count: {near_dup_group_count}")
    lines.append(f"- review_group_count: {review_group_count}")
    lines.append("")
    lines.append("Extraction Status:")
    for k in sorted(extraction_status.keys()):
        lines.append(f"- {k}: {extraction_status[k]}")
    if no_text_docs:
        lines.append("")
        lines.append("No-Text Files:")
        for d in no_text_docs[:50]:
            lines.append(f"- {d.file_name}")
        if len(no_text_docs) > 50:
            lines.append(f"- ... ({len(no_text_docs) - 50} more)")

    if errors:
        lines.append("")
        lines.append("Errors:")
        for e in errors[:50]:
            lines.append(f"- {e.get('stage','unknown')}: {e.get('file_name','')} | {e.get('error','')}")
        if len(errors) > 50:
            lines.append(f"- ... ({len(errors) - 50} more)")

    if usage:
        chat_in = usage.get("chat_in", 0)
        chat_out = usage.get("chat_out", 0)
        embed_in = usage.get("embed_in", 0)
        tokens_in = int(chat_in / 4)
        tokens_out = int(chat_out / 4)
        tokens_embed = int(embed_in / 4)
        total_tokens = tokens_in + tokens_out + tokens_embed
        lines.append("")
        lines.append("Token Estimates (approx):")
        lines.append(f"- chat_input_tokens: {tokens_in}")
        lines.append(f"- chat_output_tokens: {tokens_out}")
        lines.append(f"- embeddings_tokens: {tokens_embed}")
        lines.append(f"- total_tokens: {total_tokens}")

    report_path.write_text("\n".join(lines), encoding="utf-8")
    return report_path


def write_duplicate_group_overviews(output_dir: Path, docs: List[DocRecord]) -> None:
    """Write one unified review-group overview workbook per <Category>_Duplicate folder."""
    by_root: Dict[Path, Dict[str, List[Dict[str, Any]]]] = {}
    for d in docs:
        if not d.review_group_id:
            continue
        cat_folder = sanitize_folder(d.category)
        cluster_id = sanitize_folder(d.review_group_id)
        dup_root = output_dir / f"{cat_folder}_Duplicate"
        file_name = Path(d.moved_to).name if d.moved_to else d.file_name
        exact_score = float(d.duplicate_score) if d.duplicate_score is not None else 0.0
        near_score = float(d.near_duplicate_score) if d.near_duplicate_score is not None else 0.0
        relation = d.duplicate_relation_type or ""
        by_root.setdefault(dup_root, {}).setdefault(cluster_id, []).append(
            {
                "file_name": file_name,
                "relation": relation,
                "exact_score": exact_score,
                "near_score": near_score,
            }
        )

    for dup_root, clusters in by_root.items():
        dup_root.mkdir(parents=True, exist_ok=True)
        out_path = dup_root / "duplicate_groups_overview.xlsx"
        existing_assignments: Dict[Tuple[str, str], Dict[str, str]] = {}
        if out_path.exists():
            try:
                prev_df = pd.read_excel(out_path, sheet_name="Groups")
                for _, row in prev_df.iterrows():
                    gid = str(row.get("Group ID", "") or "").strip()
                    fname = str(row.get("FileName", "") or "").strip()
                    if not gid or not fname:
                        continue
                    if gid.startswith("Review Group "):
                        continue
                    existing_assignments[(gid, fname.lower())] = {
                        "assigned_to": str(row.get("Assigned to", "") or "").strip(),
                        "action": str(row.get("Action", "") or "").strip(),
                    }
            except Exception:
                existing_assignments = {}
        rows: List[Dict[str, Any]] = []
        relation_rank = {"exact": 0, "exact+near": 1, "near": 2}
        for cluster_id in sorted(clusters.keys()):
            rows.append(
                {
                    "Group ID": f"Review Group {cluster_id} - Assigned to:",
                    "Relation": "",
                    "FileName": "",
                    "Exact_sc": "",
                    "Near_sc": "",
                    "Assigned to": "",
                    "Action": "",
                }
            )
            members = sorted(
                clusters[cluster_id],
                key=lambda x: (
                    relation_rank.get(str(x.get("relation", "")).lower(), 3),
                    -(float(x.get("exact_score", 0.0)) if float(x.get("exact_score", 0.0)) > 0 else float(x.get("near_score", 0.0))),
                    str(x.get("file_name", "")).lower(),
                ),
            )
            for m in members:
                file_name = str(m.get("file_name", ""))
                prev = existing_assignments.get((cluster_id, file_name.lower()), {})
                rows.append(
                    {
                        "Group ID": cluster_id,
                        "Relation": m.get("relation", ""),
                        "FileName": file_name,
                        "Exact_sc": m.get("exact_score", 0.0),
                        "Near_sc": m.get("near_score", 0.0),
                        "Assigned to": prev.get("assigned_to", ""),
                        "Action": prev.get("action", ""),
                    }
                )

        df = sanitize_excel_df(
            pd.DataFrame(
                rows,
                columns=["Group ID", "Relation", "FileName", "Exact_sc", "Near_sc", "Assigned to", "Action"],
            )
        )
        with pd.ExcelWriter(out_path, engine="openpyxl") as writer:
            df.to_excel(writer, index=False, sheet_name="Groups")

        try:
            from openpyxl import load_workbook
            from openpyxl.styles import Alignment, Font

            wb = load_workbook(out_path)
            ws = wb["Groups"]
            wrap = Alignment(wrap_text=True, vertical="top")
            bold = Font(bold=True)
            widths = {"A": 34, "B": 14, "C": 32, "D": 11, "E": 11, "F": 16, "G": 16}
            for col, width in widths.items():
                ws.column_dimensions[col].width = width
            for row in range(2, ws.max_row + 1):
                cell = ws.cell(row=row, column=1)
                if str(cell.value or "").startswith("Review Group "):
                    cell.font = bold
                for col in range(1, 8):
                    ws.cell(row=row, column=col).alignment = wrap
            wb.save(out_path)
        except Exception:
            pass


def prompt_api_key_gui(title: str, label_text: str) -> Optional[str]:
    if tk is None:
        return None
    root = tk.Tk()
    root.title(title)
    root.geometry("560x240")
    apply_theme(root)

    container = tk.Frame(root, bg=THEME["bg"])
    container.pack(fill=tk.BOTH, expand=True, padx=16, pady=16)

    label = tk.Label(container, text=label_text, bg=THEME["bg"], fg=THEME["fg"], font=FONT_LABEL)
    label.pack(anchor="w", pady=(0, 8))

    entry = tk.Entry(container, show="*", width=60, bg=THEME["text_bg"], fg=THEME["fg"], insertbackground=THEME["fg"], font=FONT_BASE)
    entry.pack(fill=tk.X, pady=(0, 8))

    result: List[str] = []

    def on_ok() -> None:
        val = entry.get().strip()
        if not val:
            if messagebox:
                messagebox.showerror("Error", "API key is required.")
            return
        result.append(val)
        root.destroy()

    def on_cancel() -> None:
        root.destroy()

    btn_frame = tk.Frame(container, bg=THEME["bg"])
    btn_frame.pack(pady=12, anchor="e", fill=tk.X)
    tk.Button(btn_frame, text="Cancel", command=on_cancel, width=12, bg=THEME["btn_bg"], fg=THEME["btn_fg"], relief=tk.FLAT, height=1, font=FONT_BUTTON).pack(side=tk.RIGHT, padx=8, ipady=6)
    tk.Button(btn_frame, text="OK", command=on_ok, width=12, bg=THEME["accent"], fg="#ffffff", relief=tk.FLAT, height=1, font=FONT_BUTTON).pack(side=tk.RIGHT, padx=8, ipady=6)

    root.mainloop()
    return result[0] if result else None


def load_resume(out_dir: Path) -> Dict[str, Any]:
    resume_path = out_dir / RESUME_FILENAME
    if not resume_path.exists():
        return {"files": {}}
    try:
        with resume_path.open("r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {"files": {}}


def save_resume(out_dir: Path, resume: Dict[str, Any]) -> None:
    resume_path = out_dir / RESUME_FILENAME
    with resume_path.open("w", encoding="utf-8") as f:
        json.dump(resume, f)


def run_pipeline(
    input_dir: Path,
    output_dir: Path,
    categories: List[str],
    cfg: AzureConfig,
    dry_run: bool,
    use_resume: bool,
    ocrmypdf_enabled: bool,
    app_name: Optional[str],
    embeddings_source: str,
    append_excel: bool,
    category_path_map: Dict[str, Any],
    include_full_text_output: bool = DEFAULT_INCLUDE_FULL_TEXT_OUTPUT,
    limit: Optional[int] = None,
    no_move: bool = False,
    articles_enabled: bool = False,
    progress_cb: Optional[callable] = None,
) -> None:
    setup_logging(output_dir)
    logging.info("Starting pipeline")
    logging.info("Article generation enabled: %s", str(articles_enabled).lower())

    reset_usage()

    files, unsupported_files, staged_inputs = list_files(input_dir)
    total_files = len(files)
    input_stats = scan_input_stats(files)
    if limit is not None and limit > 0:
        files = files[:limit]
    processed_stats = scan_input_stats(files)
    if not files:
        logging.warning("No supported files found")
        try:
            unsupported_report_path = write_unsupported_files_report(output_dir, unsupported_files)
            logging.info("Wrote unsupported files report: %s", unsupported_report_path)
            if unsupported_files:
                cleanup_workbook_path = write_unsupported_cleanup_workbook(output_dir, unsupported_files)
                logging.info("Wrote unsupported cleanup workbook: %s", cleanup_workbook_path)
            report_path = write_summary_report(output_dir, [], [], unsupported_files, [], get_usage(), total_files, len(files), limit)
            logging.info("Wrote %s", report_path)
        except Exception as exc:
            logging.exception("Failed to write unsupported/summary reports: %s", exc)
        if staged_inputs is not None:
            staged_inputs.cleanup()
        return

    run_id = time.strftime("%Y%m%d%H%M%S")
    t0 = time.time()

    resume = load_resume(output_dir) if use_resume else {"files": {}}
    resume_files = resume.get("files", {})

    doc_items: List[Tuple[str, str, Optional[np.ndarray]]] = []
    article_items: List[Tuple[str, str, Optional[np.ndarray]]] = []
    docs: List[DocRecord] = []
    articles: List[ArticleRecord] = []

    doc_hashes: Dict[str, str] = {}
    doc_embeddings_by_id: Dict[str, Optional[np.ndarray]] = {}
    article_hashes: Dict[str, str] = {}
    doc_id_to_key: Dict[str, str] = {}
    article_id_to_key: Dict[str, str] = {}
    article_id_to_idx: Dict[str, int] = {}

    raw_texts: Dict[str, str] = {}
    article_texts: Dict[str, str] = {}
    extraction_statuses: Dict[str, str] = {}
    errors: List[Dict[str, str]] = []
    doc_summary_flags: Dict[str, set[str]] = {}

    iterable = tqdm(files, desc="Extracting") if tqdm else files
    for idx, input_file in enumerate(iterable, start=1):
        path = input_file.source_path
        display_path = input_file.display_path
        if progress_cb:
            progress_cb("Extracting", idx - 1, len(files))
        key = input_file.file_key
        cached = resume_files.get(key)
        doc_id = f"{run_id}-DOC-{idx:05d}"
        if cached and cached.get("doc_id"):
            doc_id = cached["doc_id"]
        logging.info("Processing %s", display_path)
        doc_id_to_key[doc_id] = key
        if cached:
            text = cached.get("text", "")
            pdf_articles = cached.get("articles_raw", []) if articles_enabled else []
            extraction_status = cached.get("extraction_status", "no_text")
            cached["doc_id"] = doc_id
        else:
            try:
                text, pdf_articles, extraction_status = process_file(
                    path,
                    ocrmypdf_enabled,
                    articles_enabled,
                    source_label=display_path,
                )
            except Exception as exc:
                logging.exception("Failed to extract %s: %s", display_path, exc)
                errors.append({"stage": "extract", "file_name": path.name, "file_path": display_path, "error": str(exc)})
                text = ""
                pdf_articles = []
                extraction_status = "no_text"
            resume_files[key] = {
                "doc_id": doc_id,
                "file_path": display_path,
                "file_name": path.name,
                "ext": path.suffix.lower(),
                "text": text,
                "articles_raw": pdf_articles,
                "extraction_status": extraction_status,
            }

        raw_texts[doc_id] = text
        extraction_statuses[doc_id] = resume_files.get(key, {}).get("extraction_status", extraction_status)
        normalized = normalize_text(text)
        hsh = exact_hash_for_file(path, key)
        doc_hashes[doc_id] = hsh

        emb_vec: Optional[np.ndarray] = None
        if embeddings_source == "full_text":
            cached_emb = cached.get("doc_embedding") if cached else None
            if cached_emb is not None:
                emb_vec = np.array(cached_emb, dtype=np.float32)
            elif normalized and len(normalized) >= MIN_EMBEDDING_CHARS and not dry_run:
                try:
                    emb = call_azure_embeddings(cfg, normalized[:MAX_CHARS_PER_CHUNK])
                    emb_vec = np.array(emb, dtype=np.float32)
                    resume_files[key]["doc_embedding"] = emb
                except Exception as exc:
                    logging.exception("Embedding failed for %s: %s", display_path, exc)
                    errors.append({"stage": "embedding", "file_name": path.name, "file_path": display_path, "error": str(exc)})

        doc_items.append((doc_id, hsh, emb_vec))
        doc_embeddings_by_id[doc_id] = emb_vec

        # Article handling (PDF only)
        if articles_enabled:
            for a_idx, (_title, body) in enumerate(pdf_articles, start=1):
                article_id = f"{doc_id}-A{a_idx:03d}"
                article_texts[article_id] = body
                ahash = exact_hash_for_text(body, f"{key}|article|{a_idx}")
                article_hashes[article_id] = ahash
                article_id_to_key[article_id] = key
                article_id_to_idx[article_id] = a_idx
                aemb_vec: Optional[np.ndarray] = None
                if embeddings_source == "full_text":
                    cached_aemb = None
                    if cached and "article_embeddings" in cached:
                        cached_aemb = cached["article_embeddings"].get(str(a_idx))
                    if cached_aemb is not None:
                        aemb_vec = np.array(cached_aemb, dtype=np.float32)
                    elif body.strip() and len(body) >= MIN_EMBEDDING_CHARS and not dry_run:
                        try:
                            aemb = call_azure_embeddings(cfg, body[:MAX_CHARS_PER_CHUNK])
                            aemb_vec = np.array(aemb, dtype=np.float32)
                            resume_files[key].setdefault("article_embeddings", {})[str(a_idx)] = aemb
                        except Exception as exc:
                            logging.exception("Article embedding failed for %s: %s", display_path, exc)
                            errors.append({"stage": "article_embedding", "file_name": path.name, "file_path": display_path, "error": str(exc)})
                article_items.append((article_id, ahash, aemb_vec))

        if use_resume:
            resume["files"] = resume_files
            save_resume(output_dir, resume)
        if progress_cb:
            progress_cb("Extracting", idx, len(files))

    if embeddings_source == "full_text":
        doc_dup_of, doc_dup_score, doc_dup_group = detect_duplicates(doc_items, DUPLICATE_THRESHOLD)
        art_dup_of, art_dup_score, art_dup_group = detect_duplicates(article_items, DUPLICATE_THRESHOLD)
    elif embeddings_source == EMBEDDINGS_SOURCE_NONE:
        doc_dup_of, doc_dup_score, doc_dup_group = detect_duplicates(doc_items, DUPLICATE_THRESHOLD)
        art_dup_of, art_dup_score, art_dup_group = detect_duplicates(article_items, DUPLICATE_THRESHOLD)
    else:
        doc_dup_of, doc_dup_score, doc_dup_group = {}, {}, {}
        art_dup_of, art_dup_score, art_dup_group = {}, {}, {}

    iterable2 = tqdm(files, desc="Summarizing") if tqdm else files
    for idx, input_file in enumerate(iterable2, start=1):
        path = input_file.source_path
        display_path = input_file.display_path
        if progress_cb:
            progress_cb("Summarizing", idx - 1, len(files))
        key = input_file.file_key
        doc_id = resume_files.get(key, {}).get("doc_id", f"{run_id}-DOC-{idx:05d}")
        text = raw_texts.get(doc_id, "")
        cached = resume_files.get(key, {})
        extraction_status = extraction_statuses.get(doc_id, "no_text")
        low_text = extraction_status != "ok" or len(text) < MIN_EXTRACTED_CHARS
        if cached.get("doc_summary") and not dry_run:
            summary = cached.get("doc_summary", {})
        elif dry_run or not text.strip() or low_text:
            summary = {}
        else:
            try:
                summary, summary_flag = summarize_document_safe(
                    cfg,
                    text,
                    categories,
                    path.name,
                    display_path,
                )
                if summary_flag:
                    doc_summary_flags.setdefault(doc_id, set()).add(summary_flag)
                resume_files[key]["doc_summary"] = summary
            except Exception as exc:
                logging.exception("Summarization failed for %s: %s", display_path, exc)
                errors.append({"stage": "summarize", "file_name": path.name, "file_path": display_path, "error": str(exc)})
                summary = {}

        category = (summary.get("category") or "uncategorized").strip()
        if low_text:
            category = UNREADABLE_CATEGORY
        if category not in categories and category not in ("Other", UNREADABLE_CATEGORY):
            category = "Other"

        tags = summary.get("tags") or []
        if isinstance(tags, str):
            tags = [t.strip() for t in tags.split(",") if t.strip()]
        # normalize/dedupe tags
        seen = set()
        norm_tags = []
        for t in tags:
            t = str(t).strip()
            tag_key = t.lower()
            if not t or tag_key in seen:
                continue
            seen.add(tag_key)
            norm_tags.append(t)
        tags = norm_tags[:MAX_TAGS]

        short_summary = (summary.get("short_summary") or "").strip()
        long_summary = (summary.get("long_summary") or "").strip()

        duplicate_of = doc_dup_of.get(doc_id, "")
        duplicate_score = doc_dup_score.get(doc_id)
        duplicate_group_id = doc_dup_group.get(doc_id, "")

        moved_to = ""
        review_flags = []
        if extraction_status != "ok":
            review_flags.append("low_text")
        if len(text) < MIN_EXTRACTED_CHARS:
            review_flags.append("short_text")
        summary_flags = doc_summary_flags.get(doc_id, set())
        if "summary_fallback_content_filter" in summary_flags:
            review_flags.append("summary_fallback_content_filter")
        if "summary_truncated_large_doc" in summary_flags:
            review_flags.append("summary_truncated_large_doc")

        docs.append(
            DocRecord(
                doc_id=doc_id,
                file_key=key,
                file_name=path.name,
                file_path=display_path,
                source_path=str(path),
                file_ext=path.suffix.lower(),
                category=category,
                tags=tags,
                short_summary=short_summary,
                long_summary=long_summary,
                word_count=len(text.split()),
                char_count=len(text),
                extraction_status=extraction_statuses.get(doc_id, "no_text"),
                review_flags=",".join(review_flags),
                duplicate_of=duplicate_of,
                duplicate_score=duplicate_score,
                duplicate_group_id=duplicate_group_id,
                near_duplicate_of="",
                near_duplicate_score=None,
                near_duplicate_group_id="",
                review_group_id="",
                duplicate_relation_type="",
                moved_to=moved_to,
            )
        )
        if progress_cb:
            progress_cb("Summarizing", idx, len(files))

        # Article summaries
        # Only for PDF (others have no articles)
        article_list = []
        # We re-split to align with doc order
        if articles_enabled and path.suffix.lower() == ".pdf":
            if cached and cached.get("articles_raw"):
                article_list = cached.get("articles_raw", [])
            else:
                try:
                    _, pages, _ = extract_text_pdf(path, ocrmypdf_enabled)
                    article_list = split_pdf_into_articles(pages, source_label=display_path)
                except Exception as exc:
                    logging.exception("Failed to split articles for %s: %s", display_path, exc)
                    errors.append({"stage": "split_articles", "file_name": path.name, "file_path": display_path, "error": str(exc)})
        for a_idx, (title, body) in enumerate(article_list, start=1):
            article_id = f"{doc_id}-A{a_idx:03d}"
            cached_summary = None
            if cached and "article_summaries" in cached:
                cached_summary = cached["article_summaries"].get(str(a_idx))
            if cached_summary is not None and not dry_run:
                art_summary = cached_summary
            elif dry_run or not body.strip():
                art_summary = ""
            else:
                try:
                    art_summary, _ = summarize_article_safe(cfg, body, path.name)
                    resume_files[key].setdefault("article_summaries", {})[str(a_idx)] = art_summary
                except Exception as exc:
                    logging.exception("Article summary failed for %s: %s", display_path, exc)
                    errors.append({"stage": "article_summarize", "file_name": path.name, "file_path": display_path, "error": str(exc)})
                    art_summary = ""
            articles.append(
                ArticleRecord(
                    doc_id=doc_id,
                    file_key=key,
                    file_name=path.name,
                    file_path=display_path,
                    article_index=a_idx,
                    article_title=title,
                    article_summary=art_summary,
                    duplicate_of=art_dup_of.get(article_id, ""),
                    duplicate_score=art_dup_score.get(article_id),
                    duplicate_group_id=art_dup_group.get(article_id, ""),
                )
            )

    if embeddings_source == "summary":
        doc_items2: List[Tuple[str, str, Optional[np.ndarray]]] = []
        article_items2: List[Tuple[str, str, Optional[np.ndarray]]] = []
        doc_embeddings_summary: Dict[str, Optional[np.ndarray]] = {}
        min_chars = min_embedding_chars_for_source(embeddings_source)

        for d in docs:
            key = doc_id_to_key.get(d.doc_id, "")
            normalized = normalize_text(raw_texts.get(d.doc_id, ""))
            emb_text = embedding_text_for_doc(embeddings_source, normalized, d.long_summary, d.short_summary)
            emb_vec: Optional[np.ndarray] = None
            cached_emb = None
            if key and key in resume_files:
                cached_emb = resume_files[key].get("doc_embedding_summary")
            if cached_emb is not None:
                emb_vec = np.array(cached_emb, dtype=np.float32)
            elif emb_text and len(emb_text) >= min_chars and not dry_run:
                try:
                    emb = call_azure_embeddings(cfg, emb_text[:MAX_CHARS_PER_CHUNK])
                    emb_vec = np.array(emb, dtype=np.float32)
                    if key:
                        resume_files[key]["doc_embedding_summary"] = emb
                except Exception as exc:
                    logging.exception("Embedding failed for %s: %s", d.file_path, exc)
            doc_items2.append((d.doc_id, doc_hashes.get(d.doc_id, ""), emb_vec))
            doc_embeddings_summary[d.doc_id] = emb_vec

        for a in articles:
            article_id = f"{a.doc_id}-A{a.article_index:03d}"
            key = article_id_to_key.get(article_id, "")
            idx = article_id_to_idx.get(article_id, None)
            emb_text = embedding_text_for_article(embeddings_source, article_texts.get(article_id, ""), a.article_summary)
            aemb_vec: Optional[np.ndarray] = None
            cached_aemb = None
            if key and idx is not None and key in resume_files:
                cached_aemb = resume_files[key].get("article_embedding_summary", {}).get(str(idx))
            if cached_aemb is not None:
                aemb_vec = np.array(cached_aemb, dtype=np.float32)
            elif emb_text and len(emb_text) >= min_chars and not dry_run:
                try:
                    aemb = call_azure_embeddings(cfg, emb_text[:MAX_CHARS_PER_CHUNK])
                    aemb_vec = np.array(aemb, dtype=np.float32)
                    if key and idx is not None:
                        resume_files[key].setdefault("article_embedding_summary", {})[str(idx)] = aemb
                except Exception as exc:
                    logging.exception("Article embedding failed for %s: %s", a.file_path, exc)
            article_items2.append((article_id, article_hashes.get(article_id, ""), aemb_vec))

        doc_dup_of, doc_dup_score, doc_dup_group = detect_duplicates(doc_items2, DUPLICATE_THRESHOLD)
        art_dup_of, art_dup_score, art_dup_group = detect_duplicates(article_items2, DUPLICATE_THRESHOLD)

        for d in docs:
            d.duplicate_of = doc_dup_of.get(d.doc_id, "")
            d.duplicate_score = doc_dup_score.get(d.doc_id)
            d.duplicate_group_id = doc_dup_group.get(d.doc_id, "")

        for a in articles:
            article_id = f"{a.doc_id}-A{a.article_index:03d}"
            a.duplicate_of = art_dup_of.get(article_id, "")
            a.duplicate_score = art_dup_score.get(article_id)
            a.duplicate_group_id = art_dup_group.get(article_id, "")
        final_doc_embeddings = doc_embeddings_summary
    else:
        final_doc_embeddings = doc_embeddings_by_id

    exact_pairs: set[Tuple[str, str]] = set()
    docs_by_id = {d.doc_id: d for d in docs}
    for d in docs:
        if not d.duplicate_of:
            continue
        if d.duplicate_of not in docs_by_id:
            continue
        exact_pairs.add(_sorted_pair(d.doc_id, d.duplicate_of))

    near_of, near_score, near_group, _near_adj, near_edges = detect_near_duplicates_docs(docs, final_doc_embeddings)
    review_groups, relation_types = build_unified_review_groups(docs, exact_pairs, near_edges)
    for d in docs:
        d.near_duplicate_of = near_of.get(d.doc_id, "")
        d.near_duplicate_score = near_score.get(d.doc_id)
        d.near_duplicate_group_id = near_group.get(d.doc_id, "")
        d.review_group_id = review_groups.get(d.doc_id, "")
        d.duplicate_relation_type = relation_types.get(d.doc_id, "")

    # Move files
    if not dry_run and not no_move:
        for move_idx, d in enumerate(docs, start=1):
            if progress_cb:
                progress_cb("Moving files", move_idx - 1, len(docs))
            src = Path(d.source_path)
            cat_folder = sanitize_folder(d.category)
            if d.review_group_id:
                cluster_id = sanitize_folder(d.review_group_id)
                dest_dir = output_dir / f"{cat_folder}_Duplicate" / cluster_id
            else:
                dest_dir = output_dir / cat_folder
            dest_dir.mkdir(parents=True, exist_ok=True)
            target = dest_dir / src.name
            if target.exists():
                stem = target.stem
                suffix = target.suffix
                suffix_idx = 1
                while True:
                    candidate = dest_dir / f"{stem}_{suffix_idx}{suffix}"
                    if not candidate.exists():
                        target = candidate
                        break
                    suffix_idx += 1
            try:
                shutil.move(str(src), str(target))
                d.moved_to = str(target)
            except Exception as exc:
                logging.exception("Failed to move %s: %s", src, exc)
                errors.append({"stage": "move", "file_name": src.name, "file_path": d.file_path, "error": str(exc)})
            if progress_cb:
                progress_cb("Moving files", move_idx, len(docs))
    if not dry_run:
        try:
            write_duplicate_group_overviews(output_dir, docs)
        except Exception as exc:
            logging.exception("Failed to write duplicate group overviews: %s", exc)
            errors.append({"stage": "write_duplicate_overview", "file_name": "", "file_path": str(output_dir), "error": str(exc)})

    full_text_rows: List[Dict[str, Any]] = []
    for d in docs:
        text = raw_texts.get(d.doc_id, "")
        full_text_rows.append(
            {
                "doc_id": d.doc_id,
                "file_key": d.file_key,
                "file_name": d.file_name,
                "file_path": d.file_path,
                "category": d.category,
                "short_summary": d.short_summary,
                "long_summary": d.long_summary,
                "tags": ", ".join(d.tags),
                "word_count": d.word_count,
                "char_count": d.char_count,
                "extraction_status": d.extraction_status,
                "review_flags": d.review_flags,
                "moved_to": d.moved_to,
                "full_text": text,
            }
        )

    # Write outputs
    peers_path = full_text_path = import_path = None
    try:
        peers_path, full_text_path, import_path = write_excels(
            output_dir,
            docs,
            articles,
            full_text_rows,
            app_name,
            append_excel,
            category_path_map,
            include_full_text_output,
            articles_enabled,
        )
        logging.info("Wrote summaries workbook: %s", peers_path)
        logging.info("Wrote import workbook: %s", import_path)
        if full_text_path is not None:
            logging.info("Wrote full-text archive: %s", full_text_path)
    except Exception as exc:
        logging.exception("Failed to write Excel outputs: %s", exc)
        errors.append({"stage": "write_excel", "file_name": "", "file_path": str(output_dir), "error": str(exc)})
    usage = get_usage()
    try:
        unsupported_report_path = write_unsupported_files_report(output_dir, unsupported_files)
        logging.info("Wrote unsupported files report: %s", unsupported_report_path)
        if unsupported_files:
            cleanup_workbook_path = write_unsupported_cleanup_workbook(output_dir, unsupported_files)
            logging.info("Wrote unsupported cleanup workbook: %s", cleanup_workbook_path)
        report_path = write_summary_report(output_dir, docs, articles, unsupported_files, errors, usage, total_files, len(files), limit)
        logging.info("Wrote %s", report_path)
    except Exception as exc:
        logging.exception("Failed to write summary report: %s", exc)

    if use_resume:
        resume["files"] = resume_files
        save_resume(output_dir, resume)

    elapsed = time.time() - t0
    save_last_run_stats(
        output_dir,
        {
            "timestamp": time.strftime("%Y-%m-%d %H:%M:%S"),
            "elapsed_sec": elapsed,
            "processed_files": len(files),
            "total_files": total_files,
            "total_size_mb": processed_stats.get("total_size_mb", 0.0),
            "ocr_enabled": ocrmypdf_enabled,
            "embeddings_source": embeddings_source,
            "chat_deployment": cfg.chat_deployment,
        },
    )
    if limit is not None and total_files > len(files):
        est_total = (elapsed / max(len(files), 1)) * total_files
        logging.info(
            "Estimate for %s files based on %s processed: ~%ss",
            total_files,
            len(files),
            int(est_total),
        )
        est_10k = (elapsed / max(len(files), 1)) * 10000
        logging.info("Estimate for 10000 files: ~%ss", int(est_10k))
    if staged_inputs is not None:
        staged_inputs.cleanup()


def run_pipeline_parallel(
    input_dir: Path,
    output_dir: Path,
    categories: List[str],
    cfg: AzureConfig,
    dry_run: bool,
    use_resume: bool,
    ocrmypdf_enabled: bool,
    app_name: Optional[str],
    embeddings_source: str,
    append_excel: bool,
    workers: int,
    category_path_map: Dict[str, Any],
    include_full_text_output: bool = DEFAULT_INCLUDE_FULL_TEXT_OUTPUT,
    limit: Optional[int] = None,
    no_move: bool = False,
    articles_enabled: bool = False,
) -> None:
    setup_logging(output_dir)
    logging.info("Starting pipeline (parallel, workers=%s)", workers)
    logging.info("Article generation enabled: %s", str(articles_enabled).lower())

    reset_usage()

    run_id = time.strftime("%Y%m%d%H%M%S")

    files, unsupported_files, staged_inputs = list_files(input_dir)
    total_files = len(files)
    if limit is not None and limit > 0:
        files = files[:limit]
    processed_stats = scan_input_stats(files)
    if not files:
        logging.warning("No supported files found")
        try:
            unsupported_report_path = write_unsupported_files_report(output_dir, unsupported_files)
            logging.info("Wrote unsupported files report: %s", unsupported_report_path)
            if unsupported_files:
                cleanup_workbook_path = write_unsupported_cleanup_workbook(output_dir, unsupported_files)
                logging.info("Wrote unsupported cleanup workbook: %s", cleanup_workbook_path)
            report_path = write_summary_report(output_dir, [], [], unsupported_files, [], get_usage(), total_files, len(files), limit)
            logging.info("Wrote %s", report_path)
        except Exception as exc:
            logging.exception("Failed to write unsupported/summary reports: %s", exc)
        if staged_inputs is not None:
            staged_inputs.cleanup()
        return

    t0 = time.time()

    resume = load_resume(output_dir) if use_resume else {"files": {}}
    resume_files = resume.get("files", {})

    raw_texts: Dict[str, str] = {}
    article_texts: Dict[str, str] = {}
    extraction_statuses: Dict[str, str] = {}
    doc_items: List[Tuple[str, str, Optional[np.ndarray]]] = []
    article_items: List[Tuple[str, str, Optional[np.ndarray]]] = []
    articles_by_doc: Dict[str, List[Tuple[str, str]]] = {}
    doc_hashes: Dict[str, str] = {}
    doc_embeddings_by_id: Dict[str, Optional[np.ndarray]] = {}
    article_hashes: Dict[str, str] = {}
    doc_id_to_key: Dict[str, str] = {}
    article_id_to_key: Dict[str, str] = {}
    article_id_to_idx: Dict[str, int] = {}
    errors: List[Dict[str, str]] = []
    doc_summary_flags: Dict[str, set[str]] = {}
    errors_lock = threading.Lock()
    state_lock = threading.Lock()

    def extract_and_embed(idx_path: Tuple[int, InputFile]) -> Tuple[int, InputFile, str, List[Tuple[str, str]], str, Optional[np.ndarray], List[Tuple[str, Optional[np.ndarray]]]]:
        idx, input_file = idx_path
        path = input_file.source_path
        display_path = input_file.display_path
        key = input_file.file_key
        with state_lock:
            cached = resume_files.get(key)
            doc_id = f"{run_id}-DOC-{idx:05d}"
            if cached and cached.get("doc_id"):
                doc_id = cached["doc_id"]
            doc_id_to_key[doc_id] = key
        if cached:
            text = cached.get("text", "")
            pdf_articles = cached.get("articles_raw", []) if articles_enabled else []
            extraction_status = cached.get("extraction_status", "no_text")
            cached["doc_id"] = doc_id
        else:
            try:
                text, pdf_articles, extraction_status = process_file(
                    path,
                    ocrmypdf_enabled,
                    articles_enabled,
                    source_label=display_path,
                )
            except Exception as exc:
                logging.exception("Failed to extract %s: %s", display_path, exc)
                with errors_lock:
                    errors.append({"stage": "extract", "file_name": path.name, "file_path": display_path, "error": str(exc)})
                text = ""
                pdf_articles = []
                extraction_status = "no_text"
            with state_lock:
                resume_files[key] = {
                    "doc_id": doc_id,
                    "file_path": display_path,
                    "file_name": path.name,
                    "ext": path.suffix.lower(),
                    "text": text,
                    "articles_raw": pdf_articles,
                    "extraction_status": extraction_status,
                }

        normalized = normalize_text(text)
        hsh = exact_hash_for_file(path, key)

        emb_vec: Optional[np.ndarray] = None
        if embeddings_source == "full_text":
            cached_emb = cached.get("doc_embedding") if cached else None
            if cached_emb is not None:
                emb_vec = np.array(cached_emb, dtype=np.float32)
            elif normalized and len(normalized) >= MIN_EMBEDDING_CHARS and not dry_run:
                try:
                    emb = call_azure_embeddings(cfg, normalized[:MAX_CHARS_PER_CHUNK])
                    emb_vec = np.array(emb, dtype=np.float32)
                    with state_lock:
                        resume_files[key]["doc_embedding"] = emb
                except Exception as exc:
                    logging.exception("Embedding failed for %s: %s", display_path, exc)

        art_embs: List[Tuple[str, Optional[np.ndarray]]] = []
        if articles_enabled:
            for a_idx, (_title, body) in enumerate(pdf_articles, start=1):
                article_id = f"{doc_id}-A{a_idx:03d}"
                aemb_vec: Optional[np.ndarray] = None
                if embeddings_source == "full_text":
                    cached_aemb = None
                    if cached and "article_embeddings" in cached:
                        cached_aemb = cached["article_embeddings"].get(str(a_idx))
                    if cached_aemb is not None:
                        aemb_vec = np.array(cached_aemb, dtype=np.float32)
                    elif body.strip() and len(body) >= MIN_EMBEDDING_CHARS and not dry_run:
                        try:
                            aemb = call_azure_embeddings(cfg, body[:MAX_CHARS_PER_CHUNK])
                            aemb_vec = np.array(aemb, dtype=np.float32)
                            with state_lock:
                                resume_files[key].setdefault("article_embeddings", {})[str(a_idx)] = aemb
                        except Exception as exc:
                            logging.exception("Article embedding failed for %s: %s", display_path, exc)
                            with errors_lock:
                                errors.append({"stage": "article_embedding", "file_name": path.name, "file_path": display_path, "error": str(exc)})
                art_embs.append((article_id, aemb_vec))

        if use_resume:
            with state_lock:
                resume["files"] = resume_files
                save_resume(output_dir, resume)

        return idx, input_file, text, pdf_articles, extraction_status, (hsh, emb_vec), art_embs

    with ThreadPoolExecutor(max_workers=workers) as ex:
        futures = [ex.submit(extract_and_embed, (i, p)) for i, p in enumerate(files, start=1)]
        for fut in as_completed(futures):
            try:
                idx, input_file, text, pdf_articles, extraction_status, doc_hash_emb, art_embs = fut.result()
            except Exception as exc:
                logging.exception("Worker failed: %s", exc)
                with errors_lock:
                    errors.append({"stage": "worker", "file_name": "", "file_path": "", "error": str(exc)})
                continue
            path = input_file.source_path
            with state_lock:
                doc_id = resume_files.get(input_file.file_key, {}).get("doc_id", f"{run_id}-DOC-{idx:05d}")
            raw_texts[doc_id] = text
            extraction_statuses[doc_id] = extraction_status
            hsh, emb_vec = doc_hash_emb
            doc_hashes[doc_id] = hsh
            doc_items.append((doc_id, hsh, emb_vec))
            doc_embeddings_by_id[doc_id] = emb_vec
            articles_by_doc[doc_id] = pdf_articles
            if articles_enabled:
                for a_idx, ((article_id, aemb_vec), (_title, body)) in enumerate(zip(art_embs, pdf_articles), start=1):
                    article_texts[article_id] = body
                    key = input_file.file_key
                    article_id_to_key[article_id] = key
                    article_id_to_idx[article_id] = a_idx
                    ahash = exact_hash_for_text(body, f"{key}|article|{a_idx}")
                    article_hashes[article_id] = ahash
                    article_items.append((article_id, ahash, aemb_vec))

    if embeddings_source == "full_text":
        doc_dup_of, doc_dup_score, doc_dup_group = detect_duplicates(doc_items, DUPLICATE_THRESHOLD)
        art_dup_of, art_dup_score, art_dup_group = detect_duplicates(article_items, DUPLICATE_THRESHOLD)
    elif embeddings_source == EMBEDDINGS_SOURCE_NONE:
        doc_dup_of, doc_dup_score, doc_dup_group = detect_duplicates(doc_items, DUPLICATE_THRESHOLD)
        art_dup_of, art_dup_score, art_dup_group = detect_duplicates(article_items, DUPLICATE_THRESHOLD)
    else:
        doc_dup_of, doc_dup_score, doc_dup_group = {}, {}, {}
        art_dup_of, art_dup_score, art_dup_group = {}, {}, {}

    docs: List[DocRecord] = []
    articles: List[ArticleRecord] = []

    def summarize_doc(idx_path: Tuple[int, InputFile]) -> Tuple[int, InputFile, Dict[str, Any], str]:
        idx, input_file = idx_path
        path = input_file.source_path
        display_path = input_file.display_path
        key = input_file.file_key
        with state_lock:
            doc_id = resume_files.get(key, {}).get("doc_id", f"{run_id}-DOC-{idx:05d}")
        text = raw_texts.get(doc_id, "")
        with state_lock:
            cached = resume_files.get(key, {})
        extraction_status = extraction_statuses.get(doc_id, "no_text")
        low_text = extraction_status != "ok" or len(text) < MIN_EXTRACTED_CHARS
        if cached.get("doc_summary") and not dry_run:
            summary = cached.get("doc_summary", {})
        elif dry_run or not text.strip() or low_text:
            summary = {}
        else:
            try:
                summary, summary_flag = summarize_document_safe(
                    cfg,
                    text,
                    categories,
                    path.name,
                    display_path,
                )
                with state_lock:
                    resume_files[key]["doc_summary"] = summary
            except Exception as exc:
                logging.exception("Summarization failed for %s: %s", display_path, exc)
                with errors_lock:
                    errors.append({"stage": "summarize", "file_name": path.name, "file_path": display_path, "error": str(exc)})
                summary = {}
                summary_flag = ""
        if dry_run or not text.strip() or low_text or cached.get("doc_summary"):
            summary_flag = ""
        if use_resume:
            with state_lock:
                resume["files"] = resume_files
                save_resume(output_dir, resume)
        return idx, input_file, summary, summary_flag

    with ThreadPoolExecutor(max_workers=workers) as ex:
        futures = [ex.submit(summarize_doc, (i, p)) for i, p in enumerate(files, start=1)]
        summaries: Dict[str, Dict[str, Any]] = {}
        for fut in as_completed(futures):
            try:
                idx, input_file, summary, summary_flag = fut.result()
            except Exception as exc:
                logging.exception("Summarize worker failed: %s", exc)
                with errors_lock:
                    errors.append({"stage": "summarize_worker", "file_name": "", "file_path": "", "error": str(exc)})
                continue
            with state_lock:
                doc_id = resume_files.get(input_file.file_key, {}).get("doc_id", f"{run_id}-DOC-{idx:05d}")
            summaries[doc_id] = summary
            if summary_flag:
                doc_summary_flags.setdefault(doc_id, set()).add(summary_flag)

    for idx, input_file in enumerate(files, start=1):
        path = input_file.source_path
        display_path = input_file.display_path
        key = input_file.file_key
        doc_id = resume_files.get(key, {}).get("doc_id", f"{run_id}-DOC-{idx:05d}")
        summary = summaries.get(doc_id, {})
        text = raw_texts.get(doc_id, "")
        extraction_status = extraction_statuses.get(doc_id, "no_text")
        low_text = extraction_status != "ok" or len(text) < MIN_EXTRACTED_CHARS

        category = (summary.get("category") or "uncategorized").strip()
        if low_text:
            category = UNREADABLE_CATEGORY
        if category not in categories and category not in ("Other", UNREADABLE_CATEGORY):
            category = "Other"

        tags = summary.get("tags") or []
        if isinstance(tags, str):
            tags = [t.strip() for t in tags.split(",") if t.strip()]
        seen = set()
        norm_tags = []
        for t in tags:
            t = str(t).strip()
            tag_key = t.lower()
            if not t or tag_key in seen:
                continue
            seen.add(tag_key)
            norm_tags.append(t)
        tags = norm_tags[:MAX_TAGS]

        short_summary = (summary.get("short_summary") or "").strip()
        long_summary = (summary.get("long_summary") or "").strip()

        duplicate_of = doc_dup_of.get(doc_id, "")
        duplicate_score = doc_dup_score.get(doc_id)
        duplicate_group_id = doc_dup_group.get(doc_id, "")

        review_flags = []
        if extraction_status != "ok":
            review_flags.append("low_text")
        if len(text) < MIN_EXTRACTED_CHARS:
            review_flags.append("short_text")
        summary_flags = doc_summary_flags.get(doc_id, set())
        if "summary_fallback_content_filter" in summary_flags:
            review_flags.append("summary_fallback_content_filter")
        if "summary_truncated_large_doc" in summary_flags:
            review_flags.append("summary_truncated_large_doc")

        docs.append(
            DocRecord(
                doc_id=doc_id,
                file_key=key,
                file_name=path.name,
                file_path=display_path,
                source_path=str(path),
                file_ext=path.suffix.lower(),
                category=category,
                tags=tags,
                short_summary=short_summary,
                long_summary=long_summary,
                word_count=len(text.split()),
                char_count=len(text),
                extraction_status=extraction_statuses.get(doc_id, "no_text"),
                review_flags=",".join(review_flags),
                duplicate_of=duplicate_of,
                duplicate_score=duplicate_score,
                duplicate_group_id=duplicate_group_id,
                near_duplicate_of="",
                near_duplicate_score=None,
                near_duplicate_group_id="",
                review_group_id="",
                duplicate_relation_type="",
                moved_to="",
            )
        )

        if articles_enabled:
            article_list = articles_by_doc.get(doc_id, [])
            for a_idx, (title, body) in enumerate(article_list, start=1):
                article_id = f"{doc_id}-A{a_idx:03d}"
                cached = resume_files.get(key, {})
                cached_summary = None
                if cached and "article_summaries" in cached:
                    cached_summary = cached["article_summaries"].get(str(a_idx))
                if cached_summary is not None and not dry_run:
                    art_summary = cached_summary
                elif dry_run or not body.strip():
                    art_summary = ""
                else:
                    try:
                        art_summary, _ = summarize_article_safe(cfg, body, path.name)
                        resume_files[key].setdefault("article_summaries", {})[str(a_idx)] = art_summary
                    except Exception as exc:
                        logging.exception("Article summary failed for %s: %s", display_path, exc)
                        with errors_lock:
                            errors.append({"stage": "article_summarize", "file_name": path.name, "file_path": display_path, "error": str(exc)})
                        art_summary = ""
                articles.append(
                    ArticleRecord(
                        doc_id=doc_id,
                        file_key=key,
                        file_name=path.name,
                        file_path=display_path,
                        article_index=a_idx,
                        article_title=title,
                        article_summary=art_summary,
                        duplicate_of=art_dup_of.get(article_id, ""),
                        duplicate_score=art_dup_score.get(article_id),
                        duplicate_group_id=art_dup_group.get(article_id, ""),
                    )
                )

    if embeddings_source == "summary":
        doc_items2: List[Tuple[str, str, Optional[np.ndarray]]] = []
        article_items2: List[Tuple[str, str, Optional[np.ndarray]]] = []
        doc_embeddings_summary: Dict[str, Optional[np.ndarray]] = {}
        min_chars = min_embedding_chars_for_source(embeddings_source)

        for d in docs:
            key = doc_id_to_key.get(d.doc_id, "")
            normalized = normalize_text(raw_texts.get(d.doc_id, ""))
            emb_text = embedding_text_for_doc(embeddings_source, normalized, d.long_summary, d.short_summary)
            emb_vec: Optional[np.ndarray] = None
            cached_emb = None
            if key and key in resume_files:
                cached_emb = resume_files[key].get("doc_embedding_summary")
            if cached_emb is not None:
                emb_vec = np.array(cached_emb, dtype=np.float32)
            elif emb_text and len(emb_text) >= min_chars and not dry_run:
                try:
                    emb = call_azure_embeddings(cfg, emb_text[:MAX_CHARS_PER_CHUNK])
                    emb_vec = np.array(emb, dtype=np.float32)
                    if key:
                        resume_files[key]["doc_embedding_summary"] = emb
                except Exception as exc:
                    logging.exception("Embedding failed for %s: %s", d.file_path, exc)
            doc_items2.append((d.doc_id, doc_hashes.get(d.doc_id, ""), emb_vec))
            doc_embeddings_summary[d.doc_id] = emb_vec

        for a in articles:
            article_id = f"{a.doc_id}-A{a.article_index:03d}"
            key = article_id_to_key.get(article_id, "")
            idx = article_id_to_idx.get(article_id, None)
            emb_text = embedding_text_for_article(embeddings_source, article_texts.get(article_id, ""), a.article_summary)
            aemb_vec: Optional[np.ndarray] = None
            cached_aemb = None
            if key and idx is not None and key in resume_files:
                cached_aemb = resume_files[key].get("article_embedding_summary", {}).get(str(idx))
            if cached_aemb is not None:
                aemb_vec = np.array(cached_aemb, dtype=np.float32)
            elif emb_text and len(emb_text) >= min_chars and not dry_run:
                try:
                    aemb = call_azure_embeddings(cfg, emb_text[:MAX_CHARS_PER_CHUNK])
                    aemb_vec = np.array(aemb, dtype=np.float32)
                    if key and idx is not None:
                        resume_files[key].setdefault("article_embedding_summary", {})[str(idx)] = aemb
                except Exception as exc:
                    logging.exception("Article embedding failed for %s: %s", a.file_path, exc)
            article_items2.append((article_id, article_hashes.get(article_id, ""), aemb_vec))

        doc_dup_of, doc_dup_score, doc_dup_group = detect_duplicates(doc_items2, DUPLICATE_THRESHOLD)
        art_dup_of, art_dup_score, art_dup_group = detect_duplicates(article_items2, DUPLICATE_THRESHOLD)

        for d in docs:
            d.duplicate_of = doc_dup_of.get(d.doc_id, "")
            d.duplicate_score = doc_dup_score.get(d.doc_id)
            d.duplicate_group_id = doc_dup_group.get(d.doc_id, "")

        for a in articles:
            article_id = f"{a.doc_id}-A{a.article_index:03d}"
            a.duplicate_of = art_dup_of.get(article_id, "")
            a.duplicate_score = art_dup_score.get(article_id)
            a.duplicate_group_id = art_dup_group.get(article_id, "")
        final_doc_embeddings = doc_embeddings_summary
    else:
        final_doc_embeddings = doc_embeddings_by_id

    exact_pairs: set[Tuple[str, str]] = set()
    docs_by_id = {d.doc_id: d for d in docs}
    for d in docs:
        if not d.duplicate_of:
            continue
        if d.duplicate_of not in docs_by_id:
            continue
        exact_pairs.add(_sorted_pair(d.doc_id, d.duplicate_of))

    near_of, near_score, near_group, _near_adj, near_edges = detect_near_duplicates_docs(docs, final_doc_embeddings)
    review_groups, relation_types = build_unified_review_groups(docs, exact_pairs, near_edges)
    for d in docs:
        d.near_duplicate_of = near_of.get(d.doc_id, "")
        d.near_duplicate_score = near_score.get(d.doc_id)
        d.near_duplicate_group_id = near_group.get(d.doc_id, "")
        d.review_group_id = review_groups.get(d.doc_id, "")
        d.duplicate_relation_type = relation_types.get(d.doc_id, "")

    if not dry_run and not no_move:
        for d in docs:
            src = Path(d.source_path)
            cat_folder = sanitize_folder(d.category)
            if d.review_group_id:
                cluster_id = sanitize_folder(d.review_group_id)
                dest_dir = output_dir / f"{cat_folder}_Duplicate" / cluster_id
            else:
                dest_dir = output_dir / cat_folder
            dest_dir.mkdir(parents=True, exist_ok=True)
            target = dest_dir / src.name
            if target.exists():
                stem = target.stem
                suffix = target.suffix
                i = 1
                while True:
                    candidate = dest_dir / f"{stem}_{i}{suffix}"
                    if not candidate.exists():
                        target = candidate
                        break
                    i += 1
            try:
                shutil.move(str(src), str(target))
                d.moved_to = str(target)
            except Exception as exc:
                logging.exception("Failed to move %s: %s", src, exc)
                with errors_lock:
                    errors.append({"stage": "move", "file_name": src.name, "file_path": d.file_path, "error": str(exc)})
    if not dry_run:
        try:
            write_duplicate_group_overviews(output_dir, docs)
        except Exception as exc:
            logging.exception("Failed to write duplicate group overviews: %s", exc)
            with errors_lock:
                errors.append({"stage": "write_duplicate_overview", "file_name": "", "file_path": str(output_dir), "error": str(exc)})

    full_text_rows: List[Dict[str, Any]] = []
    for d in docs:
        text = raw_texts.get(d.doc_id, "")
        full_text_rows.append(
            {
                "doc_id": d.doc_id,
                "file_key": d.file_key,
                "file_name": d.file_name,
                "file_path": d.file_path,
                "category": d.category,
                "short_summary": d.short_summary,
                "long_summary": d.long_summary,
                "tags": ", ".join(d.tags),
                "word_count": d.word_count,
                "char_count": d.char_count,
                "extraction_status": d.extraction_status,
                "review_flags": d.review_flags,
                "moved_to": d.moved_to,
                "full_text": text,
            }
        )

    peers_path = full_text_path = import_path = None
    try:
        peers_path, full_text_path, import_path = write_excels(
            output_dir,
            docs,
            articles,
            full_text_rows,
            app_name,
            append_excel,
            category_path_map,
            include_full_text_output,
            articles_enabled,
        )
        logging.info("Wrote summaries workbook: %s", peers_path)
        logging.info("Wrote import workbook: %s", import_path)
        if full_text_path is not None:
            logging.info("Wrote full-text archive: %s", full_text_path)
    except Exception as exc:
        logging.exception("Failed to write Excel outputs: %s", exc)
        with errors_lock:
            errors.append({"stage": "write_excel", "file_name": "", "file_path": str(output_dir), "error": str(exc)})
    usage = get_usage()
    try:
        unsupported_report_path = write_unsupported_files_report(output_dir, unsupported_files)
        logging.info("Wrote unsupported files report: %s", unsupported_report_path)
        if unsupported_files:
            cleanup_workbook_path = write_unsupported_cleanup_workbook(output_dir, unsupported_files)
            logging.info("Wrote unsupported cleanup workbook: %s", cleanup_workbook_path)
        report_path = write_summary_report(output_dir, docs, articles, unsupported_files, errors, usage, total_files, len(files), limit)
        logging.info("Wrote %s", report_path)
    except Exception as exc:
        logging.exception("Failed to write summary report: %s", exc)

    if use_resume:
        resume["files"] = resume_files
        save_resume(output_dir, resume)

    elapsed = time.time() - t0
    save_last_run_stats(
        output_dir,
        {
            "timestamp": time.strftime("%Y-%m-%d %H:%M:%S"),
            "elapsed_sec": elapsed,
            "processed_files": len(files),
            "total_files": total_files,
            "total_size_mb": processed_stats.get("total_size_mb", 0.0),
            "ocr_enabled": ocrmypdf_enabled,
            "embeddings_source": embeddings_source,
            "chat_deployment": cfg.chat_deployment,
        },
    )
    if limit is not None and total_files > len(files):
        est_total = (elapsed / max(len(files), 1)) * total_files
        logging.info(
            "Estimate for %s files based on %s processed: ~%ss",
            total_files,
            len(files),
            int(est_total),
        )
    if staged_inputs is not None:
        staged_inputs.cleanup()


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="DocAtlas document processing pipeline")
    parser.add_argument("--input", help="Input folder")
    parser.add_argument("--output", help="Output folder")
    parser.add_argument("--categories", help="Categories separated by semicolons")
    parser.add_argument("--config", help="Path to applications config JSON")
    parser.add_argument("--app", help="Application name from config")
    parser.add_argument("--dry-run", action="store_true", help="Do not call APIs or move files")
    parser.add_argument("--no-resume", action="store_true", help="Disable resume cache")
    parser.add_argument("--no-ocrmypdf", action="store_true", help="Disable OCRmyPDF and use Tesseract fallback")
    parser.add_argument("--edit-config", action="store_true", help="Open GUI editor for applications config")
    parser.add_argument("--embeddings-source", choices=["summary", "full_text", "none"], help="Use summaries, full text, or disable embeddings")
    parser.add_argument("--overwrite-excel", action="store_true", help="Overwrite Excel outputs instead of appending")
    parser.add_argument("--limit", type=int, help="Process only the first N files (for estimation)")
    parser.add_argument("--no-move", action="store_true", help="Do not move files (for estimation)")
    parser.add_argument("--charter-mode", action="store_true", help="Preview-only mode (no file moves)")
    parser.add_argument("--signal-scan", action="store_true", help="Deprecated alias for --charter-mode")
    parser.add_argument("--test-embeddings", action="store_true", help="Test embeddings endpoint and exit")
    parser.add_argument("--test-chat", action="store_true", help="Test chat endpoint and exit")
    parser.add_argument("--workers", type=int, default=1, help="Number of workers for parallel CLI processing (default: 1)")
    parser.add_argument("--articles", action="store_true", help="Enable PDF article splitting/summarization and write Articles sheet")
    parser.add_argument("--no-articles", action="store_true", help="Deprecated alias; article generation is disabled by default")
    parser.add_argument("--interactive", action="store_true", help="Prompt for missing CLI values in the terminal instead of opening the GUI")
    parser.add_argument("--category-path-map", help="Path to category_path_map.json for import Path mapping")
    parser.add_argument(
        "--include-full-text-output",
        action="store_true",
        help="Deprecated compatibility flag; full-text JSONL.GZ archive is now written by default",
    )
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    argv_tokens = sys.argv[1:]
    if args.workers < 1:
        raise ValueError("--workers must be >= 1")
    config_path = Path(args.config) if args.config else Path(__file__).with_name("applications.json")
    category_path_map_path = Path(args.category_path_map) if args.category_path_map else Path(__file__).with_name(DEFAULT_CATEGORY_PATH_MAP_FILENAME)
    category_path_map = load_category_path_map(category_path_map_path)
    app_config = load_app_config(config_path)
    is_gui_flow = (not args.interactive) and not (args.input and args.output and (args.categories or args.app))

    if args.edit_config:
        if tk is None:
            raise RuntimeError("tkinter is not available")
        root = tk.Tk()
        root.withdraw()
        edit_applications_gui(config_path, app_config, root)
        root.mainloop()
        return 0

    if args.test_embeddings:
        cfg = azure_config_from_env(require_key=False)
        if not cfg.embeddings_api_key:
            cfg.embeddings_api_key = prompt_api_key_gui(
                "DocAtlas - Enter Embeddings API Key",
                "Paste embeddings API key (not stored):",
            ) or ""
        if not cfg.embeddings_api_key:
            raise ValueError("AZURE_EMBEDDINGS_API_KEY is not set")
        try:
            emb = call_azure_embeddings(cfg, "test embedding")
            print(f"Embeddings OK. Vector length: {len(emb)}")
            return 0
        except Exception as exc:
            print(f"Embeddings test failed: {exc}")
            return 1

    if args.test_chat:
        cfg = azure_config_from_env(require_key=False)
        if not cfg.chat_api_key:
            cfg.chat_api_key = prompt_api_key_gui(
                "DocAtlas - Enter LLM API Key",
                "Paste LLM API key (not stored):",
            ) or ""
        if not cfg.chat_api_key:
            raise ValueError("AZURE_CHAT_API_KEY is not set")
        if not cfg.api_key:
            cfg.api_key = cfg.chat_api_key
        try:
            msg = [{"role": "user", "content": "Say OK"}]
            out = call_azure_chat(cfg, msg)
            print("Chat OK. Response:", out[:200])
            return 0
        except Exception as exc:
            print(f"Chat test failed: {exc}")
            return 1

    validate_app_and_category_map(app_config, category_path_map)
    if args.include_full_text_output:
        logging.warning(
            "--include-full-text-output is deprecated; DocAtlas now writes "
            "<app>__docatlas_full_text.jsonl.gz by default."
        )

    append_excel = not args.overwrite_excel
    if args.charter_mode or args.signal_scan:
        args.no_move = True

    if args.input and args.output and (args.categories or args.app):
        input_dir = Path(args.input)
        output_dir = Path(args.output)
        if args.categories:
            categories = [c.strip() for c in args.categories.split(";") if c.strip()]
            app_name = None
        elif args.app and args.app in app_config:
            categories = app_config[args.app]
            app_name = args.app
        else:
            raise ValueError("Provide --categories or a valid --app from config")
        ocrmypdf_enabled = not args.no_ocrmypdf
        embeddings_source = resolve_embeddings_source(args.embeddings_source)
        articles_enabled = bool(args.articles)
        if args.no_articles:
            logging.warning("--no-articles is deprecated and now the default behavior; use --articles to enable article generation.")
    elif args.interactive:
        (
            input_dir,
            output_dir,
            categories,
            app_name,
            ocrmypdf_enabled,
            embeddings_source,
            append_excel,
            args.no_move,
            articles_enabled,
            args.workers,
        ) = resolve_cli_interactive_inputs(args, argv_tokens, app_config)
        if args.no_articles:
            logging.warning("--no-articles is deprecated and now the default behavior; use --articles to enable article generation.")
    else:
        input_dir, output_dir = pick_directories_gui()
        categories, app_name = get_categories_gui(app_config, config_path)
        ocrmypdf_enabled = get_ocrmypdf_gui()
        embeddings_source, append_excel = get_embeddings_source_gui()
        gui_charter_mode = get_run_mode_gui()
        if gui_charter_mode is None:
            return 0
        if gui_charter_mode:
            args.no_move = True
        articles_enabled = get_articles_mode_gui()

    cfg = azure_config_from_env(require_key=(not args.dry_run and not is_gui_flow))
    if not args.dry_run:
        if not cfg.chat_api_key:
            if is_gui_flow:
                cfg.chat_api_key = prompt_api_key_gui("DocAtlas - Enter LLM API Key", "Paste LLM API key (not stored):") or ""
            if not cfg.chat_api_key:
                raise ValueError("AZURE_CHAT_API_KEY is not set")
        if embeddings_source != EMBEDDINGS_SOURCE_NONE and not cfg.embeddings_api_key:
            if is_gui_flow:
                cfg.embeddings_api_key = prompt_api_key_gui(
                    "DocAtlas - Enter Embeddings API Key",
                    "Paste embeddings API key (not stored):",
                ) or ""
            if embeddings_source != EMBEDDINGS_SOURCE_NONE and not cfg.embeddings_api_key:
                raise ValueError("AZURE_EMBEDDINGS_API_KEY is not set")
        if not cfg.api_key:
            cfg.api_key = cfg.chat_api_key
    warn_missing_ocr_deps(ocrmypdf_enabled)
    try:
        est_files, _est_unsupported_files, est_staged_inputs = list_files(input_dir)
        try:
            est_stats = scan_input_stats(est_files)
            est_sec, est_source, settings_match = quick_estimate_runtime(
                est_stats,
                output_dir,
                ocrmypdf_enabled,
                embeddings_source,
                cfg.chat_deployment,
            )
        finally:
            if est_staged_inputs is not None:
                est_staged_inputs.cleanup()
        if is_gui_flow and messagebox is not None and est_sec:
            note = ""
            if est_source == "baseline" and not settings_match:
                note = "\nNote: baseline settings differ from current run."
            msg = (
                f"Files: {est_stats.get('count', 0)}\n"
                f"Total size: {est_stats.get('total_size_mb', 0.0):.1f} MB\n"
                f"Estimated time: ~{format_duration(est_sec)}{note}"
            )
            messagebox.showinfo("DocAtlas - Estimated Runtime", msg)
    except Exception:
        pass

    if not (args.input and args.output and (args.categories or args.app)):
        # GUI progress window
        if tk is None:
            raise RuntimeError("tkinter is not available")

        q: "queue.Queue[Tuple[str, int, int]]" = queue.Queue()
        stage_times: Dict[str, float] = {}
        stage_start: Dict[str, float] = {}

        def progress_cb(stage: str, current: int, total: int) -> None:
            q.put((stage, current, total))

        def worker() -> None:
            run_pipeline(
                input_dir,
                output_dir,
                categories,
                cfg,
                args.dry_run,
                not args.no_resume,
                ocrmypdf_enabled,
                app_name,
                embeddings_source,
                append_excel,
                category_path_map,
                args.include_full_text_output,
                args.limit,
                args.no_move,
                articles_enabled,
                progress_cb,
            )
            q.put(("DONE", 1, 1))

        progress_root = tk.Tk()
        progress_root.title("DocAtlas - Processing Documents")
        progress_root.geometry("520x180")
        apply_theme(progress_root)

        frame = tk.Frame(progress_root, bg=THEME["bg"])
        frame.pack(fill=tk.BOTH, expand=True, padx=16, pady=16)

        stage_label = tk.Label(frame, text="Starting...", bg=THEME["bg"], fg=THEME["fg"], font=FONT_HEADER)
        stage_label.pack(anchor="w", pady=(0, 6))

        eta_label = tk.Label(frame, text="", bg=THEME["bg"], fg=THEME["muted"], font=FONT_SMALL)
        eta_label.pack(anchor="w", pady=(0, 10))

        prog = ttk.Progressbar(frame, length=460, mode="determinate")
        prog.pack(pady=6)

        def poll() -> None:
            try:
                while True:
                    stage, current, total = q.get_nowait()
                    if stage == "DONE":
                        stage_label.config(text="Completed")
                        eta_label.config(text="Finished")
                        prog["value"] = 100
                        progress_root.after(1200, progress_root.destroy)
                        return

                    if stage not in stage_start:
                        stage_start[stage] = time.time()
                    elapsed = time.time() - stage_start[stage]
                    if current > 0:
                        eta = (elapsed / current) * (total - current)
                        eta_label.config(text=f"ETA: ~{int(eta)}s")
                    else:
                        eta_label.config(text="ETA: estimating...")

                    stage_label.config(text=f"{stage} ({current}/{total})")
                    prog["maximum"] = max(total, 1)
                    prog["value"] = current
            except queue.Empty:
                pass
            progress_root.after(200, poll)

        t = threading.Thread(target=worker, daemon=True)
        t.start()
        poll()
        progress_root.mainloop()
    else:
        if args.workers > 1:
            run_pipeline_parallel(
                input_dir,
                output_dir,
                categories,
                cfg,
                args.dry_run,
                not args.no_resume,
                ocrmypdf_enabled,
                app_name,
                embeddings_source,
                append_excel,
                args.workers,
                category_path_map,
                args.include_full_text_output,
                args.limit,
                args.no_move,
                articles_enabled,
            )
        else:
            run_pipeline(
                input_dir,
                output_dir,
                categories,
                cfg,
                args.dry_run,
                not args.no_resume,
                ocrmypdf_enabled,
                app_name,
                embeddings_source,
                append_excel,
                category_path_map,
                args.include_full_text_output,
                args.limit,
                args.no_move,
                articles_enabled,
            )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
